# controllers/report_controller.py
"""
controllers/report_controller.py
Report generation controller that coordinates report creation.
This consolidates functionality from report_generator.py and related reporting code.
"""

import os
import re
import copy
import math
import shutil
import tempfile
import threading
import traceback
from typing import Optional, Dict, Any, List, Tuple
from datetime import datetime

import pandas as pd
import tkinter as tk
from tkinter import ttk, messagebox, filedialog

# Local imports
from utils import debug_print, show_success_message, get_save_path, plotting_sheet_test, get_plot_sheet_names, get_resource_path
from models.data_model import DataModel
from models.report_model import ReportModel, ReportConfig, ReportData
import processing


class ReportController:
    """Controller for report generation operations."""
    
    def __init__(self, data_model: DataModel, report_service: Any):
        """Initialize the report controller."""
        self.data_model = data_model
        self.report_service = report_service
        self.report_model = ReportModel()
        
        # Cross-controller references
        self.file_controller: Optional['FileController'] = None
        self.plot_controller: Optional['PlotController'] = None
        
        # GUI reference will be set when connected
        self.gui = None
        
        # Lazy loading infrastructure
        self._matplotlib_loaded = False
        self._pptx_loaded = False
        self._pil_loaded = False
        
        # Lazy loaded components
        self.plt = None
        self.Presentation = None
        self.Image = None
        
        # Report generation state
        self.temp_image_files = []  # Track temporary files for cleanup
        self.current_report_type = None
        
        print("DEBUG: ReportController initialized")
        print("DEBUG: Lazy loading system initialized for report components")
    
    def set_gui_reference(self, gui):
        """Set reference to main GUI for UI operations."""
        self.gui = gui
        print("DEBUG: ReportController connected to GUI")
    
    def set_file_controller(self, file_controller: 'FileController'):
        """Set reference to file controller."""
        self.file_controller = file_controller
        print("DEBUG: ReportController connected to FileController")
    
    def set_plot_controller(self, plot_controller: 'PlotController'):
        """Set reference to plot controller."""
        self.plot_controller = plot_controller
        print("DEBUG: ReportController connected to PlotController")
    
    def _lazy_import_matplotlib(self):
        """Lazy import matplotlib components."""
        try:
            import matplotlib.pyplot as plt
            print("TIMING: Lazy loaded matplotlib for ReportController")
            return plt
        except ImportError as e:
            print(f"ERROR: importing matplotlib: {e}")
            return None
    
    def _lazy_import_pptx(self):
        """Lazy import python-pptx components."""
        try:
            from pptx import Presentation
            from pptx.enum.text import PP_ALIGN
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            print("TIMING: Lazy loaded python-pptx for ReportController")
            return Presentation, PP_ALIGN, Inches, Pt, RGBColor
        except ImportError as e:
            print(f"ERROR: importing python-pptx: {e}")
            return None, None, None, None, None
    
    def _lazy_import_pil(self):
        """Lazy import PIL components."""
        try:
            from PIL import Image
            print("TIMING: Lazy loaded PIL for ReportController")
            return Image
        except ImportError as e:
            print(f"ERROR: importing PIL: {e}")
            return None
    
    def get_matplotlib(self):
        """Get matplotlib with lazy loading."""
        if not self._matplotlib_loaded:
            self.plt = self._lazy_import_matplotlib()
            self._matplotlib_loaded = True
        return self.plt
    
    def get_pptx_components(self):
        """Get python-pptx components with lazy loading."""
        if not self._pptx_loaded:
            components = self._lazy_import_pptx()
            if components[0] is not None:
                self.Presentation, self.PP_ALIGN, self.Inches, self.Pt, self.RGBColor = components
            self._pptx_loaded = True
        return self.Presentation, self.PP_ALIGN, self.Inches, self.Pt, self.RGBColor
    
    def get_pil(self):
        """Get PIL with lazy loading."""
        if not self._pil_loaded:
            self.Image = self._lazy_import_pil()
            self._pil_loaded = True
        return self.Image
    
    def generate_test_report(self, selected_sheet: str, sheets: Dict[str, Any], plot_options: List[str]) -> bool:
        """Generate an Excel and PowerPoint report for only the specified sheet."""
        debug_print(f"DEBUG: ReportController generating test report for {selected_sheet}")
        
        try:
            self.current_report_type = "test"
            
            # Check if sheet exists
            if selected_sheet not in sheets:
                debug_print(f"ERROR: Sheet '{selected_sheet}' not found in sheets data")
                messagebox.showerror("Error", f"Sheet '{selected_sheet}' not found.")
                return False
            
            # Show header selection dialog
            selected_headers = self._show_header_selector_dialog()
            if not selected_headers:
                debug_print("DEBUG: Header selection cancelled")
                return False
            
            debug_print(f"DEBUG: Selected headers for test report: {selected_headers}")
            
            # Get save path
            save_path = get_save_path(".xlsx")
            if not save_path:
                debug_print("DEBUG: Save cancelled")
                return False
            
            # Start report generation
            config = ReportConfig(
                report_type="test",
                output_format="excel",
                include_plots=True,
                include_images=True
            )
            self.report_model.set_config(config)
            
            report_data = ReportData(
                title=f"Test Report - {selected_sheet}",
                sheets_data={selected_sheet: sheets[selected_sheet]}
            )
            self.report_model.start_generation(report_data)
            
            # Process the sheet data
            sheet_info = sheets[selected_sheet]
            if not isinstance(sheet_info, dict) or "data" not in sheet_info:
                debug_print(f"ERROR: Invalid sheet data structure for {selected_sheet}")
                return False
            
            data = sheet_info["data"]
            debug_print(f"DEBUG: Processing sheet data with shape: {data.shape}")
            
            # Get processing function and process data
            process_function = processing.get_processing_function(selected_sheet)
            processed_data, _, full_sample_data = process_function(data)
            
            if processed_data.empty:
                debug_print(f"ERROR: No processed data for sheet {selected_sheet}")
                return False
            
            # Reorder data based on selected headers
            processed_data = self._reorder_processed_data(processed_data, selected_headers)
            debug_print(f"DEBUG: Reordered processed data shape: {processed_data.shape}")
            
            # Generate Excel report
            images_to_delete = []
            with pd.ExcelWriter(save_path, engine='xlsxwriter') as writer:
                valid_plot_options = processing.get_valid_plot_options(plot_options, full_sample_data)
                debug_print(f"DEBUG: Test report valid plot options: {valid_plot_options}")
                self._write_excel_report(writer, selected_sheet, processed_data, full_sample_data,
                                       valid_plot_options, images_to_delete)
            
            # Generate PowerPoint report
            ppt_save_path = save_path.replace('.xlsx', '.pptx')
            self._write_powerpoint_report_for_test(ppt_save_path, images_to_delete, selected_sheet,
                                                 processed_data, full_sample_data, plot_options)
            
            # Cleanup
            self._cleanup_images(images_to_delete)
            
            # Complete report generation
            self.report_model.complete_generation(save_path)
            
            # Show success message
            show_success_message("Success", 
                               f"Test report saved successfully to:\nExcel: {save_path}\nPowerPoint: {ppt_save_path}", 
                               self.gui.root if self.gui else None)
            
            debug_print("DEBUG: Test report generation completed successfully")
            return True
            
        except Exception as e:
            debug_print(f"ERROR: Test report generation failed: {e}")
            traceback.print_exc()
            if self.gui:
                messagebox.showerror("Error", f"An error occurred while generating the test report: {e}")
            return False
    
    def generate_full_report(self, filtered_sheets: Dict[str, Any], plot_options: List[str]) -> bool:
        """Generate a full report (Excel and PowerPoint) for all sheets."""
        debug_print("DEBUG: ReportController generating full report")
        
        try:
            self.current_report_type = "full"
            
            # Show header selection dialog
            selected_headers = self._show_header_selector_dialog()
            if not selected_headers:
                debug_print("DEBUG: Header selection cancelled")
                return False
            
            debug_print(f"DEBUG: Selected headers for full report: {selected_headers}")
            
            # Get save path
            save_path = get_save_path(".xlsx")
            if not save_path:
                debug_print("DEBUG: Save cancelled")
                return False
            
            # Start report generation
            config = ReportConfig(
                report_type="full",
                output_format="powerpoint",
                include_plots=True,
                include_images=True
            )
            self.report_model.set_config(config)
            
            report_data = ReportData(
                title="Full Data Report",
                sheets_data=filtered_sheets
            )
            self.report_model.start_generation(report_data)
            
            # Setup paths and tracking
            ppt_save_path = save_path.replace('.xlsx', '.pptx')
            images_to_delete = []
            total_sheets = len(filtered_sheets)
            processed_count = 0
            
            debug_print(f"DEBUG: Processing {total_sheets} sheets for full report")
            
            # Generate Excel report
            try:
                with pd.ExcelWriter(save_path, engine='xlsxwriter') as writer:
                    for sheet_name, sheet_info in filtered_sheets.items():
                        try:
                            if not isinstance(sheet_info, dict) or "data" not in sheet_info:
                                debug_print(f"DEBUG: Skipping sheet '{sheet_name}': No valid 'data' key found.")
                                continue
                            
                            data = sheet_info["data"]
                            debug_print(f"DEBUG: Processing sheet: {sheet_name}")
                            debug_print(f"DEBUG: Sheet data shape: {data.shape}")
                            
                            # Special handling for User Test Simulation
                            if sheet_name == "User Test Simulation":
                                debug_print(f"DEBUG: Processing User Test Simulation with 8-column format")
                            
                            is_plotting = plotting_sheet_test(sheet_name, data)
                            debug_print(f"DEBUG: Sheet {sheet_name} is_plotting: {is_plotting}")
                            
                            process_function = processing.get_processing_function(sheet_name)
                            
                            if is_plotting:
                                processed_data, _, full_sample_data = process_function(data)
                                debug_print(f"DEBUG: Processed data shape: {processed_data.shape}")
                                debug_print(f"DEBUG: Full sample data shape: {full_sample_data.shape}")
                                valid_plot_options = processing.get_valid_plot_options(plot_options, full_sample_data)
                                debug_print(f"DEBUG: Valid plot options: {valid_plot_options}")
                            else:
                                data = data.astype(str).replace([pd.NA], '')
                                processed_data, _, full_sample_data = process_function(data)
                                valid_plot_options = []
                            
                            if processed_data.empty or full_sample_data.empty:
                                debug_print(f"DEBUG: Skipping sheet '{sheet_name}' due to empty processed data.")
                                continue
                            
                            # Apply header selection and reordering
                            debug_print(f"DEBUG: Original processed_data columns: {processed_data.columns.tolist()}")
                            processed_data = self._reorder_processed_data(processed_data, selected_headers)
                            debug_print(f"DEBUG: Reordered processed_data columns: {processed_data.columns.tolist()}")
                            
                            self._write_excel_report(writer, sheet_name, processed_data, full_sample_data,
                                                   valid_plot_options, images_to_delete)
                            processed_count += 1
                            
                        except Exception as e:
                            debug_print(f"DEBUG: Error processing sheet '{sheet_name}': {e}")
                            traceback.print_exc()
                            continue
                
                # Update progress for Excel completion
                if self.gui and hasattr(self.gui, 'progress_dialog'):
                    self.gui.progress_dialog.update_progress_bar(50)
                    self.gui.root.update_idletasks()
                
                debug_print("DEBUG: Excel report generation completed")
                
            except Exception as e:
                debug_print(f"ERROR: Excel report generation failed: {e}")
                raise
            
            # Generate PowerPoint report
            try:
                def update_ppt_prog():
                    def callback(percent):
                        total_percent = 50 + (percent * 0.5)  # PPT is second half
                        if self.gui and hasattr(self.gui, 'progress_dialog'):
                            self.gui.progress_dialog.update_progress_bar(total_percent)
                            self.gui.root.update_idletasks()
                    return callback
                
                self._write_powerpoint_report(ppt_save_path, images_to_delete, plot_options, 
                                             selected_headers, progress_callback=update_ppt_prog())
                
                # Final progress update
                if self.gui and hasattr(self.gui, 'progress_dialog'):
                    self.gui.progress_dialog.update_progress_bar(100)
                    self.gui.root.update_idletasks()
                
                debug_print("DEBUG: PowerPoint report generation completed")
                
            except Exception as e:
                debug_print(f"ERROR: PowerPoint report generation failed: {e}")
                raise
            
            # Complete report generation
            self.report_model.complete_generation(save_path)
            
            debug_print("DEBUG: Full report generation completed successfully")
            return True
            
        except Exception as e:
            # Delete partial files
            for path in [save_path, ppt_save_path]:
                if 'save_path' in locals() and path and os.path.exists(path):
                    try:
                        os.remove(path)
                        debug_print(f"DEBUG: Cleaned up partial file: {path}")
                    except Exception as cleanup_error:
                        debug_print(f"WARNING: Failed to clean up {path}: {cleanup_error}")
            
            debug_print(f"ERROR: Full report generation failed: {e}")
            traceback.print_exc()
            return False
            
        finally:
            # Always cleanup temporary images
            self._cleanup_images(images_to_delete if 'images_to_delete' in locals() else [])
    
    def _show_header_selector_dialog(self) -> Optional[List[str]]:
        """Show header selector dialog and return selected headers."""
        debug_print("DEBUG: Showing header selector dialog")
        
        try:
            dialog = HeaderSelectorDialog(self.gui.root if self.gui else None)
            selected_headers = dialog.show()
            
            if selected_headers:
                debug_print(f"DEBUG: Headers selected: {selected_headers}")
                return selected_headers
            else:
                debug_print("DEBUG: Header selection cancelled")
                return None
                
        except Exception as e:
            debug_print(f"ERROR: Failed to show header selector dialog: {e}")
            return None
    
    def _reorder_processed_data(self, processed_data: pd.DataFrame, selected_headers: List[str]) -> pd.DataFrame:
        """Reorder processed_data columns based on selected headers."""
        debug_print(f"DEBUG: Reordering columns based on selected headers")
        
        try:
            original_columns = processed_data.columns.tolist()
            debug_print(f"DEBUG: Original columns: {original_columns}")
            debug_print(f"DEBUG: Selected headers: {selected_headers}")
            
            # Create new column order based on selected headers
            new_column_order = []
            
            # Add selected headers in order if they exist in the data
            for header in selected_headers:
                if header in original_columns:
                    new_column_order.append(header)
                    debug_print(f"DEBUG: Added header to new order: {header}")
            
            # Add any remaining columns that weren't in the selected headers
            for col in original_columns:
                if col not in new_column_order:
                    new_column_order.append(col)
                    debug_print(f"DEBUG: Added remaining column: {col}")
            
            # Reorder the DataFrame
            reordered_data = processed_data[new_column_order]
            debug_print(f"DEBUG: Reordered data shape: {reordered_data.shape}")
            debug_print(f"DEBUG: New column order: {new_column_order}")
            
            return reordered_data
            
        except Exception as e:
            debug_print(f"ERROR: Failed to reorder processed data: {e}")
            # Return original data if reordering fails
            return processed_data
    
    def _write_excel_report(self, writer, sheet_name: str, processed_data: pd.DataFrame, 
                           full_sample_data: pd.DataFrame, valid_plot_options: List[str] = None, 
                           images_to_delete: List[str] = None) -> None:
        """Write Excel report for a single sheet."""
        debug_print(f"DEBUG: Writing Excel report for sheet: {sheet_name}")
        
        try:
            if valid_plot_options is None:
                valid_plot_options = []
            if images_to_delete is None:
                images_to_delete = []
            
            # Write data to Excel
            processed_data.astype(str).replace([pd.NA], '')
            processed_data.to_excel(writer, sheet_name=sheet_name, index=False)
            
            # Add plots if this is a plotting sheet
            plot_sheet_names = get_plot_sheet_names()
            is_plotting = sheet_name in plot_sheet_names
            
            if is_plotting:
                debug_print(f"DEBUG: Adding plots to Excel for sheet: {sheet_name}")
                self._add_plots_to_excel(writer, sheet_name, full_sample_data, images_to_delete, valid_plot_options)
            
            debug_print(f"DEBUG: Excel report written successfully for sheet: {sheet_name}")
            
        except Exception as e:
            debug_print(f"ERROR: Failed to write Excel report for sheet '{sheet_name}': {e}")
            traceback.print_exc()
    
    def _add_plots_to_excel(self, writer, sheet_name: str, full_sample_data: pd.DataFrame, 
                           images_to_delete: List[str], valid_plot_options: List[str]) -> None:
        """Add plots to Excel worksheet."""
        debug_print(f"DEBUG: Adding plots to Excel worksheet: {sheet_name}")
        
        try:
            plt = self.get_matplotlib()
            if not plt:
                debug_print("ERROR: Matplotlib not available for Excel plots")
                return
            
            worksheet = writer.sheets[sheet_name]
            numeric_data = full_sample_data.apply(pd.to_numeric, errors='coerce')
            
            # Get sample names
            sample_names = None
            if hasattr(self, 'header_data') and self.header_data and 'samples' in self.header_data:
                sample_names = [sample['id'] for sample in self.header_data['samples']]
                debug_print(f"DEBUG: Extracted sample names from header_data: {sample_names}")
            
            # Check for User Test Simulation
            is_user_test_simulation = sheet_name in ["User Test Simulation", "User Simulation Test"]
            num_columns_per_sample = 8 if is_user_test_simulation else 12
            
            debug_print(f"DEBUG: Generating {len(valid_plot_options)} plots for Excel")
            
            # Generate each plot
            for i, plot_option in enumerate(valid_plot_options):
                try:
                    debug_print(f"DEBUG: Generating plot {i+1}/{len(valid_plot_options)}: {plot_option}")
                    
                    if is_user_test_simulation:
                        fig, _ = processing.plot_user_test_simulation_samples(
                            full_sample_data, num_columns_per_sample, plot_option, sample_names
                        )
                    else:
                        fig, _ = processing.plot_all_samples(
                            full_sample_data, num_columns_per_sample, plot_option, sample_names
                        )
                    
                    if fig:
                        # Save plot as image
                        plot_image_path = f"temp_plot_{sheet_name}_{plot_option}_{i}.png"
                        fig.savefig(plot_image_path, dpi=150, bbox_inches='tight')
                        
                        # Calculate position in Excel
                        start_col = len(full_sample_data.columns) + 2
                        start_row = i * 25  # Space plots vertically
                        
                        # Insert image into Excel
                        worksheet.insert_image(start_row, start_col, plot_image_path, 
                                             {'x_scale': 0.7, 'y_scale': 0.7})
                        
                        plt.close(fig)
                        
                        # Track for cleanup
                        if plot_image_path not in images_to_delete:
                            images_to_delete.append(plot_image_path)
                        
                        debug_print(f"DEBUG: Successfully added plot {plot_option} to Excel")
                    
                except Exception as e:
                    debug_print(f"ERROR: Failed to generate plot '{plot_option}' for sheet '{sheet_name}': {e}")
                    traceback.print_exc()
            
            debug_print(f"DEBUG: Completed adding plots to Excel for sheet: {sheet_name}")
            
        except Exception as e:
            debug_print(f"ERROR: Failed to add plots to Excel: {e}")
            traceback.print_exc()
    
    def _write_powerpoint_report_for_test(self, ppt_save_path: str, images_to_delete: List[str], 
                                         sheet_name: str, processed_data: pd.DataFrame, 
                                         full_sample_data: pd.DataFrame, plot_options: List[str]) -> None:
        """Write PowerPoint report for a single test sheet."""
        debug_print(f"DEBUG: Writing PowerPoint report for test sheet: {sheet_name}")
        
        try:
            Presentation, PP_ALIGN, Inches, Pt, RGBColor = self.get_pptx_components()
            if not Presentation:
                debug_print("ERROR: python-pptx not available for PowerPoint generation")
                return
            
            # Create presentation
            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            
            # Get images for this sheet
            current_file = getattr(self.gui, 'current_file', None)
            image_paths = []
            if (self.gui and hasattr(self.gui, 'sheet_images') and current_file and 
                current_file in self.gui.sheet_images and sheet_name in self.gui.sheet_images[current_file]):
                image_paths = self.gui.sheet_images[current_file][sheet_name]
            
            # Create main content slide
            main_slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # Add background and logo
            self._setup_slide_background(main_slide, prs, Inches)
            
            # Add title
            title_shape = main_slide.shapes.add_textbox(Inches(0.45), Inches(-0.04), 
                                                       Inches(10.72), Inches(0.64))
            text_frame = title_shape.text_frame
            text_frame.clear()
            p = text_frame.add_paragraph()
            p.text = sheet_name
            p.alignment = PP_ALIGN.LEFT
            run = p.runs[0]
            run.font.name = "Montserrat"
            run.font.size = Pt(32)
            run.font.bold = True
            
            # Determine if this is a plotting sheet
            is_plotting = plotting_sheet_test(sheet_name, processed_data)
            
            if is_plotting:
                # Add table and plots
                table_width = Inches(8.07)
                self._add_table_to_slide(main_slide, processed_data, table_width, is_plotting, Inches)
                
                # Add plots
                valid_plot_options = processing.get_valid_plot_options(plot_options, full_sample_data)
                if valid_plot_options:
                    self._add_plots_to_slide(main_slide, sheet_name, full_sample_data, 
                                           valid_plot_options, images_to_delete, Inches)
            else:
                # Add full-width table
                table_width = Inches(13.03)
                self._add_table_to_slide(main_slide, processed_data, table_width, is_plotting, Inches)
            
            # Add image slide if images exist
            if image_paths:
                valid_image_paths = [path for path in image_paths if os.path.exists(path)]
                if valid_image_paths:
                    debug_print(f"DEBUG: Adding image slide with {len(valid_image_paths)} images")
                    img_slide = prs.slides.add_slide(prs.slide_layouts[6])
                    self._setup_image_slide(prs, img_slide, sheet_name, Inches, PP_ALIGN, Pt)
                    self._add_images_to_slide(img_slide, valid_image_paths, Inches)
            
            # Clean presentation tables
            processing.clean_presentation_tables(prs)
            
            # Save presentation
            prs.save(ppt_save_path)
            debug_print(f"DEBUG: PowerPoint test report saved successfully: {ppt_save_path}")
            
        except Exception as e:
            debug_print(f"ERROR: Failed to write PowerPoint test report: {e}")
            traceback.print_exc()
    
    def _write_powerpoint_report(self, ppt_save_path: str, images_to_delete: List[str], 
                                plot_options: List[str], selected_headers: List[str], 
                                progress_callback=None) -> None:
        """Write full PowerPoint report with all sheets."""
        debug_print(f"DEBUG: Writing full PowerPoint report: {ppt_save_path}")
        
        try:
            Presentation, PP_ALIGN, Inches, Pt, RGBColor = self.get_pptx_components()
            if not Presentation:
                debug_print("ERROR: python-pptx not available for PowerPoint generation")
                return
            
            # Create presentation
            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            
            # Calculate total slides for progress tracking
            total_slides = 1  # Cover slide
            if self.gui and hasattr(self.gui, 'filtered_sheets'):
                for sheet_name in self.gui.filtered_sheets.keys():
                    total_slides += 1  # Main slide
                    # Check if sheet has images
                    if (hasattr(self.gui, 'sheet_images') and hasattr(self.gui, 'current_file') and
                        self.gui.current_file in self.gui.sheet_images and 
                        sheet_name in self.gui.sheet_images.get(self.gui.current_file, {})):
                        total_slides += 1  # Image slide
            
            processed_slides = 0
            
            # Create cover slide
            cover_slide = prs.slides.add_slide(prs.slide_layouts[6])
            bg_path = get_resource_path("resources/Cover_Page_Logo.jpg")
            if os.path.exists(bg_path):
                cover_slide.shapes.add_picture(bg_path, left=Inches(0), top=Inches(0),
                                             width=prs.slide_width, height=prs.slide_height)
            
            # Add cover title
            textbox_title = cover_slide.shapes.add_textbox(
                left=Inches((prs.slide_width.inches - 12) / 2),
                top=Inches(2.35),
                width=Inches(12),
                height=Inches(0.88)
            )
            text_frame = textbox_title.text_frame
            text_frame.clear()
            p = text_frame.add_paragraph()
            p.text = "Data Analysis Report"
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0]
            run.font.name = "Montserrat"
            run.font.size = Pt(48)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            
            processed_slides += 1
            if progress_callback:
                progress_callback(processed_slides, total_slides)
            
            # Process each sheet
            if self.gui and hasattr(self.gui, 'filtered_sheets'):
                for sheet_name, sheet_info in self.gui.filtered_sheets.items():
                    try:
                        debug_print(f"DEBUG: Processing PowerPoint slide for sheet: {sheet_name}")
                        
                        if not isinstance(sheet_info, dict) or "data" not in sheet_info:
                            debug_print(f"DEBUG: Skipping sheet '{sheet_name}': No valid data")
                            continue
                        
                        data = sheet_info["data"]
                        is_plotting = plotting_sheet_test(sheet_name, data)
                        
                        # Process data
                        process_function = processing.get_processing_function(sheet_name)
                        if is_plotting:
                            processed_data, _, full_sample_data = process_function(data)
                            valid_plot_options = processing.get_valid_plot_options(plot_options, full_sample_data)
                        else:
                            data = data.astype(str).replace([pd.NA], '')
                            processed_data, _, full_sample_data = process_function(data)
                            valid_plot_options = []
                        
                        if processed_data.empty:
                            debug_print(f"DEBUG: Skipping sheet '{sheet_name}': Empty processed data")
                            continue
                        
                        # Apply header reordering
                        processed_data = self._reorder_processed_data(processed_data, selected_headers)
                        
                        # Create main slide for this sheet
                        slide = prs.slides.add_slide(prs.slide_layouts[6])
                        self._setup_slide_background(slide, prs, Inches)
                        
                        # Add title
                        title_shape = slide.shapes.add_textbox(Inches(0.45), Inches(-0.04), 
                                                             Inches(10.72), Inches(0.64))
                        text_frame = title_shape.text_frame
                        text_frame.clear()
                        p = text_frame.add_paragraph()
                        p.text = sheet_name
                        p.alignment = PP_ALIGN.LEFT
                        run = p.runs[0]
                        run.font.name = "Montserrat"
                        run.font.size = Pt(32)
                        run.font.bold = True
                        
                        # Add content based on sheet type
                        if is_plotting:
                            table_width = Inches(8.07)
                            self._add_table_to_slide(slide, processed_data, table_width, is_plotting, Inches)
                            if valid_plot_options:
                                self._add_plots_to_slide(slide, sheet_name, full_sample_data, 
                                                       valid_plot_options, images_to_delete, Inches)
                            else:
                                debug_print(f"DEBUG: No valid plot options for sheet '{sheet_name}'. Skipping plots.")
                        else:
                            table_width = Inches(13.03)
                            self._add_table_to_slide(slide, processed_data, table_width, is_plotting, Inches)
                        
                        # Move title to front
                        if slide.shapes.title:
                            title_shape = slide.shapes.title
                            spTree = slide.shapes._spTree
                            spTree.remove(title_shape._element)
                            spTree.append(title_shape._element)
                        
                        processed_slides += 1
                        if progress_callback:
                            progress_callback(processed_slides, total_slides)
                        
                        # Add image slide if images exist
                        if (hasattr(self.gui, 'sheet_images') and hasattr(self.gui, 'current_file') and
                            self.gui.current_file in self.gui.sheet_images and 
                            sheet_name in self.gui.sheet_images.get(self.gui.current_file, {})):
                            
                            image_paths = self.gui.sheet_images[self.gui.current_file][sheet_name]
                            valid_image_paths = [path for path in image_paths if os.path.exists(path)]
                            
                            if valid_image_paths:
                                debug_print(f"DEBUG: Adding image slide for sheet: {sheet_name}")
                                img_slide = prs.slides.add_slide(prs.slide_layouts[6])
                                self._setup_image_slide(prs, img_slide, sheet_name, Inches, PP_ALIGN, Pt)
                                self._add_images_to_slide(img_slide, valid_image_paths, Inches)
                                
                                processed_slides += 1
                                if progress_callback:
                                    progress_callback(processed_slides, total_slides)
                        
                    except Exception as sheet_error:
                        debug_print(f"ERROR: Error processing sheet '{sheet_name}' for PowerPoint: {sheet_error}")
                        processed_slides += 1
                        traceback.print_exc()
                        continue
            
            # Clean presentation tables
            processing.clean_presentation_tables(prs)
            
            # Save presentation
            prs.save(ppt_save_path)
            debug_print(f"DEBUG: Full PowerPoint report saved successfully: {ppt_save_path}")
            
        except Exception as e:
            debug_print(f"ERROR: Failed to write full PowerPoint report: {e}")
            traceback.print_exc()
    
    def _setup_slide_background(self, slide, prs, Inches):
        """Setup slide background and logo."""
        try:
            # Add background
            background_path = get_resource_path("resources/ccell_background.png")
            if os.path.exists(background_path):
                slide.shapes.add_picture(background_path, Inches(0), Inches(0),
                                       width=prs.slide_width, height=prs.slide_height)
            
            # Add logo
            logo_path = get_resource_path("resources/ccell_logo_full.png")
            if os.path.exists(logo_path):
                slide.shapes.add_picture(logo_path, Inches(11.21), Inches(0.43),
                                       width=Inches(1.57), height=Inches(0.53))
            
        except Exception as e:
            debug_print(f"ERROR: Failed to setup slide background: {e}")
    
    def _setup_image_slide(self, prs, slide, sheet_name, Inches, PP_ALIGN, Pt):
        """Setup image slide with background and title."""
        try:
            # Add background
            background_path = get_resource_path("resources/ccell_background.png")
            if os.path.exists(background_path):
                slide.shapes.add_picture(background_path, Inches(0), Inches(0),
                                       width=prs.slide_width, height=prs.slide_height)
            
            # Add logo
            logo_path = get_resource_path("resources/ccell_logo_full.png")
            if os.path.exists(logo_path):
                slide.shapes.add_picture(logo_path, Inches(11.21), Inches(0.43),
                                       width=Inches(1.57), height=Inches(0.53))
            
            # Add title
            title_shape = slide.shapes.add_textbox(Inches(0.45), Inches(-0.04), 
                                                 Inches(10.72), Inches(0.64))
            text_frame = title_shape.text_frame
            text_frame.clear()
            p = text_frame.add_paragraph()
            p.text = f"{sheet_name} - Images"
            p.alignment = PP_ALIGN.LEFT
            run = p.runs[0]
            run.font.name = "Montserrat"
            run.font.size = Pt(32)
            run.font.bold = True
            
        except Exception as e:
            debug_print(f"ERROR: Failed to setup image slide: {e}")
    
    def _add_table_to_slide(self, slide, data: pd.DataFrame, table_width, is_plotting: bool, Inches):
        """Add data table to PowerPoint slide."""
        debug_print(f"DEBUG: Adding table to slide with {data.shape[0]} rows and {data.shape[1]} columns")
        
        try:
            # Calculate table dimensions
            rows, cols = data.shape
            if rows == 0 or cols == 0:
                debug_print("WARNING: Empty data, skipping table")
                return
            
            # Add header row
            table_rows = rows + 1
            
            # Position table
            left = Inches(0.45)
            top = Inches(0.85)
            height = Inches(min(5.5, table_rows * 0.3))  # Limit height
            
            # Create table
            table = slide.shapes.add_table(table_rows, cols, left, top, table_width, height).table
            
            # Set column headers
            for col_idx, column_name in enumerate(data.columns):
                cell = table.cell(0, col_idx)
                cell.text = str(column_name)
                # Format header
                cell.fill.solid()
                cell.fill.fore_color.rgb = (0, 0, 0)  # Black background
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = (255, 255, 255)  # White text
                        run.font.bold = True
            
            # Fill data rows
            for row_idx in range(min(rows, 20)):  # Limit to 20 rows for slide space
                for col_idx in range(cols):
                    cell = table.cell(row_idx + 1, col_idx)
                    cell_value = data.iloc[row_idx, col_idx]
                    cell.text = str(cell_value) if pd.notna(cell_value) else ""
            
            debug_print(f"DEBUG: Table added successfully with {min(rows, 20)} data rows")
            
        except Exception as e:
            debug_print(f"ERROR: Failed to add table to slide: {e}")
            traceback.print_exc()
    
    def _add_plots_to_slide(self, slide, sheet_name: str, full_sample_data: pd.DataFrame, 
                           valid_plot_options: List[str], images_to_delete: List[str], Inches):
        """Add plots to PowerPoint slide."""
        debug_print(f"DEBUG: Adding {len(valid_plot_options)} plots to slide for sheet: {sheet_name}")
        
        try:
            plt = self.get_matplotlib()
            if not plt:
                debug_print("ERROR: Matplotlib not available for PowerPoint plots")
                return
            
            # Get sample names
            sample_names = None
            if hasattr(self, 'header_data') and self.header_data and 'samples' in self.header_data:
                sample_names = [sample['id'] for sample in self.header_data['samples']]
            
            # Check for User Test Simulation
            is_user_test_simulation = sheet_name in ["User Test Simulation", "User Simulation Test"]
            num_columns_per_sample = 8 if is_user_test_simulation else 12
            
            # Position plots on the right side of the slide
            plot_x = Inches(8.78)
            plot_top = Inches(0.85)
            plot_height = Inches(1.72)
            plot_spacing = Inches(1.9)
            
            for i, plot_option in enumerate(valid_plot_options[:3]):  # Limit to 3 plots per slide
                try:
                    debug_print(f"DEBUG: Generating plot {i+1}: {plot_option}")
                    
                    if is_user_test_simulation:
                        fig, _ = processing.plot_user_test_simulation_samples(
                            full_sample_data, num_columns_per_sample, plot_option, sample_names
                        )
                    else:
                        fig, _ = processing.plot_all_samples(
                            full_sample_data, num_columns_per_sample, plot_option, sample_names
                        )
                    
                    if fig:
                        # Save plot as temporary image
                        plot_image_path = f"temp_ppt_plot_{sheet_name}_{plot_option}_{i}.png"
                        fig.savefig(plot_image_path, dpi=150, bbox_inches='tight')
                        
                        # Add image to slide
                        plot_y = plot_top + (i * plot_spacing)
                        slide.shapes.add_picture(plot_image_path, plot_x, plot_y, 
                                               Inches(2.29), plot_height)
                        
                        plt.close(fig)
                        
                        # Track for cleanup
                        if plot_image_path not in images_to_delete:
                            images_to_delete.append(plot_image_path)
                        
                        debug_print(f"DEBUG: Successfully added plot {plot_option} to slide")
                    
                except Exception as e:
                    debug_print(f"ERROR: Failed to generate plot '{plot_option}' for sheet '{sheet_name}': {e}")
                    traceback.print_exc()
            
        except Exception as e:
            debug_print(f"ERROR: Failed to add plots to slide: {e}")
            traceback.print_exc()
    
    def _add_images_to_slide(self, slide, image_paths: List[str], Inches):
        """Add images to PowerPoint slide."""
        debug_print(f"DEBUG: Adding {len(image_paths)} images to slide")
        
        try:
            Image = self.get_pil()
            if not Image:
                debug_print("ERROR: PIL not available for image processing")
                return
            
            # Calculate grid layout
            max_images_per_slide = 6
            images_to_show = image_paths[:max_images_per_slide]
            
            # Grid configuration
            cols = 3
            rows = 2
            
            # Image dimensions and spacing
            img_width = Inches(4.0)
            img_height = Inches(2.5)
            start_x = Inches(0.5)
            start_y = Inches(1.5)
            spacing_x = Inches(4.3)
            spacing_y = Inches(2.8)
            
            for i, img_path in enumerate(images_to_show):
                try:
                    if not os.path.exists(img_path):
                        debug_print(f"WARNING: Image file not found: {img_path}")
                        continue
                    
                    # Calculate position
                    row = i // cols
                    col = i % cols
                    x = start_x + (col * spacing_x)
                    y = start_y + (row * spacing_y)
                    
                    # Process image if needed (crop, resize)
                    processed_image_path = self._process_image_for_slide(img_path, Image)
                    
                    # Add image to slide
                    slide.shapes.add_picture(processed_image_path, x, y, img_width, img_height)
                    
                    # Track processed image for cleanup if it's different from original
                    if processed_image_path != img_path and processed_image_path not in self.temp_image_files:
                        self.temp_image_files.append(processed_image_path)
                    
                    debug_print(f"DEBUG: Added image {i+1} to slide: {os.path.basename(img_path)}")
                    
                except Exception as e:
                    debug_print(f"ERROR: Failed to add image {img_path} to slide: {e}")
                    continue
            
            debug_print(f"DEBUG: Successfully added {len(images_to_show)} images to slide")
            
        except Exception as e:
            debug_print(f"ERROR: Failed to add images to slide: {e}")
            traceback.print_exc()
    
    def _process_image_for_slide(self, img_path: str, Image) -> str:
        """Process image for PowerPoint slide (resize, crop if needed)."""
        try:
            # Check if image needs cropping
            should_crop = False
            if (self.gui and hasattr(self.gui, 'image_crop_states') and 
                img_path in self.gui.image_crop_states):
                should_crop = self.gui.image_crop_states[img_path]
            
            if should_crop:
                # Use image controller for processing if available
                if hasattr(self.gui, 'image_controller'):
                    processed_img = self.gui.image_controller.process_image_with_crop(img_path, should_crop)
                    if processed_img:
                        # Save processed image
                        temp_path = f"temp_processed_{os.path.basename(img_path)}"
                        processed_img.save(temp_path)
                        return temp_path
            
            # Return original path if no processing needed
            return img_path
            
        except Exception as e:
            debug_print(f"ERROR: Failed to process image for slide: {e}")
            return img_path  # Return original on error
    
    def _cleanup_images(self, images_to_delete: List[str]):
        """Clean up temporary image files."""
        debug_print(f"DEBUG: Cleaning up {len(images_to_delete)} temporary image files")
        
        try:
            # Clean up images from the provided list
            for image_path in set(images_to_delete):
                if os.path.exists(image_path):
                    try:
                        os.remove(image_path)
                        debug_print(f"DEBUG: Cleaned up: {image_path}")
                    except OSError as cleanup_error:
                        debug_print(f"WARNING: Failed to delete image {image_path}: {cleanup_error}")
            
            # Clean up temp files tracked by the controller
            for temp_file in self.temp_image_files[:]:  # Copy list to avoid modification during iteration
                try:
                    if os.path.exists(temp_file):
                        os.remove(temp_file)
                        debug_print(f"DEBUG: Cleaned up temp file: {temp_file}")
                    self.temp_image_files.remove(temp_file)
                except Exception as e:
                    debug_print(f"WARNING: Failed to remove temp file {temp_file}: {e}")
            
            debug_print("DEBUG: Image cleanup completed")
            
        except Exception as e:
            debug_print(f"ERROR: Failed to cleanup images: {e}")
    
    def get_report_status(self) -> Dict[str, Any]:
        """Get current report generation status."""
        return {
            'is_generating': self.report_model.is_generating,
            'progress': self.report_model.generation_progress,
            'report_type': self.current_report_type,
            'temp_files': len(self.temp_image_files)
        }
    
    def cancel_report_generation(self) -> bool:
        """Cancel ongoing report generation."""
        debug_print("DEBUG: Cancelling report generation")
        
        try:
            self.report_model.is_generating = False
            self.current_report_type = None
            
            # Cleanup any temporary files
            self._cleanup_images([])
            
            debug_print("DEBUG: Report generation cancelled")
            return True
            
        except Exception as e:
            debug_print(f"ERROR: Failed to cancel report generation: {e}")
            return False


class HeaderSelectorDialog:
    """Dialog for selecting and ordering report headers."""
    
    def __init__(self, parent):
        self.parent = parent
        self.result = None
        self.dialog = None
        
        # Default header order - matching requirements
        self.default_headers = [
            "Sample Name", "Media", "Viscosity", "Puffing Regime", 
            "Voltage, Resistance, Power", "Average TPM", 
            "Standard Deviation", "Normalized TPM", "Draw Pressure",
            "Usage Efficiency", "Initial Oil Mass", "Burn", "Clog", "Notes"
        ]
        
        self.header_vars = {}
        self.order_vars = {}
    
    def show(self) -> Optional[List[str]]:
        """Show dialog and return selected headers in order."""
        debug_print("DEBUG: Showing HeaderSelectorDialog")
        
        try:
            self.dialog = tk.Toplevel(self.parent)
            self.dialog.title("Select Report Headers")
            self.dialog.geometry("300x500")
            self.dialog.grab_set()
            
            # Instructions - centered
            instruction_frame = ttk.Frame(self.dialog)
            instruction_frame.pack(fill="x", pady=10)
            ttk.Label(instruction_frame, text="Select headers and set order (1=leftmost):").pack()
            
            # Header selection frame
            main_frame = ttk.Frame(self.dialog)
            main_frame.pack(fill="both", expand=True, padx=10, pady=5)
            
            # Create checkboxes and order entries for each header
            for i, header in enumerate(self.default_headers):
                frame = ttk.Frame(main_frame)
                frame.pack(fill="x", pady=2)
                
                # Checkbox
                var = tk.BooleanVar(value=True)  # Default selected
                self.header_vars[header] = var
                checkbox = ttk.Checkbutton(frame, text=header, variable=var)
                checkbox.pack(side="left")
                
                # Order entry
                order_var = tk.StringVar(value=str(i + 1))
                self.order_vars[header] = order_var
                order_entry = ttk.Entry(frame, textvariable=order_var, width=3)
                order_entry.pack(side="right")
                ttk.Label(frame, text="Order:").pack(side="right", padx=(0, 5))
            
            # Buttons
            button_frame = ttk.Frame(self.dialog)
            button_frame.pack(fill="x", pady=10)
            
            ttk.Button(button_frame, text="OK", command=self._on_ok).pack(side="right", padx=5)
            ttk.Button(button_frame, text="Cancel", command=self._on_cancel).pack(side="right")
            
            # Wait for result
            self.dialog.wait_window()
            
            debug_print(f"DEBUG: HeaderSelectorDialog result: {self.result}")
            return self.result
            
        except Exception as e:
            debug_print(f"ERROR: Failed to show header selector dialog: {e}")
            return None
    
    def _on_ok(self):
        """Handle OK button click."""
        try:
            # Get selected headers with their orders
            selected_headers = []
            header_order_pairs = []
            
            for header, var in self.header_vars.items():
                if var.get():  # If selected
                    try:
                        order = int(self.order_vars[header].get())
                        header_order_pairs.append((order, header))
                    except ValueError:
                        # Invalid order, skip this header
                        continue
            
            # Sort by order and extract headers
            header_order_pairs.sort(key=lambda x: x[0])
            selected_headers = [header for _, header in header_order_pairs]
            
            if not selected_headers:
                messagebox.showwarning("Warning", "Please select at least one header.")
                return
            
            self.result = selected_headers
            self.dialog.destroy()
            
        except Exception as e:
            debug_print(f"ERROR: Error in header selector OK: {e}")
            messagebox.showerror("Error", "An error occurred while processing headers.")
    
    def _on_cancel(self):
        """Handle Cancel button click."""
        self.result = None
        self.dialog.destroy()