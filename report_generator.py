# report_generator.py
import os
import traceback
import shutil
import pandas as pd
import math
import matplotlib.pyplot as plt
import tkinter as tk
from tkinter import ttk
from PIL import Image
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from datetime import datetime
import processing
from utils import get_save_path, plotting_sheet_test, get_plot_sheet_names, debug_print, show_success_message
from resource_utils import get_resource_path
from tkinter import messagebox  # For showing info/errors

class HeaderSelectorDialog:
    def __init__(self, parent):
        self.parent = parent
        self.result = None
        self.dialog = None

        # Default header order - matching your new requirements
        self.default_headers = [
            "Sample Name", "Media", "Viscosity", "Puffing Regime",
            "Voltage, Resistance, Power", "Average TPM",
            "Standard Deviation", "Normalized TPM", "Draw Pressure","Usage Efficiency", "Initial Oil Mass","Burn", "Clog", "Notes"
        ]

        self.header_vars = {}
        self.order_vars = {}

    def show(self):
        """Show dialog and return selected headers in order"""
        self.dialog = tk.Toplevel(self.parent)
        self.dialog.title("Select Report Headers")
        self.dialog.geometry("300x500")
        self.dialog.grab_set()

        # Instructions - centered
        instruction_frame = ttk.Frame(self.dialog)
        instruction_frame.pack(fill="x", pady=10)
        ttk.Label(instruction_frame, text="Select headers and set order (1=leftmost):").pack()

        # Main content frame for centering the scrollable area
        content_frame = ttk.Frame(self.dialog)
        content_frame.pack(fill="both", expand=True, padx=20, pady=(0, 10))

        # Scrollable frame with explicit height to prevent button overlap
        canvas = tk.Canvas(content_frame, height=350)
        scrollbar = ttk.Scrollbar(content_frame, orient="vertical", command=canvas.yview)
        scrollable_frame = ttk.Frame(canvas)

        # Header selection
        for i, header in enumerate(self.default_headers):
            frame = ttk.Frame(scrollable_frame)
            frame.pack(fill="x", padx=10, pady=2)

            # Checkbox - new headers default to unchecked
            default_checked = header not in ["Usage Efficiency", "Initial Oil Mass", "Normalized TPM"]
            var = tk.BooleanVar(value=default_checked)
            ttk.Checkbutton(frame, text=header, variable=var).pack(side="left")
            self.header_vars[header] = var

            # Order spinbox
            order_var = tk.StringVar(value=str(i+1))
            ttk.Spinbox(frame, from_=1, to=15, textvariable=order_var, width=5).pack(side="right")
            self.order_vars[header] = order_var

        scrollable_frame.bind("<Configure>", lambda e: canvas.configure(scrollregion=canvas.bbox("all")))
        canvas.create_window((0, 0), window=scrollable_frame, anchor="nw")
        canvas.configure(yscrollcommand=scrollbar.set)

        canvas.pack(side="left", fill="both", expand=True)
        scrollbar.pack(side="right", fill="y")

        # Buttons frame - at the bottom, properly separated from content
        button_frame = ttk.Frame(self.dialog)
        button_frame.pack(side="bottom", fill="x", pady=10)

        # Center the buttons using a container frame
        button_container = ttk.Frame(button_frame)
        button_container.pack()

        ttk.Button(button_container, text="Cancel", command=self.cancel).pack(side="left", padx=5)
        ttk.Button(button_container, text="OK", command=self.ok).pack(side="left", padx=5)

        # Center the dialog window on screen
        self.dialog.update_idletasks()
        width = self.dialog.winfo_width()
        height = self.dialog.winfo_height()
        x = (self.dialog.winfo_screenwidth() // 2) - (width // 2)
        y = (self.dialog.winfo_screenheight() // 2) - (height // 2)
        self.dialog.geometry(f"{width}x{height}+{x}+{y}")

        self.dialog.wait_window()
        return self.result

    def ok(self):
        # Get selected headers with their order
        selected = []
        for header in self.default_headers:
            if self.header_vars[header].get():
                try:
                    order = int(self.order_vars[header].get())
                    selected.append((header, order))
                except ValueError:
                    continue

        # Sort by order and return header names
        selected.sort(key=lambda x: x[1])
        self.result = [header for header, _ in selected]
        self.dialog.destroy()

    def cancel(self):
        self.result = None
        self.dialog.destroy()


class ReportGenerator:
    def __init__(self, gui):
        """
        Initialize the ReportGenerator with a reference to the main Tk window.
        This is used for scheduling GUI callbacks.
        """
        self.gui = gui
        self.root = gui.root

    def generate_full_report(self, filtered_sheets: dict, plot_options: list) -> None:
        """
        Generate a full report (Excel and PowerPoint) for all sheets.
        This function assumes that progress reporting is handled externally.

        Args:
            filtered_sheets (dict): A dict mapping sheet names to sheet info.
            plot_options (list): The list of plot options.
        """

        header_dialog = HeaderSelectorDialog(self.gui.root)
        selected_headers = header_dialog.show()

        if not selected_headers:
            debug_print("DEBUG: Header selection cancelled")
            raise ValueError("Header selection cancelled")

        debug_print(f"DEBUG: Selected headers for full report: {selected_headers}")

        save_path = get_save_path(".xlsx")
        if not save_path:
            raise ValueError("Save cancelled")

        ppt_save_path = save_path.replace('.xlsx', '.pptx')
        images_to_delete = []
        total_sheets = len(filtered_sheets)
        processed_count = 0

        try:
            with pd.ExcelWriter(save_path, engine='xlsxwriter') as writer:
                for sheet_name, sheet_info in filtered_sheets.items():
                    try:
                        if not isinstance(sheet_info, dict) or "data" not in sheet_info:
                            debug_print(f"DEBUG: Skipping sheet '{sheet_name}': No valid 'data' key found.")
                            continue

                        data = sheet_info["data"]  #  Extract the actual DataFrame
                        debug_print(f"DEBUG: Processing sheet: {sheet_name}")
                        debug_print(f"DEBUG: Sheet data shape: {data.shape}")

                        # Special handling for User Test Simulation
                        if sheet_name == "User Test Simulation":
                            debug_print(f"DEBUG: Processing User Test Simulation with 8-column format")
                            debug_print(f"DEBUG: Data columns: {data.columns.tolist()}")
                            debug_print(f"DEBUG: Data preview:\n{data.head()}")

                        is_plotting = plotting_sheet_test(sheet_name, data)
                        debug_print(f"DEBUG: Sheet {sheet_name} is_plotting: {is_plotting}")

                        process_function = processing.get_processing_function(sheet_name)
                        debug_print(f"DEBUG: Using processing function: {process_function.__name__}")

                        if is_plotting:
                            processed_data, _, full_sample_data = process_function(data)
                            debug_print(f"DEBUG: Processed data shape: {processed_data.shape}")
                            debug_print(f"DEBUG: Full sample data shape: {full_sample_data.shape}")
                            valid_plot_options = processing.get_valid_plot_options(plot_options, full_sample_data)
                            debug_print(f"DEBUG: Valid plot options: {valid_plot_options}")
                        else:
                            data = data.astype(str).replace([pd.NA], '')
                            processed_data, _, full_sample_data = process_function(data)
                            valid_plot_options = []

                        if processed_data.empty or full_sample_data.empty:
                            debug_print(f"DEBUG: Skipping sheet '{sheet_name}' due to empty processed data.")
                            continue

                        # Apply header selection and reordering
                        debug_print(f"DEBUG: Original processed_data columns: {processed_data.columns.tolist()}")
                        processed_data = self.reorder_processed_data(processed_data, selected_headers)
                        debug_print(f"DEBUG: Reordered processed_data columns: {processed_data.columns.tolist()}")


                        self.write_excel_report(writer, sheet_name, processed_data, full_sample_data,
                                                  valid_plot_options, images_to_delete)
                        processed_count += 1

                    except Exception as e:
                        debug_print(f"DEBUG: Error processing sheet '{sheet_name}': {e}")
                        import traceback
                        traceback.print_exc()
                        continue

            try:
                def update_ppt_prog():
                    def callback(percent):
                        total_percent = 50 + (percent * 0.5)  # PPT is second half
                        self.gui.progress_dialog.update_progress_bar(total_percent)
                        self.gui.root.update_idletasks()
                    return callback

                self.write_powerpoint_report(ppt_save_path, images_to_delete, plot_options, selected_headers, progress_callback=update_ppt_prog())

                # Final progress update
                self.gui.progress_dialog.update_progress_bar(100)
                self.gui.root.update_idletasks()

            except Exception as e:
                debug_print(f"DEBUG: Error writing PowerPoint report: {e}")
                raise

        except Exception as e:
            # Delete partial files
            for path in [save_path, ppt_save_path]:
                if path and os.path.exists(path):
                    os.remove(path)
            raise
        finally:
            self.cleanup_images(images_to_delete)

    def generate_test_report(self, selected_sheet: str, sheets: dict, plot_options: list) -> None:
        """
        Generate an Excel and PowerPoint report for only the specified sheet.

        Args:
            selected_sheet (str): The name of the sheet for the test report.
            sheets (dict): A dict mapping sheet names to raw sheet data.
            plot_options (list): List of plot options.
        """

        header_dialog = HeaderSelectorDialog(self.gui.root)
        selected_headers = header_dialog.show()

        if not selected_headers:
            debug_print("DEBUG: Header selection cancelled")
            return

        debug_print(f"DEBUG: Selected headers: {selected_headers}")

        save_path = get_save_path(".xlsx")
        if not save_path:
            return

        images_to_delete = []
        try:
            if not selected_sheet:
                messagebox.showerror("Error", "No sheet is selected for the test report.")
                return

            sheet_info = sheets.get(selected_sheet, {})
            data = sheet_info.get("data", pd.DataFrame())  # Ensure it's a DataFrame

            debug_print(f"DEBUG: Test report for {selected_sheet}")
            debug_print(f"DEBUG: Data shape: {data.shape}")

            if data is None or data.empty:
                messagebox.showwarning("Warning", f"Sheet '{selected_sheet}' is empty.")
                return

            # Special handling for User Test Simulation
            if selected_sheet == "User Test Simulation":
                debug_print(f"DEBUG: Test report processing User Test Simulation with 8-column format")
                debug_print(f"DEBUG: Data columns: {data.columns.tolist()}")

            process_function = processing.get_processing_function(selected_sheet)
            debug_print(f"DEBUG: Using processing function: {process_function.__name__}")

            processed_data, _, full_sample_data = process_function(data)
            debug_print(f"DEBUG: Test report processed data shape: {processed_data.shape}")
            debug_print(f"DEBUG: Test report full sample data shape: {full_sample_data.shape}")

            if processed_data.empty or full_sample_data.empty:
                messagebox.showwarning("Warning", f"Sheet '{selected_sheet}' did not yield valid processed data.")
                return

            processed_data = self.reorder_processed_data(processed_data, selected_headers)
            debug_print(f"DEBUG: Reordered processed data shape: {processed_data.shape}")

            with pd.ExcelWriter(save_path, engine='xlsxwriter') as writer:
                valid_plot_options = processing.get_valid_plot_options(plot_options, full_sample_data)
                debug_print(f"DEBUG: Test report valid plot options: {valid_plot_options}")
                self.write_excel_report(writer, selected_sheet, processed_data, full_sample_data,
                                          valid_plot_options, images_to_delete)

            ppt_save_path = save_path.replace('.xlsx', '.pptx')
            self.write_powerpoint_report_for_test(ppt_save_path, images_to_delete, selected_sheet,
                                                  processed_data, full_sample_data, plot_options)

            self.cleanup_images(images_to_delete)
            show_success_message("Success", f"Test report saved successfully to:\nExcel: {save_path}\nPowerPoint: {ppt_save_path}", self.gui.root)
        except Exception as e:
            debug_print(f"DEBUG: Test report generation error: {e}")
            import traceback
            traceback.print_exc()
            messagebox.showerror("Error", f"An error occurred while generating the test report: {e}")

    def reorder_processed_data(self, processed_data, selected_headers):
        """
        Reorder processed_data columns based on selected headers.

        Args:
            processed_data: DataFrame with original column order
            selected_headers: List of header names in desired order

        Returns:
            DataFrame with reordered columns
        """
        debug_print(f"DEBUG: Reordering columns. Original: {processed_data.columns.tolist()}")
        debug_print(f"DEBUG: Selected headers: {selected_headers}")

        # Create new DataFrame with selected columns in order
        reordered_data = pd.DataFrame()

        for header in selected_headers:
            if header in processed_data.columns:
                reordered_data[header] = processed_data[header]
                debug_print(f"DEBUG: Added column: {header}")
            else:
                debug_print(f"DEBUG: Missing column: {header}, adding empty")
                reordered_data[header] = [''] * len(processed_data)

        debug_print(f"DEBUG: Final reordered columns: {reordered_data.columns.tolist()}")
        return reordered_data

    def add_plots_to_slide(self, slide, sheet_name: str, full_sample_data: pd.DataFrame, valid_plot_options: list, images_to_delete: list) -> None:
        # New cascading plot layout starting under the table
        plot_start_left = Inches(0.02)
        plot_top = Inches(5.26)
        plot_height = Inches(2.0)
        current_left = plot_start_left

        debug_print(f"DEBUG: Starting cascading plot layout at left={plot_start_left}, top={plot_top}, height={plot_height}")

        numeric_data = full_sample_data.apply(pd.to_numeric, errors='coerce')
        if numeric_data.isna().all(axis=0).all():
            debug_print(f"No numeric data available for plotting in sheet '{sheet_name}'.")
            return
        sample_names = None
        if hasattr(self, 'header_data') and self.header_data and 'samples' in self.header_data:
            sample_names = [sample['id'] for sample in self.header_data['samples']]
            debug_print(f"DEBUG: Extracted sample names from header_data: {sample_names}")

        # Determine if this is User Test Simulation
        is_user_test_simulation = sheet_name in ["User Test Simulation", "User Simulation Test"]
        num_columns_per_sample = 8 if is_user_test_simulation else 12

        for i, plot_option in enumerate(valid_plot_options):
            plot_image_path = f"{sheet_name}_{plot_option}_plot.png"
            try:
                # Correct argument order
                fig, sample_names_returned = processing.plot_all_samples(numeric_data, plot_option, num_columns_per_sample, sample_names)

                # Calculate plot width based on aspect ratio to maintain proportions
                if is_user_test_simulation and hasattr(fig, 'is_split_plot') and fig.is_split_plot:
                    # For User Test Simulation split plots
                    plt.savefig(plot_image_path, dpi=150, bbox_inches='tight')
                    debug_print(f"DEBUG: Saved User Test Simulation split plot: {plot_image_path}")

                    # Wider aspect ratio for split plots
                    plot_width = Inches(3.5)  # Maintain aspect ratio for split plots
                    debug_print(f"DEBUG: Split plot {i}: positioning at left={current_left}, width={plot_width}")
                    slide.shapes.add_picture(plot_image_path, current_left, plot_top, plot_width, plot_height)
                else:
                    # Standard single plots - maintain current aspect ratio (2.29/1.72 ≈ 1.33)
                    plt.savefig(plot_image_path, dpi=150)
                    debug_print(f"DEBUG: Saved standard plot: {plot_image_path}")

                    # Calculate width to maintain aspect ratio with 2" height
                    aspect_ratio = 2.29 / 1.72  # Current aspect ratio from original code
                    plot_width = plot_height * aspect_ratio  # Maintain aspect ratio
                    debug_print(f"DEBUG: Standard plot {i}: positioning at left={current_left}, width={plot_width}")
                    slide.shapes.add_picture(plot_image_path, current_left, plot_top, plot_width, plot_height)

                # Move to next position for cascade effect
                current_left += plot_width
                debug_print(f"DEBUG: Next plot will start at left={current_left}")

                plt.close()
                if plot_image_path not in images_to_delete:
                    images_to_delete.append(plot_image_path)
            except Exception as e:
                print(f"Error generating plot '{plot_option}' for sheet '{sheet_name}': {e}")
                import traceback
                traceback.print_exc()

    def add_plots_to_slide_old(self, slide, sheet_name: str, full_sample_data: pd.DataFrame, valid_plot_options: list, images_to_delete: list) -> None:
        plot_top = Inches(1.21)
        left_column_x = Inches(8.43)
        right_column_x = Inches(10.84)
        numeric_data = full_sample_data.apply(pd.to_numeric, errors='coerce')
        if numeric_data.isna().all(axis=0).all():
            debug_print(f"No numeric data available for plotting in sheet '{sheet_name}'.")
            return
        sample_names = None
        if hasattr(self, 'header_data') and self.header_data and 'samples' in self.header_data:
            sample_names = [sample['id'] for sample in self.header_data['samples']]
            debug_print(f"DEBUG: Extracted sample names from header_data: {sample_names}")

        # Determine if this is User Test Simulation
        is_user_test_simulation = sheet_name in ["User Test Simulation", "User Simulation Test"]
        num_columns_per_sample = 8 if is_user_test_simulation else 12

        for i, plot_option in enumerate(valid_plot_options):
            plot_image_path = f"{sheet_name}_{plot_option}_plot.png"
            try:
                # FIXED: Correct argument order
                fig, sample_names_returned = processing.plot_all_samples(numeric_data, plot_option, num_columns_per_sample, sample_names)

                if is_user_test_simulation and hasattr(fig, 'is_split_plot') and fig.is_split_plot:
                    # For User Test Simulation split plots, save with adjusted size and spacing
                    plt.savefig(plot_image_path, dpi=150, bbox_inches='tight')
                    debug_print(f"DEBUG: Saved User Test Simulation split plot: {plot_image_path}")

                    # Adjust positioning for split plots (they're wider)
                    plot_x = left_column_x if i % 2 == 0 else right_column_x - Inches(1.0)  # Shift left for wider plots
                    if i % 2 != 0:
                        plot_top += Inches(2.2)  # More vertical spacing for taller plots

                    # Add with larger dimensions for split plots
                    slide.shapes.add_picture(plot_image_path, plot_x, plot_top, Inches(3.5), Inches(2.0))
                else:
                    # Standard single plots
                    plt.savefig(plot_image_path, dpi=150)
                    debug_print(f"DEBUG: Saved standard plot: {plot_image_path}")

                    plot_x = left_column_x if i % 2 == 0 else right_column_x
                    if i % 2 != 0:
                        plot_top += Inches(1.83)
                    slide.shapes.add_picture(plot_image_path, plot_x, plot_top, Inches(2.29), Inches(1.72))

                plt.close()
                if plot_image_path not in images_to_delete:
                    images_to_delete.append(plot_image_path)
            except Exception as e:
                print(f"Error generating plot '{plot_option}' for sheet '{sheet_name}': {e}")
                import traceback
                traceback.print_exc()

    def write_excel_report(self, writer, sheet_name: str, processed_data, full_sample_data, valid_plot_options=[], images_to_delete=None) -> None:
        try:
            processed_data.astype(str).replace([pd.NA], '')
            processed_data.to_excel(writer, sheet_name=sheet_name, index=False)
            plot_sheet_names = get_plot_sheet_names()
            is_plotting = sheet_name in plot_sheet_names
            if is_plotting:
                self.add_plots_to_excel(writer, sheet_name, full_sample_data, images_to_delete, valid_plot_options)
        except Exception as e:
            print(f"Error writing Excel report for sheet '{sheet_name}': {e}")
            traceback.print_exc()

    def write_powerpoint_report_for_test(self, ppt_save_path: str, images_to_delete: list, sheet_name: str, processed_data, full_sample_data, plot_options: list) -> None:
        try:
            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)

            current_file = self.gui.current_file
            image_paths = self.gui.sheet_images.get(current_file, {}).get(sheet_name, [])

            # Create main content slide
            main_slide = prs.slides.add_slide(prs.slide_layouts[6])

            # Add background and logo
            background_path = get_resource_path("resources/ccell_background.png")
            main_slide.shapes.add_picture(background_path, Inches(0), Inches(0),
                                          width=prs.slide_width, height=prs.slide_height)
            logo_path = get_resource_path("resources/ccell_logo_full.png")
            main_slide.shapes.add_picture(logo_path, Inches(11.21), Inches(0.43),
                                          width=Inches(1.57), height=Inches(0.53))

            # Add title
            title_shape = main_slide.shapes.add_textbox(Inches(0.45), Inches(-0.04),
                                                       Inches(10.72), Inches(0.64))
            text_frame = title_shape.text_frame

            text_frame.clear()

            p = text_frame.add_paragraph()
            p.text = sheet_name
            p.font.name = "Montserrat"
            p.font.size = Pt(32)
            p.font.bold = True

            # Add table/plots
            plot_sheet_names = get_plot_sheet_names()
            is_plotting = sheet_name in plot_sheet_names

            # For test reports, ALWAYS use processed data for tables
            table_width = Inches(8.07) if is_plotting else Inches(13.03)
            self.add_table_to_slide(main_slide,
                                   processed_data,  # <-- ALWAYS use processed_data for test reports
                                   table_width,
                                   is_plotting)

            # Only add plots if it's a plotting sheet AND we have valid plot options
            if is_plotting:
                valid_plot_options = processing.get_valid_plot_options(plot_options, full_sample_data)
                if valid_plot_options:
                    self.add_plots_to_slide(main_slide, sheet_name, full_sample_data,
                                           valid_plot_options, images_to_delete)

            # Create image slide if needed
            if image_paths:
                img_slide = prs.slides.add_slide(prs.slide_layouts[6])
                self.setup_image_slide(prs, img_slide, sheet_name)
                self.add_images_to_slide(img_slide, image_paths)

            processing.clean_presentation_tables(prs)
            prs.save(ppt_save_path)
            debug_print(f"PowerPoint test report saved to {ppt_save_path}")

        except Exception as e:
            print(f"Error generating test PowerPoint: {e}")
            if os.path.exists(ppt_save_path):
                os.remove(ppt_save_path)
            raise
            traceback.print_exc()

    def write_powerpoint_report(self, ppt_save_path: str, images_to_delete: list, plot_options: list, selected_headers: list = None, progress_callback = None) -> None:
        try:
            processed_slides = 0
            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)

            total_slides = 1  # Cover slide
            for sheet_name in self.gui.filtered_sheets.keys():
                total_slides += 1  # Main slide
                if sheet_name in self.gui.sheet_images and self.gui.sheet_images[self.gui.current_file][sheet_name]:
                    total_slides += 1  # Image slide

            cover_slide = prs.slides.add_slide(prs.slide_layouts[6])
            bg_path = get_resource_path("resources/Cover_Page_Logo.jpg")
            cover_slide.shapes.add_picture(bg_path, left=Inches(0), top=Inches(0),
                                             width=prs.slide_width, height=prs.slide_height)
            textbox_title = cover_slide.shapes.add_textbox(
                left=Inches((prs.slide_width.inches - 12) / 2),
                top=Inches(2.35),
                width=Inches(12),
                height=Inches(0.88)
            )
            text_frame = textbox_title.text_frame
            text_frame.margin_top = 0
            text_frame.margin_bottom = 0
            for para in list(text_frame.paragraphs):
                text_frame._element.remove(para._element)
            p = text_frame.add_paragraph()
            if self.gui.file_path:
                p.text = f"{os.path.splitext(os.path.basename(self.gui.file_path))[0]} Standard Test Report"
            else:
                p.text = "Standard Test Report"
            p.alignment = PP_ALIGN.CENTER
            p.space_before = Pt(0)
            p.space_after = Pt(0)
            text_frame.word_wrap = True
            run = p.runs[0] if p.runs else p.add_run()
            run.font.name = "Montserrat"
            run.font.size = Pt(46)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            cover_slide.shapes._spTree.remove(textbox_title._element)
            cover_slide.shapes._spTree.append(textbox_title._element)
            textbox_sub = cover_slide.shapes.add_textbox(
                left=Inches(5.73),
                top=Inches(4.05),
                width=Inches(1.87),
                height=Inches(0.37)
            )
            sub_frame = textbox_sub.text_frame
            sub_frame.margin_top = 0
            sub_frame.margin_bottom = 0
            for para in list(sub_frame.paragraphs):
                sub_frame._element.remove(para._element)
            p2 = sub_frame.add_paragraph()
            p2.text = datetime.today().strftime("%d %B %Y")
            p2.alignment = PP_ALIGN.CENTER
            p2.space_before = Pt(0)
            p2.space_after = Pt(0)
            sub_frame.word_wrap = False
            run2 = p2.runs[0] if p2.runs else p2.add_run()
            run2.font.name = "Montserrat"
            run2.font.size = Pt(16)
            run2.font.bold = True
            run2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            cover_slide.shapes._spTree.remove(textbox_sub._element)
            cover_slide.shapes._spTree.append(textbox_sub._element)
            logo_path = get_resource_path("resources/ccell_logo_full_white.png")
            cover_slide.shapes.add_picture(logo_path, left=Inches(5.7), top=Inches(6.12),
                                             width=Inches(1.93), height=Inches(0.66))

            # start setting up main slides
            background_path = get_resource_path("resources/ccell_background.png")
            logo_path = get_resource_path("resources/ccell_logo_full.png")
            plot_sheet_names = get_plot_sheet_names()

            processed_slides += 1
            self._update_ppt_progress(processed_slides, total_slides)

            for sheet_name in self.gui.filtered_sheets.keys():
                try:
                    data = self.gui.sheets.get(sheet_name)
                    if data is None or data.empty:
                        # Try to get data from filtered_sheets instead
                        sheet_info = self.gui.filtered_sheets.get(sheet_name)
                        if sheet_info and "data" in sheet_info:
                            data = sheet_info["data"]
                        else:
                            debug_print(f"Skipping sheet '{sheet_name}': No data available.")
                            continue

                    process_function = processing.get_processing_function(sheet_name)
                    processed_data, _, full_sample_data = process_function(data)

                    if processed_data.empty:
                        debug_print(f"Skipping sheet '{sheet_name}': Processed data is empty.")
                        continue

                    # Apply header selection and reordering if selected_headers is provided
                    if selected_headers:
                        debug_print(f"DEBUG: Applying header reordering for sheet '{sheet_name}'")
                        debug_print(f"DEBUG: Original columns: {processed_data.columns.tolist()}")
                        processed_data = self.reorder_processed_data(processed_data, selected_headers)
                        debug_print(f"DEBUG: Reordered columns: {processed_data.columns.tolist()}")

                    try:
                        slide_layout = prs.slide_layouts[6]
                        slide = prs.slides.add_slide(slide_layout)
                    except IndexError:
                        raise ValueError(f"Slide layout 6 not found in the PowerPoint template.")
                    slide.shapes.add_picture(background_path, left=Inches(0), top=Inches(0),
                                             width=Inches(13.33), height=Inches(7.5))
                    slide.shapes.add_picture(logo_path, left=Inches(11.21), top=Inches(0.43),
                                             width=Inches(1.57), height=Inches(0.53))

                    title_shape = slide.shapes.add_textbox(Inches(0.45), Inches(-0.04), Inches(10.72), Inches(0.64))
                    text_frame = title_shape.text_frame
                    text_frame.clear()
                    p = text_frame.add_paragraph()
                    p.text = sheet_name
                    p.alignment = PP_ALIGN.LEFT
                    run = p.runs[0]
                    run.font.name = "Montserrat"
                    run.font.size = Pt(32)
                    run.font.bold = True

                    is_plotting = sheet_name in plot_sheet_names
                    if not is_plotting:
                        if processed_data.empty:
                            debug_print(f"Skipping non-plotting sheet '{sheet_name}' due to empty data.")
                            continue
                        table_width = Inches(13.03)
                        # For non-plotting sheets, use processed_data (which now has header reordering applied)
                        self.add_table_to_slide(slide, processed_data, table_width, is_plotting)
                    else:
                        table_width = Inches(8.07)
                        # For plotting sheets, use processed_data (which now has header reordering applied)
                        self.add_table_to_slide(slide, processed_data, table_width, is_plotting)
                        valid_plot_options = processing.get_valid_plot_options(plot_options, full_sample_data)
                        if valid_plot_options:
                            self.add_plots_to_slide(slide, sheet_name, full_sample_data, valid_plot_options, images_to_delete)
                        else:
                            debug_print(f"No valid plot options for sheet '{sheet_name}'. Skipping plots.")
                    if slide.shapes.title:
                        title_shape = slide.shapes.title
                        spTree = slide.shapes._spTree
                        spTree.remove(title_shape._element)
                        spTree.append(title_shape._element)

                    processed_slides += 1

                    if progress_callback:
                        progress_callback(processed_slides, total_slides)

                    self._update_ppt_progress(processed_slides, total_slides)

                    # Add image slide if images exist
                    if self.gui.current_file in self.gui.sheet_images and sheet_name in self.gui.sheet_images.get(self.gui.current_file, {}):
                        debug_print("Images Exist! Adding a slide...")
                        current_file = self.gui.current_file
                        image_paths = self.gui.sheet_images.get(current_file, {}).get(sheet_name, [])
                        valid_image_paths = [path for path in image_paths if os.path.exists(path)]
                        img_slide = prs.slides.add_slide(prs.slide_layouts[6])
                        self.setup_image_slide(prs, img_slide, sheet_name)
                        self.add_images_to_slide(img_slide, valid_image_paths)

                        # Update progress after image slide
                        processed_slides += 1
                        self._update_ppt_progress(processed_slides, total_slides)

                except Exception as sheet_error:
                    debug_print(f"Error processing sheet '{sheet_name}': {sheet_error}")
                    processed_slides += 1
                    traceback.print_exc()
                    continue
            processing.clean_presentation_tables(prs)
            prs.save(ppt_save_path)
            debug_print(f"PowerPoint report saved successfully at {ppt_save_path}.")
        except Exception as e:
            print(f"Error writing PowerPoint report: {e}")
            traceback.print_exc()
        finally:
            for image_path in set(images_to_delete):
                if os.path.exists(image_path):
                    try:
                        os.remove(image_path)
                    except OSError as cleanup_error:
                        print(f"Error deleting image {image_path}: {cleanup_error}")

    def write_powerpoint_report_old(self, ppt_save_path: str, images_to_delete: list, plot_options: list, progress_callback = None) -> None:
        try:
            processed_slides = 0
            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)

            total_slides = 1  # Cover slide
            for sheet_name in self.gui.filtered_sheets.keys():
                total_slides += 1  # Main slide
                if sheet_name in self.gui.sheet_images and self.gui.sheet_images[self.gui.current_file][sheet_name]:
                    total_slides += 1  # Image slide


            cover_slide = prs.slides.add_slide(prs.slide_layouts[6])
            bg_path = get_resource_path("resources/Cover_Page_Logo.jpg")
            cover_slide.shapes.add_picture(bg_path, left=Inches(0), top=Inches(0),
                                             width=prs.slide_width, height=prs.slide_height)
            textbox_title = cover_slide.shapes.add_textbox(
                left=Inches((prs.slide_width.inches - 12) / 2),
                top=Inches(2.35),
                width=Inches(12),
                height=Inches(0.88)
            )
            text_frame = textbox_title.text_frame
            text_frame.margin_top = 0
            text_frame.margin_bottom = 0
            for para in list(text_frame.paragraphs):
                text_frame._element.remove(para._element)
            p = text_frame.add_paragraph()
            if self.gui.file_path:
                p.text = f"{os.path.splitext(os.path.basename(self.gui.file_path))[0]} Standard Test Report"
            else:
                p.text = "Standard Test Report"
            p.alignment = PP_ALIGN.CENTER
            p.space_before = Pt(0)
            p.space_after = Pt(0)
            text_frame.word_wrap = True
            run = p.runs[0] if p.runs else p.add_run()
            run.font.name = "Montserrat"
            run.font.size = Pt(46)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            cover_slide.shapes._spTree.remove(textbox_title._element)
            cover_slide.shapes._spTree.append(textbox_title._element)
            textbox_sub = cover_slide.shapes.add_textbox(
                left=Inches(5.73),
                top=Inches(4.05),
                width=Inches(1.87),
                height=Inches(0.37)
            )
            sub_frame = textbox_sub.text_frame
            sub_frame.margin_top = 0
            sub_frame.margin_bottom = 0
            for para in list(sub_frame.paragraphs):
                sub_frame._element.remove(para._element)
            p2 = sub_frame.add_paragraph()
            p2.text = datetime.today().strftime("%d %B %Y")
            p2.alignment = PP_ALIGN.CENTER
            p2.space_before = Pt(0)
            p2.space_after = Pt(0)
            sub_frame.word_wrap = False
            run2 = p2.runs[0] if p2.runs else p2.add_run()
            run2.font.name = "Montserrat"
            run2.font.size = Pt(16)
            run2.font.bold = True
            run2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            cover_slide.shapes._spTree.remove(textbox_sub._element)
            cover_slide.shapes._spTree.append(textbox_sub._element)
            logo_path = get_resource_path("resources/ccell_logo_full_white.png")
            cover_slide.shapes.add_picture(logo_path, left=Inches(5.7), top=Inches(6.12),
                                             width=Inches(1.93), height=Inches(0.66))

            # start setting up main slides

            background_path = get_resource_path("resources/ccell_background.png")
            logo_path = get_resource_path("resources/ccell_logo_full.png")
            plot_sheet_names = get_plot_sheet_names()

            processed_slides += 1
            self._update_ppt_progress(processed_slides, total_slides)

            for sheet_name in self.gui.filtered_sheets.keys():
                try:
                    data = self.gui.sheets.get(sheet_name)
                    if data is None or data.empty:
                        # Try to get data from filtered_sheets instead
                        sheet_info = self.gui.filtered_sheets.get(sheet_name)
                        if sheet_info and "data" in sheet_info:
                            data = sheet_info["data"]
                        else:
                            debug_print(f"Skipping sheet '{sheet_name}': No data available.")
                            continue

                    process_function = processing.get_processing_function(sheet_name)
                    processed_data, _, full_sample_data = process_function(data)

                    if processed_data.empty:
                        debug_print(f"Skipping sheet '{sheet_name}': Processed data is empty.")
                        continue
                    try:
                        slide_layout = prs.slide_layouts[6]
                        slide = prs.slides.add_slide(slide_layout)
                    except IndexError:
                        raise ValueError(f"Slide layout 5 not found in the PowerPoint template.")
                    slide.shapes.add_picture(background_path, left=Inches(0), top=Inches(0),
                                             width=Inches(13.33), height=Inches(7.5))
                    slide.shapes.add_picture(logo_path, left=Inches(11.21), top=Inches(0.43),
                                             width=Inches(1.57), height=Inches(0.53))

                    title_shape = slide.shapes.add_textbox(Inches(0.45), Inches(-0.04), Inches(10.72), Inches(0.64)) # hard coded
                    text_frame = title_shape.text_frame
                    text_frame.clear()
                    p = text_frame.add_paragraph()
                    p.text = sheet_name
                    p.alignment = PP_ALIGN.LEFT
                    run = p.runs[0]
                    run.font.name = "Montserrat"
                    run.font.size = Pt(32)
                    run.font.bold = True

                    is_plotting = sheet_name in plot_sheet_names
                    if not is_plotting:
                        if processed_data.empty:
                            debug_print(f"Skipping non-plotting sheet '{sheet_name}' due to empty data.")
                            continue
                        table_width = Inches(13.03)
                        self.add_table_to_slide(slide, full_sample_data, table_width, is_plotting)
                    else:
                        table_width = Inches(8.07)
                        self.add_table_to_slide(slide, processed_data, table_width, is_plotting)
                        valid_plot_options = processing.get_valid_plot_options(plot_options, full_sample_data)
                        if valid_plot_options:
                            self.add_plots_to_slide(slide, sheet_name, full_sample_data, valid_plot_options, images_to_delete)
                        else:
                            debug_print(f"No valid plot options for sheet '{sheet_name}'. Skipping plots.")
                    if slide.shapes.title:
                        title_shape = slide.shapes.title
                        spTree = slide.shapes._spTree
                        spTree.remove(title_shape._element)
                        spTree.append(title_shape._element)


                    processed_slides += 1

                    if progress_callback:
                        progress_callback(processed_slides, total_slides)

                    self._update_ppt_progress(processed_slides, total_slides)

                    # Add image slide if images exist

                    if self.gui.current_file in self.gui.sheet_images and sheet_name in self.gui.sheet_images.get(self.gui.current_file, {}):
                        debug_print("Images Exist! Adding a slide...")
                        current_file = self.gui.current_file
                        image_paths = self.gui.sheet_images.get(current_file, {}).get(sheet_name, [])
                        valid_image_paths = [path for path in image_paths if os.path.exists(path)]
                        img_slide = prs.slides.add_slide(prs.slide_layouts[6])
                        self.setup_image_slide(prs, img_slide, sheet_name)
                        self.add_images_to_slide(img_slide, valid_image_paths)

                        # Update progress after image slide
                        processed_slides += 1
                        self._update_ppt_progress(processed_slides, total_slides)


                except Exception as sheet_error:
                    debug_print(f"Error processing sheet '{sheet_name}': {sheet_error}")
                    processed_slides += 1
                    traceback.print_exc()
                    continue
            processing.clean_presentation_tables(prs)
            prs.save(ppt_save_path)
            debug_print(f"PowerPoint report saved successfully at {ppt_save_path}.")
        except Exception as e:
            print(f"Error writing PowerPoint report: {e}")
            traceback.print_exc()
              # Skip this sheet and continue with others
        finally:
            for image_path in set(images_to_delete):
                if os.path.exists(image_path):
                    try:
                        os.remove(image_path)
                    except OSError as cleanup_error:
                        print(f"Error deleting image {image_path}: {cleanup_error}")

    def _update_ppt_progress(self, processed: int, total: int):
        """Update progress for PowerPoint phase (40% of total)"""
        base_progress = 60  # Excel phase already completed
        ppt_progress = int((processed / total) * 40)  # 40% allocated to PPT
        self.gui.progress_dialog.update_progress_bar(base_progress + ppt_progress)
        self.gui.root.update_idletasks()

    def add_plots_to_excel(self, writer, sheet_name: str, full_sample_data, images_to_delete: list, valid_plot_options: list) -> None:
        try:
            worksheet = writer.sheets[sheet_name]
            numeric_data = full_sample_data.apply(pd.to_numeric, errors='coerce')
            sample_names = None
            if hasattr(self, 'header_data') and self.header_data and 'samples' in self.header_data:
                sample_names = [sample['id'] for sample in self.header_data['samples']]
                debug_print(f"DEBUG: Extracted sample names from header_data: {sample_names}")
            if numeric_data.isna().all().all():
                return

            # Determine if this is User Test Simulation
            is_user_test_simulation = sheet_name in ["User Test Simulation", "User Simulation Test"]
            num_columns_per_sample = int(8 if is_user_test_simulation else 12)

            for i, plot_option in enumerate(valid_plot_options):
                plot_image_path = f"{sheet_name}_{plot_option}_plot.png"
                try:
                    # Correct argument order
                    fig, sample_names_returned = processing.plot_all_samples(numeric_data, plot_option, num_columns_per_sample, sample_names)

                    if is_user_test_simulation and hasattr(fig, 'is_split_plot') and fig.is_split_plot:
                        # For User Test Simulation split plots, save with higher DPI and better format
                        plt.savefig(plot_image_path, dpi=200, bbox_inches='tight')
                        debug_print(f"DEBUG: Saved User Test Simulation split plot for Excel: {plot_image_path}")

                        # Adjust Excel positioning for wider split plots
                        col_offset = 10 + (i % 2) * 15  # More spacing for wider plots
                        row_offset = 2 + (i // 2) * 25  # More vertical spacing
                    else:
                        # Standard plots
                        plt.savefig(plot_image_path, dpi=300)
                        debug_print(f"DEBUG: Saved standard plot for Excel: {plot_image_path}")

                        col_offset = 10 + (i % 2) * 10
                        row_offset = 2 + (i // 2) * 20

                    plt.close()
                    images_to_delete.append(plot_image_path)
                    worksheet.insert_image(row_offset, col_offset, plot_image_path)

                except Exception as e:
                    print(f"Error generating plot '{plot_option}' for sheet '{sheet_name}': {e}")
                    import traceback
                    traceback.print_exc()
        except Exception as e:
            print(f"Error adding plots to Excel for sheet '{sheet_name}': {e}")
            import traceback
            traceback.print_exc()

    def add_logo_to_slide(self, slide) -> None:
        logo_path = get_resource_path('resources/ccell_logo_full.png')
        left = Inches(7.75)
        top = Inches(0.43)
        height = Inches(0.53)
        width = Inches(1.57)
        if os.path.exists(logo_path):
            slide.shapes.add_picture(logo_path, left, top, width, height)
        else:
            print(f"Logo image not found at {logo_path}")

    def cleanup_images(self, images_to_delete: list) -> None:
        for image_path in set(images_to_delete):
            try:
                os.remove(image_path)
            except FileNotFoundError:
                # The file was already removed or never created; that's fine.
                pass
            except OSError as e:
                print(f"Error deleting file {image_path}: {e}")

    def add_table_to_slide(self, slide, processed_data, table_width, is_plotting) -> bool:
        processed_data = processed_data.fillna(' ')
        processed_data = processed_data.astype(str)
        rows, cols = processed_data.shape
        if cols > 20 and rows > 30 and not is_plotting:
            debug_print(f"Skipping table creation for slide. Number of columns ({cols}) exceeds 20, number of rows ({rows}) exceeds 30, and it is not a plotting sheet.")
            return
        table_left = Inches(0.15)
        table_top = Inches(1.19)
        table_width = Inches(12.99)  # Fixed width
        debug_print(f"DEBUG: Creating table with width {table_width} and auto-adjusting height")
        # Auto-adjust height based on content - let PowerPoint handle it
        table_shape = slide.shapes.add_table(rows+1, cols, table_left, table_top, table_width,
                                              Inches(0.3 * (rows + 1))).table  # Reduced height multiplier for tighter fit
        header_font = {'name': 'Arial', 'size': Pt(10), 'bold': True}
        cell_font = {'name': 'Arial', 'size': Pt(8), 'bold': False}
        for col_idx, col_name in enumerate(processed_data.columns):
            cell = table_shape.cell(0, col_idx)
            cell.text = str(col_name)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = header_font['name']
                paragraph.font.size = header_font['size']
                paragraph.font.bold = header_font['bold']
                paragraph.alignment = PP_ALIGN.CENTER
        for row_idx in range(rows):
            for col_idx in range(cols):
                cell = table_shape.cell(row_idx + 1, col_idx)
                cell.text = processed_data.iat[row_idx, col_idx]
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.name = cell_font['name']
                    paragraph.font.size = cell_font['size']
                    paragraph.font.bold = cell_font['bold']
                    paragraph.alignment = PP_ALIGN.CENTER
        return True

    def add_images_to_slide(self, slide, image_paths: list) -> None:
        """Adds images to the slide in a grid layout."""
        start_left = Inches(0.05)
        start_top = Inches(1.2)
        available_width = Inches(13.18)  # 13.33 - 0.15
        available_height = Inches(6.3)   # 7.5 - 1.2

        num_images = len(image_paths)
        if num_images == 0:
            return

        # Calculate grid dimensions to best fit available space
        cols = math.ceil(math.sqrt(num_images * (available_width / available_height)))
        rows = math.ceil(num_images / cols)

        # Adjust to prevent excessive columns/rows
        cols = min(cols, 4)  # Max 4 columns for better layout
        rows = math.ceil(num_images / cols)

        # Calculate cell dimensions with margins
        horizontal_margin = Inches(0.2)
        vertical_margin = Inches(0.2)

        cell_width = (available_width - (cols - 1) * horizontal_margin) / cols
        cell_height = (available_height - (rows - 1) * vertical_margin) / rows

        # Position each image in grid
        for i, img_path in enumerate(image_paths):
            # Get original image dimensions
            with Image.open(img_path) as img:
                img_width, img_height = img.size
                aspect_ratio = img_width / img_height

            # Calculate display dimensions
            max_width = cell_width
            max_height = cell_height

            # Determine which dimension to constrain
            if (cell_width / cell_height) > aspect_ratio:
                # Constrain by height
                display_height = min(cell_height, max_height)
                display_width = display_height * aspect_ratio
            else:
                # Constrain by width
                display_width = min(cell_width, max_width)
                display_height = display_width / aspect_ratio

            # Center the image in the cell
            row = i // cols
            col = i % cols
            x_offset = (cell_width - display_width) / 2
            y_offset = (cell_height - display_height) / 2

            left = start_left + col * (cell_width + horizontal_margin) + x_offset
            top = start_top + row * (cell_height + vertical_margin) + y_offset

            slide.shapes.add_picture(img_path, left=left, top=top, width=display_width, height = display_height)

    def setup_image_slide(self, prs, slide, sheet_name):
        """Sets up the background, logo, and title for an image slide."""
        # Add background
        background_path = get_resource_path("resources/ccell_background.png")
        slide.shapes.add_picture(background_path, Inches(0), Inches(0),
                                 width=prs.slide_width, height=prs.slide_height)

        # Add logo
        logo_path = get_resource_path("resources/ccell_logo_full.png")
        slide.shapes.add_picture(logo_path, Inches(11.21), Inches(0.43),
                                 width=Inches(1.57), height=Inches(0.53))

        # Add title
        title_shape = slide.shapes.add_textbox(Inches(0.45), Inches(-0.04), Inches(10.72), Inches(0.64))
        text_frame = title_shape.text_frame
        text_frame.clear()
        p = text_frame.add_paragraph()
        p.text = sheet_name
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.name = "Montserrat"
        run.font.size = Pt(32)
        run.font.bold = True

    def clean_presentation_tables(presentation):
        """
        Clean all tables in the given PowerPoint presentation by:
        - Removing rows where all cells are empty.
        - Removing columns where all cells, including the header, are empty.

        Args:
            presentation (pptx.Presentation): The PowerPoint presentation to process.
        """
        for slide in presentation.slides:
            for shape in slide.shapes:
                if not shape.has_table:
                    continue  # Skip if the shape is not a table

                table = shape.table

                # Get the number of rows and columns
                num_rows = len(table.rows)
                num_cols = len(table.columns)

                # Step 1: Remove empty rows
                rows_to_keep = []
                for row_idx in range(num_rows):
                    is_empty_row = all(
                        not cell.text.strip() for cell in table.rows[row_idx].cells
                    )
                    if not is_empty_row:
                        rows_to_keep.append(row_idx)

                # Keep only the non-empty rows
                for row_idx in reversed(range(num_rows)):
                    if row_idx not in rows_to_keep:
                        table._tbl.remove(table.rows[row_idx]._tr)  # Remove row from XML

                # Step 2: Remove empty columns
                num_rows = len(table.rows)  # Updated number of rows after cleaning
                cols_to_keep = []
                for col_idx in range(num_cols):
                    # Check if all cells in the column (including the header) are empty or 'nan'
                    is_empty_col = all(
                        not table.cell(row_idx, col_idx).text.strip() or table.cell(row_idx, col_idx).text.strip() == "nan"
                        for row_idx in range(num_rows)
                    )
                    if not is_empty_col:
                        cols_to_keep.append(col_idx)

                # Keep only the non-empty columns
                for col_idx in reversed(range(num_cols)):
                    if col_idx not in cols_to_keep:
                        for row in table.rows:
                            row._tr.remove(row.cells[col_idx]._tc)  # Remove column cell from XML