"""
Export Manager
Handles export operations including plot images, table images, and PowerPoint reports
"""

import tkinter as tk
from tkinter import messagebox, filedialog
import matplotlib.pyplot as plt
import os
import tempfile
from datetime import datetime
from utils import debug_print, show_success_message


class ExportManager:
    """Manages export operations for sensory data."""
    
    def __init__(self, sensory_window, sample_manager, plot_manager):
        """Initialize the export manager with reference to main window."""
        self.sensory_window = sensory_window
        self.sample_manager = sample_manager
        self.plot_manager = plot_manager
        
    def save_plot_as_image(self):
        """Save the current spider plot as an image file."""
        debug_print("DEBUG: Starting plot image save")

        if not self.sensory_window.samples:
            messagebox.showwarning("Warning", "No samples to save! Please add samples first.")
            return

        try:
            filename = filedialog.asksaveasfilename(
                title="Save Plot as Image",
                defaultextension=".png",
                filetypes=[
                    ("PNG files", "*.png"),
                    ("JPG files", "*.jpg"),
                    ("PDF files", "*.pdf"),
                    ("SVG files", "*.svg"),
                    ("All files", "*.*")
                ]
            )

            if filename:
                # Ensure we have the latest plot
                self.plot_manager.update_plot()

                # Save the figure with high DPI for quality
                self.sensory_window.fig.savefig(filename, dpi=300, bbox_inches='tight',
                               facecolor='white', edgecolor='none')

                debug_print(f"DEBUG: Plot saved successfully to {filename}")
                show_success_message("Success", f"Plot saved successfully as {os.path.basename(filename)}", self.sensory_window.window)

        except Exception as e:
            debug_print(f"DEBUG: Error saving plot: {e}")
            messagebox.showerror("Error", f"Failed to save plot: {str(e)}")

    def save_table_as_image(self):
        """Save the sensory data table as an image."""
        debug_print("DEBUG: Starting table image save with comments")

        if not self.sensory_window.samples:
            messagebox.showwarning("Warning", "No data to save! Please add samples first.")
            return

        try:
            filename = filedialog.asksaveasfilename(
                title="Save Table as Image",
                defaultextension=".png",
                filetypes=[
                    ("PNG files", "*.png"),
                    ("JPG files", "*.jpg"),
                    ("PDF files", "*.pdf"),
                    ("All files", "*.*")
                ]
            )

            if filename:
                # Create table data with attributes as headers and samples as rows
                table_data = []

                # Header row with attributes + comments
                headers = ["Sample"] + self.sensory_window.metrics + ["Additional Comments"]
                table_data.append(headers)

                # Data rows - one per sample
                for sample_name, sample_data in self.sensory_window.samples.items():
                    row = [sample_name]
                    for metric in self.sensory_window.metrics:
                        row.append(str(sample_data.get(metric, "N/A")))
                    # Add comments - get from sample data or leave blank if empty
                    comments = sample_data.get("comments", "").strip()
                    row.append(comments if comments else "")
                    table_data.append(row)

                # Create figure for table with wider width to accommodate comments
                fig, ax = plt.subplots(figsize=(16, max(6, len(self.sensory_window.samples) * 0.5)))
                ax.axis('tight')
                ax.axis('off')

                # Create table
                table = ax.table(cellText=table_data[1:], colLabels=table_data[0],
                               cellLoc='center', loc='center')
                table.auto_set_font_size(False)
                table.set_fontsize(9)
                table.scale(1.4, 2.2)

                # Style the table
                for i in range(len(headers)):
                    table[(0, i)].set_facecolor('#4CAF50')
                    table[(0, i)].set_text_props(weight='bold', color='white')

                comments_col_idx = len(headers) - 1
                for row_idx in range(len(table_data)):
                    if row_idx == 0:  # Header
                        continue
                    cell = table[(row_idx, comments_col_idx)]
                    cell.set_width(0.3)
                    cell.set_text_props(ha='left', va='top', wrap=True)

                # Add title
                ax.set_title("Sensory Evaluation Results", fontsize=16, fontweight='bold', pad=20)

                # Save with high quality
                fig.savefig(filename, dpi=300, bbox_inches='tight',
                           facecolor='white', edgecolor='none')
                plt.close(fig)

                debug_print(f"DEBUG: Table with comments saved successfully to {filename}")
                show_success_message("Success", f"Table saved successfully as {os.path.basename(filename)}", self.sensory_window.window)

        except Exception as e:
            debug_print(f"DEBUG: Error saving table: {e}")
            import traceback
            traceback.print_exc()
            messagebox.showerror("Error", f"Failed to save table: {str(e)}")

    def generate_powerpoint_report(self):
        """Generate a PowerPoint report using the same template as generate_test_report."""
        debug_print("DEBUG: Starting PowerPoint report generation")

        if not self.sensory_window.samples:
            messagebox.showwarning("Warning", "No data to export! Please add samples first.")
            return

        try:
            # Get save location
            filename = filedialog.asksaveasfilename(
                title="Save PowerPoint Report",
                defaultextension=".pptx",
                filetypes=[("PowerPoint files", "*.pptx"), ("All files", "*.*")]
            )

            if not filename:
                return

            debug_print(f"DEBUG: Creating PowerPoint report at {filename}")

            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor
            from datetime import datetime
            import tempfile

            # Create presentation with same template structure as existing reports
            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)

            # Create main content slide
            main_slide = prs.slides.add_slide(prs.slide_layouts[6])

            # Add background and logo using existing template structure
            from resource_utils import get_resource_path

            background_path = get_resource_path("resources/ccell_background.png")
            if os.path.exists(background_path):
                main_slide.shapes.add_picture(background_path, Inches(0), Inches(0),
                                            width=prs.slide_width, height=prs.slide_height)
                debug_print("DEBUG: Background added successfully")
            else:
                debug_print("DEBUG: Background not found, using plain slide")

            logo_path = get_resource_path("resources/ccell_logo_full.png")
            if os.path.exists(logo_path):
                main_slide.shapes.add_picture(logo_path, Inches(11.21), Inches(0.43),
                                            width=Inches(1.57), height=Inches(0.53))
                debug_print("DEBUG: Logo added successfully")

            # Add title
            title_shape = main_slide.shapes.add_textbox(Inches(0.45), Inches(-0.04),
                                                       Inches(10.72), Inches(0.64))
            text_frame = title_shape.text_frame
            text_frame.clear()

            p = text_frame.add_paragraph()
            p.text = "Sensory Evaluation Report"
            p.font.name = "Montserrat"
            p.font.size = Pt(32)
            p.font.bold = True

            # Create table data with proper structure (attributes as headers + comments)
            table_data = []
            headers = ["Sample"] + self.sensory_window.metrics + ["Additional Comments"]

            # Add header data if available
            header_info = []
            for field in self.sensory_window.header_fields:
                if field in self.sensory_window.header_vars and self.sensory_window.header_vars[field].get():
                    header_info.append(f"{field}: {self.sensory_window.header_vars[field].get()}")

            # Add data rows with current comments (including any just typed)
            for sample_name, sample_data in self.sensory_window.samples.items():
                row = [sample_name]
                for metric in self.sensory_window.metrics:
                    row.append(str(sample_data.get(metric, "N/A")))
                # Include current comments
                comments = sample_data.get("comments", "").strip()
                row.append(comments if comments else "")
                table_data.append(row)

            if table_data:
                table_shape = main_slide.shapes.add_table(
                    len(table_data) + 1, len(headers),  # +1 for header row
                    Inches(0.45), Inches(1.5),
                    Inches(6.5), Inches(4.5)
                )
                table = table_shape.table

                # Set header row
                for col_idx, header in enumerate(headers):
                    cell = table.cell(0, col_idx)
                    cell.text = header
                    cell.text_frame.paragraphs[0].font.bold = True
                    cell.text_frame.paragraphs[0].font.size = Pt(10)

                # Set data rows
                for row_idx, row_data in enumerate(table_data, 1):
                    for col_idx, cell_value in enumerate(row_data):
                        cell = table.cell(row_idx, col_idx)
                        cell.text = str(cell_value)

                        # Special formatting for comments column
                        if col_idx == len(headers) - 1:
                            cell.text_frame.paragraphs[0].font.size = Pt(8)
                            # Set text alignment for comments
                            for paragraph in cell.text_frame.paragraphs:
                                paragraph.alignment = PP_ALIGN.LEFT
                        else:
                            cell.text_frame.paragraphs[0].font.size = Pt(9)

                debug_print(f"DEBUG: Table with comments added - {len(table_data)} rows and {len(headers)} columns")

            # Create plot
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
                plot_image_path = tmp_file.name

            try:
                # Create a special figure for PowerPoint with better legend handling
                fig_ppt, ax_ppt = plt.subplots(figsize=(8, 6), subplot_kw=dict(projection='polar'))
                fig_ppt.patch.set_facecolor('white')

                # Get selected samples for plotting
                selected_samples = []
                if hasattr(self.sample_manager, 'sample_checkboxes'):
                    for sample_name, checkbox_var in self.sample_manager.sample_checkboxes.items():
                        if checkbox_var.get() and sample_name in self.sensory_window.samples:
                            selected_samples.append(sample_name)

                # If no samples selected, select all
                if not selected_samples:
                    selected_samples = list(self.sensory_window.samples.keys())

                if selected_samples:
                    # Setup the spider plot
                    num_metrics = len(self.sensory_window.metrics)
                    angles = np.linspace(0, 2 * np.pi, num_metrics, endpoint=False).tolist()
                    angles += angles[:1]  # Complete the circle

                    # Set up the plot
                    ax_ppt.set_theta_offset(np.pi / 2)
                    ax_ppt.set_theta_direction(-1)
                    ax_ppt.set_thetagrids(np.degrees(angles[:-1]), self.sensory_window.metrics, fontsize=10)
                    ax_ppt.set_ylim(0, 9)
                    ax_ppt.set_yticks(range(1, 10))
                    ax_ppt.set_yticklabels(range(1, 10))
                    ax_ppt.grid(True, alpha=0.3)

                    # Colors for different samples
                    colors = ['red', 'green', 'blue', 'orange', 'purple', 'brown', 'pink', 'cyan', 'magenta', 'lime']
                    line_styles = ['-', '--', '-.', ':']

                    # Plot each selected sample
                    for i, sample_name in enumerate(selected_samples):
                        sample_data = self.sensory_window.samples[sample_name]
                        values = [sample_data.get(metric, 5) for metric in self.sensory_window.metrics]
                        values += values[:1]  # Complete the circle

                        color = colors[i % len(colors)]
                        line_style = line_styles[i % len(line_styles)]

                        # Plot the line and markers
                        ax_ppt.plot(angles, values, 'o', linewidth=2.5, label=sample_name,
                                   color=color, linestyle=line_style, markersize=8, alpha=0.8)
                        # Fill the area
                        ax_ppt.fill(angles, values, alpha=0.1, color=color)

                    ax_ppt.legend(loc='upper right', bbox_to_anchor=(1.1, 1.1), fontsize=9)

                    # Set title
                    ax_ppt.set_title('Sensory Profile Comparison', fontsize=12, fontweight='bold', pad=15)

                # Save the PowerPoint-specific plot
                fig_ppt.savefig(plot_image_path, dpi=300, bbox_inches='tight',
                               facecolor='white', edgecolor='none')
                plt.close(fig_ppt)

                main_slide.shapes.add_picture(plot_image_path,
                                            Inches(7.2), Inches(1.5),
                                            Inches(5.8), Inches(4.5))
                debug_print("DEBUG: Plot with proper legend positioning added to PowerPoint slide")

            finally:
                # Clean up temporary file
                if os.path.exists(plot_image_path):
                    os.remove(plot_image_path)

            # Add header information as text box if available
            if header_info:
                info_shape = main_slide.shapes.add_textbox(Inches(0.45), Inches(6.2),
                                                         Inches(12.0), Inches(1.0))
                info_frame = info_shape.text_frame
                info_frame.clear()

                p = info_frame.add_paragraph()
                p.text = " | ".join(header_info)
                p.font.size = Pt(10)
                p.font.name = "Montserrat"

            # Save the presentation
            prs.save(filename)
            debug_print(f"DEBUG: PowerPoint saved successfully to {filename}")
            show_success_message("Success", f"PowerPoint report saved successfully as {os.path.basename(filename)}", self.sensory_window.window)

        except Exception as e:
            debug_print(f"DEBUG: Error generating PowerPoint report: {e}")
            import traceback
            traceback.print_exc()
            messagebox.showerror("Error", f"Failed to generate PowerPoint report: {str(e)}")

