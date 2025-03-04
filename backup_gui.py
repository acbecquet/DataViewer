"""
GUI module for the DataViewer Application. Developed by Charlie Becquet.
Provides a Tkinter-based interface for interacting with Excel data, generating reports,
and plotting graphs.
"""
import os
import time
import threading
import copy
import shutil
from typing import Dict, Optional, List, Any
from datetime import datetime
import matplotlib
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import psutil
import tkinter as tk
import backup_processing as processing
from tkinter import ttk, filedialog, messagebox, Toplevel, Label, Button
from tkintertable import TableCanvas, TableModel
from matplotlib.backends.backend_tkagg import FigureCanvasTkAgg, NavigationToolbar2Tk
from matplotlib.widgets import CheckButtons
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from utils import (
    get_resource_path,
    autofit_columns_in_excel,
    clean_columns,
    wrap_text,
    round_values,
    get_save_path,
    is_standard_file,
    plotting_sheet_test
)
matplotlib.use('TkAgg')  # Ensure Matplotlib uses TkAgg backend

APP_BACKGROUND_COLOR = '#0504AA'
BUTTON_COLOR = '#4169E1'
FONT = ('Arial', 12)
PLOT_CHECKBOX_TITLE = "Click Checkbox to \nAdd/Remove Item \nFrom Plot"



class TestingGUI:
    """Main GUI class for the Standardized Testing application."""

    def __init__(self, root):
        self.report_thread = None # initialize report thread
        self.root = root
        self.root.title("Standardized Testing GUI")
        self.initialize_variables()
        self.configure_ui()
        self.show_startup_menu()
        # Bind the window close event
        self.root.protocol("WM_DELETE_WINDOW", self.on_app_close)

    # Initialization and Configuration
    def initialize_variables(self) -> None:
        """Initialize variables used throughout the GUI."""
        self.sheets: Dict[str, pd.DataFrame] = {}
        self.filtered_sheets: Dict[str, pd.DataFrame] = {}
        self.all_filtered_sheets: List[Dict[str, Any]] = []
        self.current_file: Optional[str] = None
        self.selected_sheet = tk.StringVar()
        self.selected_plot_type = tk.StringVar(value="TPM")
        self.file_path: Optional[str] = None
        self.plot_dropdown = None
        self.threads = [] # track active threads
        self.canvas = None
        self.figure = None
        self.axes = None
        self.lines = None
        self.lock = threading.Lock()
        self.line_labels = []
        self.original_lines_data = []
        self.checkbox_cid = None
        self.plot_options = ["TPM", "Draw Pressure", "Resistance", "Power Efficiency", "TPM (Bar)"]
        self.check_buttons = None

    def setup_ui(self, hide_alerts: bool = False) -> None:
        """
        Set up the main user interface.

        Args:
            hide_alerts (bool): Suppress alert messages for missing data.
        """
        self.root.title(
            f"Loaded: {os.path.basename(self.file_path)}"
            if self.file_path
            else "No File Loaded"
        )

        # Add buttons and dropdowns to the top frame
        if not hasattr(self, 'buttons_added') or not self.buttons_added:
            self.add_report_buttons(self.top_frame)  # Add buttons
            self.add_trend_analysis_button(self.top_frame)  # Add Trend Analysis button
            self.buttons_added = True  # Mark as added to prevent duplicates
        
 
        self.add_or_update_file_dropdown()  

        if self.file_path and self.filtered_sheets:
            self.update_ui_for_current_file()           
        elif not hide_alerts:
            
            messagebox.showinfo(
                "Info", "No valid file loaded. Please use the Add Data button to select a file."
            )

    def show_startup_menu(self) -> None:
        """Display a startup menu with 'New' and 'Load' options."""
        startup_menu = Toplevel(self.root)
        startup_menu.title("Welcome")
        startup_menu.geometry("300x150")
        startup_menu.transient(self.root)
        startup_menu.grab_set()  # Make the window modal

        label = ttk.Label(startup_menu, text="Welcome to DataViewer by SDR!", font=FONT,background=APP_BACKGROUND_COLOR, foreground="white")
        label.pack(pady=10)

        # New Button
        new_button = ttk.Button(
            startup_menu,
            text="New",
            command=lambda: self.create_new_template(startup_menu)
        )
        new_button.pack(pady=5)

        # Load Button
        load_button = ttk.Button(
            startup_menu,
            text="Load",
            command=lambda: self.start_file_loading(startup_menu)
        )
        load_button.pack(pady=5)

        self.center_window(startup_menu)

    def start_file_loading(self, startup_menu: Toplevel) -> None:
        """Handle the 'Load' button click in the startup menu."""
        startup_menu.destroy()
        self.load_initial_file()

    def create_new_template(self, startup_menu: Toplevel) -> None:
        """
        Handle the 'New' button click to create a new template file.
        Args:
            startup_menu (Toplevel): The startup menu window.
        """
        try:
            # Destroy the startup menu
            startup_menu.destroy()

            # Get the path to the template file
            template_path = get_resource_path("resources/Standardized Test Template - LATEST VERSION - 2025 Jan.xlsx")
            if not os.path.exists(template_path):
                messagebox.showerror("Error", "Template file not found. Please check the resources folder.")
                return

            # Ask the user where to save the new file
            new_file_path = filedialog.asksaveasfilename(
                title="Save New Test File As",
                defaultextension=".xlsx",
                filetypes=[("Excel files", "*.xlsx")]
            )
            if not new_file_path:
                return  # User canceled the save dialog

            # Copy the template file to the new location
            shutil.copy(template_path, new_file_path)

            # Load the new file into the application
            self.file_path = new_file_path
            self.load_excel_file(new_file_path)

            # Clear and reinitialize frames
            self.clear_display_frame(is_plotting_sheet=False, is_empty_sheet=False)

            # Add the new file to the internal tracking structure
            self.all_filtered_sheets.append({
                "file_name": os.path.basename(new_file_path),
                "file_path": new_file_path,
                "filtered_sheets": copy.deepcopy(self.filtered_sheets)
            })

            # Update the UI
            self.set_active_file(os.path.basename(new_file_path))
            self.setup_ui()

            messagebox.showinfo("Success", "New template created and loaded successfully.")
        except Exception as e:
            messagebox.showerror("Error", f"An error occurred while creating a new template: {e}")

    def center_window(self, window: tk.Toplevel, width: Optional[int] = None, height: Optional[int] = None) -> None:
        """
        Center a given Tkinter window on the screen.
    
        Args:
            window (tk.Toplevel): The Tkinter window to be centered.
            width (Optional[int]): Desired width of the window. If not provided, the current width is used.
            height (Optional[int]): Desired height of the window. If not provided, the current height is used.
        """
        
        window.update_idletasks()  # Ensure all geometry calculations are updated

        # Get the current dimensions of the window if width/height are not specified
        window_width = width or window.winfo_width()
        window_height = height or window.winfo_height()

        # Get the screen dimensions
        screen_width = window.winfo_screenwidth()
        screen_height = window.winfo_screenheight()

        # Calculate the position for centering
        position_x = (screen_width - window_width) // 2
        position_y = (screen_height - window_height) // 2

        # Set the geometry with the calculated position
        window.geometry(f"{window_width}x{window_height}+{position_x}+{position_y}")

    def configure_ui(self) -> None:
        """Configure the UI appearance and set application properties."""
        icon_path = get_resource_path('resources/ccell_icon.png')
        self.root.iconphoto(False, tk.PhotoImage(file=icon_path))
        self.set_app_colors()
        self.set_window_size(0.75, 0.75)
        self.center_window(self.root)

        self.setup_frames()

        self.add_menu()
        self.add_or_update_file_dropdown()
    
    def add_or_update_file_dropdown(self) -> None:
        """Add a file selection dropdown or update its values if it already exists."""
        if not hasattr(self, 'file_dropdown') or not self.file_dropdown:
            # Create the dropdown for selecting files
            dropdown_frame = ttk.Frame(self.top_frame, width = 400, height = 50)
            dropdown_frame.pack(side="left", pady=2, padx=5)

            file_label = ttk.Label(dropdown_frame, text="Select File:", font=FONT, foreground="white", background = APP_BACKGROUND_COLOR)
            file_label.pack(side="left", padx=(0, 5))

            self.file_dropdown_var = tk.StringVar()
            self.file_dropdown = ttk.Combobox(
                dropdown_frame,
                textvariable=self.file_dropdown_var,
                state="readonly",
                font=FONT,
                width = 10
            )
            self.file_dropdown.pack(side="left", fill="x", expand=True, padx=(5, 10))

            # Bind selection event to update the displayed sheet
            self.file_dropdown.bind("<<ComboboxSelected>>", self.on_file_selection)

        # Update the dropdown values with the loaded files
        self.update_file_dropdown()

    def add_or_update_plot_dropdown(self) -> None:
        """Add a file selection dropdown or update its values if it already exists."""
        if not hasattr(self, 'plot_dropdown') or not self.plot_dropdown:
            # Create the dropdown for selecting plots
            dropdown_frame = ttk.Frame(self.plot_frame, width = 400, height = 50)
            dropdown_frame.pack(side="right", pady=2, padx=5)

            file_label = ttk.Label(dropdown_frame, text="Select File:", font=FONT, foreground="white", background = APP_BACKGROUND_COLOR)
            file_label.pack(side="left", padx=(0, 5))

            self.file_dropdown_var = tk.StringVar()
            self.file_dropdown = ttk.Combobox(
                dropdown_frame,
                textvariable=self.file_dropdown_var,
                state="readonly",
                font=FONT,
                width = 10
            )
            self.file_dropdown.pack(side="left", fill="x", expand=True, padx=(5, 10))

            # Bind selection event to update the displayed sheet
            self.file_dropdown.bind("<<ComboboxSelected>>", self.on_file_selection)

        # Update the dropdown values with the loaded files
        self.update_file_dropdown()

    def update_file_dropdown(self) -> None:
        """Update the file dropdown with loaded file names."""
        file_names = [file_data["file_name"] for file_data in self.all_filtered_sheets]
        self.file_dropdown["values"] = file_names
        if file_names:
            self.file_dropdown_var.set(file_names[-1])
            self.file_dropdown.update_idletasks()  # Force immediate UI refresh

    def populate_or_update_sheet_dropdown(self) -> None:
        """Populate or update the dropdown for sheet selection."""
        if not hasattr(self, 'drop_down_menu') or not self.drop_down_menu:
            # Create a dropdown for sheet selection
            self.drop_down_menu = ttk.Combobox(
                self.top_frame,
                textvariable=self.selected_sheet,
                state="readonly",
                font=FONT
            )
            self.drop_down_menu.pack(anchor = 'n',pady=(10, 10))

            # Bind selection event to update the displayed sheet
            self.drop_down_menu.bind(
                "<<ComboboxSelected>>",
                lambda event: self.update_displayed_sheet(self.selected_sheet.get())
            )
        self.drop_down_menu
        # Update dropdown values with sheet names
        all_sheet_names = list(self.filtered_sheets.keys())
        self.drop_down_menu["values"] = all_sheet_names

        current_selection = self.selected_sheet.get()
        if current_selection in all_sheet_names:
            self.selected_sheet.set(current_selection)
        elif all_sheet_names:
            # if the current selection is invalid, default to the first sheet
            self.selected_sheet.set(all_sheet_names[0])

    def build_aggregate_plots_dict(self) -> dict:
        """
        Iterate over filtered_sheets, and for each plotting sheet, compute the aggregate metrics,
        generate the aggregate plot figure, and store it in a dictionary.
        Returns a dictionary mapping sheet names to their corresponding Matplotlib figure.
        """
        aggregate_plots = {}
        # Get list of plotting sheet names from processing module
        plot_sheet_names = processing.get_plot_sheet_names()
        for sheet_name, sheet_info in self.filtered_sheets.items():
            # Process only plotting sheets
            if sheet_name not in plot_sheet_names:
                continue
            full_sample_data = sheet_info["data"]
            if full_sample_data.empty:
                continue
            # Get aggregate metrics for this sheet using the aggregate function.
            agg_df = processing.aggregate_sheet_metrics(full_sample_data, num_columns_per_sample=12)
            # Optionally, add a column for the sheet name.
            agg_df["Sheet"] = sheet_name
            # Generate the aggregate plot figure.
            fig = processing.plot_aggregate_trends(agg_df)
            aggregate_plots[sheet_name] = fig
        return aggregate_plots
 
    def collect_trend_data_aggregate(self, is_plotting_sheet: bool = True) -> pd.DataFrame:
        """
        Collect aggregate metrics only from plotting sheets (if is_plotting_sheet is True) for trend analysis.
        For each sheet in filtered_sheets that qualifies as a plotting sheet, this function calls the aggregation function
        and adds a 'Sheet' column to indicate which sheet the sample comes from.
    
        Returns:
            pd.DataFrame: DataFrame with columns: "Sheet", "Sample Name", "Average TPM", "Total Puffs".
        """
        aggregate_data = []
        # Get the list of plotting sheet names
        plot_sheet_names = processing.get_plot_sheet_names()
    
        for sheet_name, sheet_info in self.filtered_sheets.items():
            # If is_plotting_sheet is True, process only those sheets that are plotting sheets
            if is_plotting_sheet and sheet_name not in plot_sheet_names:
                continue
        
            full_sample_data = sheet_info["data"]
            if not full_sample_data.empty:
                agg_df = processing.aggregate_sheet_metrics(full_sample_data, num_columns_per_sample=12)
                agg_df["Sheet"] = sheet_name
                aggregate_data.append(agg_df)
    
        if aggregate_data:
            return pd.concat(aggregate_data, ignore_index=True)
        else:
            return pd.DataFrame()

    def on_file_selection(self, event) -> None:
        """Handle file selection from the dropdown."""
        selected_file = self.file_dropdown_var.get()
        if not selected_file or selected_file == self.current_file:
            return

        try:
            self.set_active_file(selected_file)
            self.update_ui_for_current_file()
        except ValueError as e:
            messagebox.showerror("Error", str(e))

    def add_trend_analysis_button(self, parent_frame: ttk.Frame) -> None:
        """Add a 'Trend Analysis' button to the UI."""
        trend_button = ttk.Button(
            parent_frame,
            text="Trend Analysis", state = "enabled",
            command=self.open_trend_analysis_window
        )
        trend_button.pack(side="right", padx=(5, 10), pady=(5, 5))  # disabled for now

    def open_trend_analysis_window(self) -> None:
        """
        Open a new window for trend analysis. This window creates a dropdown populated with the
        names of plotting sheets (from filtered_sheets) and a dedicated frame to display the
        corresponding aggregate plot. When the user selects a different sheet from the dropdown,
        the plot is regenerated and updated.
        """
        trend_window = Toplevel(self.root)
        trend_window.title("Trend Analysis")
        trend_window.geometry("800x600")
        self.center_window(trend_window)
    
        # Create a frame at the top for the dropdown menu
        dropdown_frame = ttk.Frame(trend_window)
        dropdown_frame.pack(pady=10)
    
        # Get a list of plotting sheet names from filtered_sheets.
        # (Assuming your get_plot_sheet_names() returns the valid names.)
        valid_plot_sheet_names = [sheet_name for sheet_name in self.filtered_sheets.keys() 
                                  if sheet_name in processing.get_plot_sheet_names()]
    
        if not valid_plot_sheet_names:
            messagebox.showerror("Error", "No valid plotting sheets found for trend analysis.")
            trend_window.destroy()
            return

        # Set up a StringVar and dropdown for selecting aggregate plots.
        self.selected_agg_plot = tk.StringVar(value=valid_plot_sheet_names[0])
        dropdown = ttk.Combobox(dropdown_frame, textvariable=self.selected_agg_plot,
                                values=valid_plot_sheet_names, state="readonly", font=FONT)
        dropdown.pack()
        # Bind selection changes to update the plot.
        dropdown.bind("<<ComboboxSelected>>", lambda event: self.update_aggregate_plot(plot_frame))
    
        # Create a frame to hold the aggregate plot.
        plot_frame = ttk.Frame(trend_window)
        plot_frame.pack(fill="both", expand=True)
    
        # Initially, generate and display the aggregate plot for the first valid plotting sheet.
        self.update_aggregate_plot(plot_frame)

    def display_aggregate_plot(self, parent, fig: plt.Figure) -> None:
        """
        Given a sheet name (from aggregate_plots), embed its corresponding figure in the parent window.
        """
        # Clear previous widgets in parent if needed.
        for widget in parent.winfo_children():
            # Optionally, you might want to skip the dropdown widget by checking its type.
            if isinstance(widget, ttk.Combobox):
                continue
            widget.destroy()

        
        canvas = FigureCanvasTkAgg(fig, master=parent)
        canvas.draw()
        canvas.get_tk_widget().pack(fill="both", expand=True)

        # Add a toolbar if desired.
        toolbar = NavigationToolbar2Tk(canvas, parent)
        toolbar.update()
        toolbar.pack(side="bottom", fill="x")

    def update_aggregate_plot(self, parent) -> None:
        """
        Update the displayed aggregate plot based on the selected dropdown value.
        Regenerate the aggregate plot figure from scratch for the selected sheet.
        """
        selected_sheet = self.selected_agg_plot.get()
        

        # Retrieve the corresponding sheet's full sample data
        sheet_info = self.filtered_sheets.get(selected_sheet)
        if not sheet_info:
            messagebox.showerror("Error", f"Sheet '{selected_sheet}' not found.")
            return

        full_sample_data = sheet_info["data"]
        print(f"Selected sheet: {selected_sheet}, shape: {full_sample_data.shape}")
        if full_sample_data.empty:
            messagebox.showerror("Error", f"No data available for sheet '{selected_sheet}'.")
            return

        # Regenerate aggregate metrics and a fresh plot figure
        agg_df = processing.aggregate_sheet_metrics(full_sample_data, num_columns_per_sample=12)
        agg_df["Sheet"] = selected_sheet
        new_fig = processing.plot_aggregate_trends(agg_df)
    
        # (Optionally, update your cached dictionary if you still wish to cache figures)
        self.aggregate_plots = self.aggregate_plots if hasattr(self, 'aggregate_plots') else {}
        self.aggregate_plots[selected_sheet] = new_fig
    
        # Display the newly generated figure.
        self.display_aggregate_plot(parent, new_fig)


    def stop_threads(self):
        """Stop all active threads."""
        for thread in self.threads:
            if thread.is_alive():
                print(f"Stopping thread: {thread.name}")
                # Threads are generally non-stoppable; they must exit gracefully.

    def load_initial_file(self) -> None:
        """Start a thread to load the initial Excel file"""

        thread = threading.Thread(target=self._load_initial_file_thread, daemon = True)
        self.threads.append(thread)
        thread.start()

    def _load_initial_file_thread(self):
        """Worker thread for loading the initial Excel file."""    
        try:
            file_paths = filedialog.askopenfilenames(title="Select Excel File(s)", filetypes=[("Excel files", "*.xlsx *.xls")]
            )

            if not file_paths:
                messagebox.showinfo("Info", "No files were selected")
                return

            
            self.show_progress_bar("Loading file...")
            
            for index, file_path in enumerate(file_paths, start=1):
                try:
                    self.update_progress_bar(int((index / len(file_paths)) * 100))
                    file_name = os.path.basename(file_path)
                    
                    self.load_excel_file(file_path)
                    with self.lock:
                        self.all_filtered_sheets.append({
                            "file_name": file_name,
                            "file_path": file_path,
                            "filtered_sheets": copy.deepcopy(self.filtered_sheets)
                        })
                except Exception as e:
                    print(f"Error loading file {file_path}: {e}")  # Log the error for debugging
                    messagebox.showwarning(
                        "Warning",
                        f"Failed to load file: {os.path.basename(file_path)}. Error: {e}"
                    )               

            # Automatically select the last loaded file
            if self.all_filtered_sheets:
                with self.lock:
                    last_file = self.all_filtered_sheets[-1]
                self.set_active_file(last_file["file_name"])

                self.root.after(0,self.setup_ui)
    
            self.update_progress_bar(100)

        except Exception as e:
            messagebox.showerror("Error", f"Error loading file: {e}")
        finally:
            self.hide_progress_bar()

    def set_active_file(self, file_name: str) -> None:
        """Set the active file based on the given file name."""
        for file_data in self.all_filtered_sheets:
            
            if file_data["file_name"] == file_name:
                self.current_file = file_name
                self.file_path = file_data.get("file_path", None)
                self.filtered_sheets = file_data["filtered_sheets"]
                if self.file_path is None:
                    raise ValueError(f"No file path associated with the file '{file_name}'.")
                break
        else:
            raise ValueError(f"File '{file_name}' not found.")

    def update_ui_for_current_file(self) -> None:
        """Update UI components to reflect the currently active file."""
        if not self.current_file:
            return

        # Update file dropdown selection
        self.file_dropdown_var.set(self.current_file)

        # Update sheet dropdown
        self.populate_or_update_sheet_dropdown()

        # Update displayed sheet
        current_sheet = self.selected_sheet.get()
        if current_sheet not in self.filtered_sheets:
            # If the selected sheet doesn't exist, default to the first sheet
            first_sheet = list(self.filtered_sheets.keys())[0] if self.filtered_sheets else None
            if first_sheet:
                self.selected_sheet.set(first_sheet)
                self.update_displayed_sheet(first_sheet)
        else:
            # Display the currently selected sheet
            self.update_displayed_sheet(current_sheet)

    def show_progress_bar(self, message: str) -> None:
        """Display a progress bar in a separate window with a white font."""
        if hasattr(self, 'progress_window') and self.progress_window.winfo_exists():
            return  # Prevent multiple instances of the progress bar

        # Create a new top-level window for the progress bar
        self.progress_window = Toplevel(self.root)
        self.progress_window.title("Progress")
        self.progress_window.geometry("400x100")
        self.progress_window.resizable(False, False)
        self.progress_window.configure(bg="#0504AA")  # Set the background color


        #Center the progress window relative to the main application

        self.center_window(self.progress_window)

        # Add a label with white font for the message
        self.progress_label = Label(self.progress_window, text=message, fg="white", bg="#0504AA", font=("Arial", 12))
        self.progress_label.pack(pady=10)

        # Add the progress bar
        self.progress_bar = ttk.Progressbar(self.progress_window, orient='horizontal', mode='determinate', length=300)
        self.progress_bar.pack(pady=10)
        self.progress_bar['value'] = 0  # Initialize progress bar

        # Disable interactions with the main window while the progress window is active
        self.progress_window.transient(self.root)  # Make it a child window
        self.progress_window.grab_set()  # Prevent interactions with the main app

    def update_progress_bar(self, value: int) -> None:
        """Update the progress bar value."""
        if hasattr(self, 'progress_bar') and self.progress_bar.winfo_exists():
            self.progress_bar['value'] = value  # Update progress value
            self.progress_window.update_idletasks()  # Refresh the progress window

    def hide_progress_bar(self) -> None:
        """Destroy the progress bar window after completion."""
        if hasattr(self, 'progress_window') and self.progress_window.winfo_exists():
            self.progress_window.grab_release()  # Release grab from the main app
            self.progress_window.destroy()  # Close the progress window
            del self.progress_window  # Remove the attribute to prevent stale references

    def cleanup_resources(self) -> None:
        """Clean up any resources used by the application."""
        # Stop any background threads or clean up temporary files
        for thread in self.threads:
            if thread.is_alive():
                # Threads should be daemon threads, so they will automatically close
                print(f"Stopping thread {thread.name}")
        

    def on_app_close(self):
        """Handle application shutdown."""
        try:
            for thread in self.threads:
                if thread.is_alive():
                    print(f"Stopping thread {thread.name}")
            self.root.destroy()
            os._exit(0)
        except Exception as e:
            print(f"Error during shutdown: {e}")
            os._exit(1)

    def add_menu(self) -> None:
        """Create a top-level menu with Help and About options."""
        menubar = tk.Menu(self.root)
        
        # Help menu
        helpmenu = tk.Menu(menubar, tearoff=0)
        helpmenu.add_command(label="Help", command=self.show_help)
        helpmenu.add_separator()
        helpmenu.add_command(label="About", command=self.show_about)
        menubar.add_cascade(label="Help", menu=helpmenu)
        
        self.root.config(menu=menubar)

    def show_help(self):
        """Display help dialog."""
        messagebox.showinfo("Help", "This program is designed to be used with excel data according to the SDR Standardized Testing Template.\n \nClick 'Generate Test Report' to create an excel report of a single test, or click 'Generate Full Report' to generate both an excel file and powerpoint file of all the contents within the file.")

    def show_about(self) -> None:
        """Display about dialog."""
        messagebox.showinfo("About", "SDR DataView Alpha Version 1.0\nDeveloped by Charlie Becquet")

    def set_app_colors(self) -> None:
        """Set consistent color theme and fonts for the application."""
        style = ttk.Style()
        self.root.configure(bg=APP_BACKGROUND_COLOR)  # Set background color
        style.configure('TLabel', background=APP_BACKGROUND_COLOR, font=FONT)
        style.configure('TButton', background=BUTTON_COLOR, font=FONT, padding=6)
        style.configure('TCombobox', font=FONT)
        style.map('TCombobox', background=[('readonly', APP_BACKGROUND_COLOR)])
        # Set color for all widgets, loop through children
        for widget in self.root.winfo_children():
            widget.configure(bg='#EFEFEF')

    def ask_open_file(self) -> Optional[str]:
        """Prompt the user to select an Excel file."""
        return filedialog.askopenfilename(title="Select Standardized Testing File",
                                          filetypes=[("Excel files", "*.xlsx *.xls")])

    def add_data(self):
        """Handle adding a new data file directly and update UI accordingly."""
        file_paths = filedialog.askopenfilenames(
            title="Select Excel File(s)", filetypes=[("Excel files", "*.xlsx *.xls")]
        )
        if not file_paths:
            messagebox.showinfo("Info", "No files were selected")
            return

        for file_path in file_paths:
            self.load_excel_file(file_path)
            self.all_filtered_sheets.append({
                "file_name": os.path.basename(file_path),
                "file_path": file_path,
                "filtered_sheets": copy.deepcopy(self.filtered_sheets)
            })
        self.update_file_dropdown()
        # Set the active file to the last loaded file and update UI
        last_file = self.all_filtered_sheets[-1]
        self.set_active_file(last_file["file_name"])
        self.update_ui_for_current_file()
        messagebox.showinfo("Success", f"Data from {len(file_paths)} file(s) added successfully.")



        def handle_file_selection(current_sheet):
            # Open file dialog to select a file
            file_paths = filedialog.askopenfilenames(
                title="Select Excel File(s)", filetypes=[("Excel files", "*.xlsx *.xls")]
            )
            if not file_paths:
                return

            try:
                for file_path in file_paths:
                    # Load the file using the existing function
                    self.load_excel_file(file_path)

                    # Store the new data into the all_filtered_sheets array
                    self.all_filtered_sheets.append({
                        "file_name": os.path.basename(file_path),
                        "file_path": file_path,
                        "filtered_sheets": copy.deepcopy(self.filtered_sheets)
                    })
                self.update_file_dropdown()

                # Update the sheet dropdown based on the selected sheet in the Add Data menu
                selected_sheet = selected_sheet_var.get()  # Get the selected sheet from the dropdown in the Add Data popup
                # Logic to determine which sheet to display
                if selected_sheet == "All Tests":
                    display_sheet = current_sheet  # Display the passed current sheet
                else:
                    display_sheet = selected_sheet  # Display the specifically selected sheet


                # Update the displayed sheet and plots
                self.file_dropdown_var.set(os.path.basename(file_path))  # Update the file dropdown to the new file
                self.selected_sheet.set(display_sheet)  # Set the sheet dropdown to the selected sheet
                self.update_displayed_sheet(display_sheet)  # Update plots and tables for the selected sheet

                messagebox.showinfo("Success", f"Data from {len(file_paths)} files added successfully.")
                popup.destroy()
            except Exception as e:
                messagebox.showerror("Error", f"Failed to load files: {e}")

        # Add "Select File" button
        select_file_button = ttk.Button(popup, text="Select File", command=lambda: handle_file_selection(current_sheet))
        select_file_button.pack(pady=5)

        # Add "Cancel" button
        cancel_button = ttk.Button(popup, text="Cancel", command=popup.destroy)
        cancel_button.pack(pady = 5)

    def load_excel_file(self, file_path: str) -> None:
        """
        Load the selected Excel file and process its sheets.
        For legacy files (non-standard), convert the file to the proper format,
        move the new formatted file into a folder called 'legacy data', and load that copy.
        """
        try:
            # Ensure the file is a valid Excel file.
            if not processing.is_valid_excel_file(os.path.basename(file_path)):
                raise ValueError(f"Invalid Excel file selected: {file_path}")

            if not is_standard_file(file_path):
                # Legacy (non-standard) file: convert it to the proper format.
                # Convert the legacy file using the conversion function.
                legacy_dir = os.path.join(os.path.abspath("."), "legacy data")
                if not os.path.exists(legacy_dir):
                    os.makedirs(legacy_dir)

                standard_df = processing.convert_legacy_file_using_template(file_path)
            
                # Now use the formatted DataFrame (standard_df) for display.
                key = f"Legacy_{os.path.basename(file_path)}"
                self.filtered_sheets = {key: {"data": standard_df, "is_empty": standard_df.empty}}
                self.selected_sheet.set(key)
                self.update_displayed_sheet(key)
            else:
                # Standard file: load and process as usual.
                print("Standard File. Processing")
                self.sheets = processing.load_excel_file(file_path)
                self.filtered_sheets = {
                    name: {"data": data, "is_empty": data.empty}
                    for name, data in self.sheets.items()
                }
                first_sheet = list(self.filtered_sheets.keys())[0]
                self.selected_sheet.set(first_sheet)
                self.update_displayed_sheet(first_sheet)
        except Exception as e:
            messagebox.showerror("Error", f"Error occurred while loading file: {e}")

       
    def setup_frames(self, is_plotting_sheet: bool = False) -> None:
        """Set up or reinitialize the primary frames for UI elements."""
        # Reinitialize top_frame
        if not hasattr(self, 'top_frame') or not self.top_frame:
            self.top_frame = ttk.Frame(self.root)
            self.top_frame.pack(side="top", fill="x", pady=(10, 5), padx=10)

        # Reinitialize display_frame
        if not hasattr(self, 'display_frame') or not self.display_frame:
            self.display_frame = ttk.Frame(self.root)
            self.display_frame.pack(fill="both", expand=True, padx=10, pady=10)

        # Create table_frame (always)
        self.table_frame = ttk.Frame(self.display_frame)
        if is_plotting_sheet:
            self.table_frame.pack(side="left", fill="both", expand=True, padx=(5, 10))
            # Also create plot_frame if needed
            self.plot_frame = ttk.Frame(self.display_frame)
            self.plot_frame.pack(side="right", fill="both", expand=True, padx=(10, 5))
        else:
            # If not plotting, have table_frame fill the entire display_frame
            self.table_frame.pack(fill="both", expand=True)
            self.plot_frame = None


    def populate_sheet_dropdown(self) -> None:
        """Set up the dropdown for sheet selection and handle its events."""
        all_sheet_names = list(self.filtered_sheets.keys())

        # Create a dropdown menu for selecting sheets
        self.drop_down_menu = ttk.Combobox(
            self.top_frame,
            textvariable=self.selected_sheet,
            values=list(self.filtered_sheets.keys()),
            state="readonly"  # Prevent users from typing in custom values
        )
        self.drop_down_menu.pack(pady=(10, 10))
        self.drop_down_menu.place(relx=0.5, rely=0.5, anchor="center")

        # Bind selection event
        self.drop_down_menu.bind(
            "<<ComboboxSelected>>",
            lambda event: self.update_displayed_sheet(self.selected_sheet.get())
        )

        # Automatically select the first sheet if available
        if all_sheet_names:
            first_sheet = all_sheet_names[0]
            self.selected_sheet.set(first_sheet)
            self.update_displayed_sheet(first_sheet)
        else:
            messagebox.showerror("Error", "No sheets found in the file.")
       
    def set_window_size(self, width_ratio: float, height_ratio: float) -> None:
        """
        Set the window size as a percentage of the screen dimensions.

        Args:
            width_ratio (float): Width as a fraction of screen width.
            height_ratio (float): Height as a fraction of screen height.
        """
        
        screen_width = self.root.winfo_screenwidth()
        screen_height = self.root.winfo_screenheight()
        window_width = int(screen_width * width_ratio)
        window_height = int(screen_height * height_ratio)
        position_x = (screen_width - window_width) // 2
        position_y = (screen_height - window_height) // 2
        self.root.geometry(f"{window_width}x{window_height}+{position_x}+{position_y}")

    def add_report_buttons(self, parent_frame: ttk.Frame) -> None:
        """Add buttons for generating reports and align them in the parent frame."""

        # Add Data Button
        add_data_btn = ttk.Button(parent_frame, text="Add Data", command=self.add_data)
        add_data_btn.pack(side="right", padx=(5, 5), pady=(5, 5))

        # Generate Full Report button
        full_report_btn = ttk.Button(parent_frame, text="Generate Full Report", command=self.generate_full_report)
        full_report_btn.pack(side="left", padx=(5, 5), pady=(5, 5))

        # Generate Test Report button
        test_report_btn = ttk.Button(parent_frame, text="Generate Test Report", command=self.generate_test_report)
        test_report_btn.pack(side="left", padx=(5, 5), pady=(5, 5))

    def embed_plot_in_frame(self, fig: plt.Figure, frame: ttk.Frame) -> FigureCanvasTkAgg:
        """
        Helper method to embed a Matplotlib plot into a Tkinter frame and add checkboxes for toggling lines.
    
        Args:
            fig (plt.Figure): The Matplotlib figure to embed.
            frame (ttk.Frame): The Tkinter frame to embed the figure into.

        Returns:
            FigureCanvasTkAgg: The Matplotlib canvas object that has been embedded.
        """
        # Clear any existing widgets in the frame
        for widget in frame.winfo_children():
            widget.pack_forget()

        if self.figure: 
            plt.close(self.figure)  # Close only active figure, clear previous figures to prevent memory bloat

        fig.subplots_adjust(right=0.82) # Adjust right margin of the plot
        
        # Embed the figure into the frame
        self.canvas = FigureCanvasTkAgg(fig, master=frame)
        self.canvas.draw()
        self.canvas.get_tk_widget().pack(fill='both', expand=True)

        # Add Matplotlib toolbar
        toolbar = NavigationToolbar2Tk(self.canvas, frame)
        toolbar.update()
            # Retain figure and axes for future modifications
        self.figure = fig
        self.axes = fig.gca()
        self.lines = self.axes.lines  # Get lines in the current axes

        # Bind mouse and scroll events
        
        self.canvas.mpl_connect("scroll_event", lambda event: self.zoom(event))
        # Reapply checkboxes
        self.add_checkboxes()
        

        return self.canvas


    def generate_full_report(self): 
        save_path = get_save_path(".xlsx")
        if not save_path:
            return

        # Show the progress bar with an appropriate message.
        self.show_progress_bar("Generating full report...")

        ppt_save_path = save_path.replace('.xlsx', '.pptx')
        images_to_delete = []

        try:
            total_sheets = len(self.filtered_sheets)
            processed_count = 0
            with pd.ExcelWriter(save_path, engine='xlsxwriter') as writer:
                for sheet_name in self.filtered_sheets:
                    try:
                        sheet_info = self.filtered_sheets.get(sheet_name)
                        data = sheet_info["data"]
                        is_plotting = plotting_sheet_test(sheet_name, data)
                        
                        process_function = processing.get_processing_function(sheet_name)
                        if is_plotting:
                            processed_data, _, full_sample_data = process_function(data)
                            valid_plot_options = processing.get_valid_plot_options(self.plot_options, full_sample_data)
                        else:
                            data = data.astype(str).replace([np.nan, pd.NA], '')
                            processed_data, _, full_sample_data = process_function(data)
                            valid_plot_options = []

                        if processed_data.empty or full_sample_data.empty:
                            print(f"Skipping sheet '{sheet_name}' due to empty processed data.")
                            continue

                        # Write the Excel portion for this sheet.
                        self.write_excel_report(
                            writer, sheet_name, processed_data, full_sample_data, valid_plot_options, images_to_delete
                        )

                        # If the sheet is plotting, schedule plotting calls in the main thread.
                        if is_plotting and valid_plot_options:
                            # Prepare numeric data once.
                            numeric_data = full_sample_data.apply(pd.to_numeric, errors='coerce')
                            for i, plot_option in enumerate(valid_plot_options):
                                plot_image_path = f"{sheet_name}_{plot_option}_plot.png"
                                # Define a helper function that does the plotting and saving.
                                def generate_plot(po=plot_option, nd=numeric_data, sn=sheet_name):
                                    processing.plot_all_samples(nd, 12, po)
                                    plt.savefig(plot_image_path)
                                    plt.close()
                                # Schedule the plotting on the main thread.
                                self.root.after(0, generate_plot)
                                if plot_image_path not in images_to_delete:
                                    images_to_delete.append(plot_image_path)

                    except Exception as e:
                        print(f"Error processing sheet '{sheet_name}': {e}")
                        continue

                    processed_count += 1
                    progress_value = int((processed_count / total_sheets) * 100)
                    self.update_progress_bar(progress_value)

            # Now generate the PowerPoint report (which calls your usual method).
            try:
                self.write_powerpoint_report(ppt_save_path, images_to_delete)
            except Exception as e:
                print(f"Error writing PowerPoint report: {e}")

            messagebox.showinfo("Success", "Full report saved successfully.")
        except Exception as e:
            messagebox.showerror("Error", f"An error occurred: {e}")
        finally:
            # hide the progress bar regardless of success or failure.
            self.hide_progress_bar()

    def add_plots_to_slide(self, slide, sheet_name: str, full_sample_data: pd.DataFrame, valid_plot_options: list, images_to_delete: list) -> None:
        """
        Add plots to a PowerPoint slide based on valid plot options for the given sheet.

        Args:
            slide (pptx.slide.Slide): The PowerPoint slide to which plots will be added.
            sheet_name (str): The name of the sheet being processed.
            full_sample_data (pd.DataFrame): Data used for plotting.
            valid_plot_options (list): List of valid plot options for the sheet.
            images_to_delete (list): List to track generated plot image paths for cleanup.
        """
        plot_top = Inches(1.21)  # Initial vertical position
        left_column_x = Inches(8.43)  # Left column position
        right_column_x = Inches(10.84)  # Right column position


        numeric_data = full_sample_data.apply(pd.to_numeric, errors='coerce')

        # Ensure there's numeric data before proceeding
        if numeric_data.isna().all(axis=0).all():  # Check if all data is NaN
            print(f"No numeric data available for plotting in sheet '{sheet_name}'.")
            return

        for i, plot_option in enumerate(valid_plot_options):
            # Set the selected plot type
            self.selected_plot_type.set(plot_option)
            plot_image_path = f"{sheet_name}_{plot_option}_plot.png"
            
            # Generate and save the plot
            try:
            
                processing.plot_all_samples(numeric_data, 12, plot_option)
                plt.savefig(plot_image_path)
                plt.close()

                # Avoid adding duplicate entries
                if plot_image_path not in images_to_delete:  # <- Add this block here
                    images_to_delete.append(plot_image_path)

                # Alternate between left and right columns
                plot_x = left_column_x if i % 2 == 0 else right_column_x
                if i % 2 != 0:  # Move down for every second image
                    plot_top += Inches(1.83)

                # Add the plot image to the slide
                slide.shapes.add_picture(plot_image_path, plot_x, plot_top, Inches(2.29), Inches(1.72))
            except Exception as e:
                print(f"Error generating plot '{plot_option}' for sheet '{sheet_name}': {e}")

    def write_powerpoint_report_for_test(self, ppt_save_path: str, images_to_delete: list, sheet_name: str, processed_data: pd.DataFrame, full_sample_data: pd.DataFrame) -> None:
        """
        Generate a PowerPoint report for only the specified sheet.
    
        Args:
            ppt_save_path (str): Path to save the PowerPoint file.
            images_to_delete (list): List of images to clean up after saving.
            sheet_name (str): The name of the sheet to process.
            processed_data (pd.DataFrame): Processed data for the sheet.
            full_sample_data (pd.DataFrame): Full sample data (for plotting) for the sheet.
        """
        try:
            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)

            # (Optional) Create a Cover Slide first
            # If you want the cover slide to be included in the test report as well,
            # copy your cover slide code here. Otherwise, you can skip it.
            # For example, you might decide that the test report does not need a cover slide.

            # Create the slide for the current sheet using a blank layout
            try:
                slide_layout = prs.slide_layouts[5]  # Use blank slide layout
                slide = prs.slides.add_slide(slide_layout)
            except IndexError:
                raise ValueError("Slide layout 5 not found in the PowerPoint template.")

            # Add background image
            background_path = get_resource_path("resources/ccell_background.png")
            slide.shapes.add_picture(
                background_path,
                left=Inches(0),
                top=Inches(0),
                width=prs.slide_width,
                height=prs.slide_height
            )

            # Add logo image (upper right)
            logo_path = get_resource_path("resources/ccell_logo_full.png")
            slide.shapes.add_picture(
                logo_path,
                left=Inches(11.21),
                top=Inches(0.43),
                width=Inches(1.57),
                height=Inches(0.53)
            )

            # Add a title textbox for the sheet
            if slide.shapes.title:
                title_shape = slide.shapes.title
            else:
                # If no placeholder exists, add a new textbox
                title_shape = slide.shapes.add_textbox(Inches(0.45), Inches(0.45), Inches(10.72), Inches(0.64))
            text_frame = title_shape.text_frame
            text_frame.margin_top = 0
            text_frame.margin_bottom = 0
            # Remove any default paragraphs
            for para in list(text_frame.paragraphs):
                text_frame._element.remove(para._element)
            p = text_frame.add_paragraph()
            p.text = sheet_name
            p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(0)
            p.space_after = Pt(0)
            text_frame.word_wrap = True
            run = p.runs[0] if p.runs else p.add_run()
            run.font.name = "Montserrat"
            run.font.size = Pt(32)
            run.font.bold = True

            # Bring the title textbox to the front
            spTree = slide.shapes._spTree
            spTree.remove(title_shape._element)
            spTree.append(title_shape._element)

            # Determine if the sheet is a plotting sheet
            plot_sheet_names = processing.get_plot_sheet_names()
            is_plotting = sheet_name in plot_sheet_names

            if is_plotting:
                # For plotting sheets, add the table (using the processed data) and then add plots.
                table_width = Inches(8.07)
                self.add_table_to_slide(slide, processed_data, table_width, is_plotting)
                valid_plot_options = processing.get_valid_plot_options(self.plot_options, full_sample_data)
                if valid_plot_options:
                    self.add_plots_to_slide(slide, sheet_name, full_sample_data, valid_plot_options, images_to_delete)
                else:
                    print(f"No valid plot options for sheet '{sheet_name}'. Skipping plots.")
            else:
                # For non-plotting sheets, add the table using the full sample data.
                table_width = Inches(13.03)
                self.add_table_to_slide(slide, full_sample_data, table_width, is_plotting)

            processing.clean_presentation_tables(prs)
            prs.save(ppt_save_path)
            print(f"PowerPoint test report saved successfully at {ppt_save_path}.")

        except Exception as e:
            print(f"Error writing PowerPoint test report: {e}")


    def generate_test_report(self):
        """
        Generate an Excel and PowerPoint report for only the currently displayed sheet.
        The sheet is selected from the test dropdown.
        """
        save_path = get_save_path(".xlsx")
        if not save_path:
            return
        # Show the progress bar with an initial message.
        self.show_progress_bar("Generating test report...")
        images_to_delete = []
        try:
            # Get the currently selected sheet name from the test dropdown.
            current_sheet = self.selected_sheet.get()
            if not current_sheet:
                messagebox.showerror("Error", "No sheet is selected for the test report.")
                return

            # Retrieve data for the selected sheet.
            data = self.sheets.get(current_sheet)
            if data is None or data.empty:
                messagebox.showwarning("Warning", f"Sheet '{current_sheet}' is empty.")
                return

            # Process the sheet data.
            process_function = processing.get_processing_function(current_sheet)
            processed_data, _, full_sample_data = process_function(data)
            if processed_data.empty or full_sample_data.empty:
                messagebox.showwarning("Warning", f"Sheet '{current_sheet}' did not yield valid processed data.")
                return

            # Update progress after processing the sheet.
            self.update_progress_bar(30)

            # Write the Excel report for the current sheet.
            with pd.ExcelWriter(save_path, engine='xlsxwriter') as writer:
                valid_plot_options = processing.get_valid_plot_options(self.plot_options, full_sample_data)
                self.write_excel_report(writer, current_sheet, processed_data, full_sample_data, valid_plot_options, images_to_delete)
            self.update_progress_bar(60)
            
            # Determine the PPT file path.
            ppt_save_path = save_path.replace('.xlsx', '.pptx')
            # Write the PowerPoint report for the current sheet.
            self.write_powerpoint_report_for_test(ppt_save_path, images_to_delete, current_sheet, processed_data, full_sample_data)

            # Clean up any temporary images.
            self.cleanup_images(images_to_delete)
            self.update_progress_bar(100)

            messagebox.showinfo("Success", f"Test report saved successfully to:\nExcel: {save_path}\nPowerPoint: {ppt_save_path}")

        except Exception as e:
            messagebox.showerror("Error", f"An error occurred while generating the test report: {e}")

        finally:
            # Hide the progress bar once the process is complete.
            self.hide_progress_bar()

    def add_logo_to_slide(self, slide) -> None:
        """
        Add a logo to the top-right corner of the slide.

        Args:
            slide (pptx.slide.Slide): The slide to which the logo will be added.
        """
        logo_path = get_resource_path('resources/ccell_logo_full.png')
        left = Inches(7.75)
        top = Inches(0.43)
        height = Inches(0.53)
        width = Inches(1.57)

        if os.path.exists(logo_path):
            slide.shapes.add_picture(logo_path, left, top, width, height)
        else:
            print(f"Logo image not found at {logo_path}")

    def add_plots_to_excel(self, writer, sheet_name, full_sample_data, images_to_delete, valid_plot_options):
        """Add plots to the Excel sheet based on the valid plot options."""
        try:
            worksheet = writer.sheets[sheet_name]
            # Filter numeric data for plotting
            numeric_data = full_sample_data.apply(pd.to_numeric, errors = 'coerce')

            # Ensure there's numeric data before proceeding
            
            if numeric_data.isna().all().all():  # Check if all data is NaN
                return
            # Add each plot to the Excel sheet
            for i, plot_option in enumerate(valid_plot_options):
                plot_image_path = f"{sheet_name}_{plot_option}_plot.png"

                try:
                    # Generate and save the plot
                    processing.plot_all_samples(numeric_data, 12, plot_option)
                    plt.savefig(plot_image_path, dpi=300)
                    plt.close()

                    # Track the image for cleanup
                    images_to_delete.append(plot_image_path)

                    # Insert the image into the worksheet
                    col_offset = 10 + (i % 2) * 10  # Alternate columns
                    row_offset = 2 + (i // 2) * 20  # Adjust rows for every two images
                    worksheet.insert_image(row_offset, col_offset, plot_image_path)
                except Exception as e:
                    print(f"Error generating plot '{plot_option}' for sheet '{sheet_name}': {e}")
        except Exception as e:
            print(f"Error adding plots to Excel for sheet '{sheet_name}': {e}")


    def write_excel_report(self, writer, sheet_name, processed_data, full_sample_data, valid_plot_options=[], images_to_delete=None):
        """Write the Excel portion of the report."""
        try:
            # Write processed data to Excel
            processed_data.astype(str).replace([np.nan,pd.NA],'') # faster null handling
            processed_data.to_excel(writer, sheet_name=sheet_name, index=False)
            plot_sheet_names = processing.get_plot_sheet_names()
            is_plotting = sheet_name in plot_sheet_names
            # Add plots to Excel if valid plot options are available
            if is_plotting:
                self.add_plots_to_excel(writer, sheet_name, full_sample_data, images_to_delete, valid_plot_options)
        except Exception as e:
            print(f"Error writing Excel report for sheet '{sheet_name}': {e}")

    def write_powerpoint_report(self, ppt_save_path: str, images_to_delete: list) -> None:
        """
        Generate the PowerPoint report and save it to the specified path.

        Args:
            ppt_save_path (str): Path to save the PowerPoint file.
            images_to_delete (list): List of images to clean up after saving.
        """
        try:
            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)

            # =============================
            # Create the Cover Slide
            # =============================
            cover_slide = prs.slides.add_slide(prs.slide_layouts[6])  # Use a blank layout

            # Add the cover background image (fills the entire slide)
            bg_path = get_resource_path("resources/Cover_Page_Logo.jpg")  # adjust extension if needed
            cover_slide.shapes.add_picture(
                bg_path,
                left=Inches(0),
                top=Inches(0),
                width=prs.slide_width,
                height=prs.slide_height
            )

            # -- Add the Title TextBox --
            # Desired dimensions: width 12", height 3.01", top = 0.88", centered horizontally.
            textbox_title = cover_slide.shapes.add_textbox(
                left=Inches((prs.slide_width.inches - 12) / 2),
                top=Inches(2.35),
                width=Inches(12),
                height=Inches(0.88)
            )
            text_frame = textbox_title.text_frame
            text_frame.margin_top = 0
            text_frame.margin_bottom = 0
            # Remove any default paragraphs
            for para in list(text_frame.paragraphs):
                text_frame._element.remove(para._element)
            p = text_frame.add_paragraph()
            # Set title text using the loaded file name if available
            if self.file_path:
                p.text = f"{os.path.splitext(os.path.basename(self.file_path))[0]} Standard Test Report"
            else:
                p.text = "Standard Test Report"

            p.alignment = PP_ALIGN.CENTER
            p.space_before = Pt(0)
            p.space_after = Pt(0)
            text_frame.word_wrap = True  # Force the text to stay within slide 
            run = p.runs[0] if p.runs else p.add_run()
            run.font.name = "Montserrat"
            run.font.size = Pt(46)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # White font color

            # Bring the title textbox to the front
            cover_slide.shapes._spTree.remove(textbox_title._element)
            cover_slide.shapes._spTree.append(textbox_title._element)

            # -- Add the Subtitle TextBox --
            # Desired dimensions: left 5.73", top 4.05", width 1.87", height 0.37"
            textbox_sub = cover_slide.shapes.add_textbox(
                left=Inches(5.73),
                top=Inches(4.05),
                width=Inches(1.87),
                height=Inches(0.37)
            )
            sub_frame = textbox_sub.text_frame
            sub_frame.margin_top = 0
            sub_frame.margin_bottom = 0
            # Remove any default paragraphs
            for para in list(sub_frame.paragraphs):
                sub_frame._element.remove(para._element)
            p2 = sub_frame.add_paragraph()
            # Format the current date as "DD Month YYYY"
            p2.text = datetime.today().strftime("%d %B %Y")
            p2.alignment = PP_ALIGN.CENTER
            p2.space_before = Pt(0)
            p2.space_after = Pt(0)
            sub_frame.word_wrap = False
            run2 = p2.runs[0] if p2.runs else p2.add_run()
            run2.font.name = "Montserrat"
            run2.font.size = Pt(16)
            run2.font.bold = True
            run2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # White font color

            # Bring the subtitle textbox to the front
            cover_slide.shapes._spTree.remove(textbox_sub._element)
            cover_slide.shapes._spTree.append(textbox_sub._element)

            # -- Add the Logo Image --
            # Desired dimensions: left 5.7", top 6.12", width 1.93", height 0.66"
            logo_path = get_resource_path("resources/ccell_logo_full_white.png")
            cover_slide.shapes.add_picture(
                logo_path,
                left=Inches(5.7),
                top=Inches(6.12),
                width=Inches(1.93),
                height=Inches(0.66)
            )

        
            # =============================
            # Process the Remaining Sheets
            # =============================
            # get image paths from resources
            background_path = get_resource_path("resources/ccell_background.png")
            logo_path = get_resource_path("resources/ccell_logo_full.png")
            plot_sheet_names = processing.get_plot_sheet_names()
            # Iterate through the sheets
            for sheet_name in self.filtered_sheets.keys():
                try:
                    # Retrieve sheet data
                    data = self.sheets.get(sheet_name)
                    if data is None or data.empty:
                        print(f"Skipping sheet '{sheet_name}': No data available.")
                        continue

                    # Process the sheet data
                    process_function = processing.get_processing_function(sheet_name)
                    processed_data, _, full_sample_data = process_function(data)

                    # Validate processed data
                    if processed_data.empty:
                        print(f"Skipping sheet '{sheet_name}': Processed data is empty.")
                        continue

                    # Create a new slide
                    try:
                        slide_layout = prs.slide_layouts[5]  # Use blank slide layout
                        slide = prs.slides.add_slide(slide_layout)
                    except IndexError:
                        raise ValueError(f"Slide layout 5 not found in the PowerPoint template.")

                    # Add background image
                    slide.shapes.add_picture(
                        background_path,
                        left=Inches(0),
                        top=Inches(0),
                        width=Inches(13.33),
                        height=Inches(7.5)
                    )

                    # Add logo
                    slide.shapes.add_picture(
                        logo_path,
                        left=Inches(11.21),
                        top=Inches(0.43),
                        width=Inches(1.57),
                        height=Inches(0.53)
                    )

                    # Add a title to the slide (or update existing title placeholder)
                    if slide.shapes.title:
                        title_shape = slide.shapes.title
                        # Set title text, position, size, and font
                        title_shape.left = Inches(0.45)
                        title_shape.top = Inches(0.45)
                        title_shape.width = Inches(10.72)
                        title_shape.height = Inches(0.64)

                        text_frame = title_shape.text_frame
                        text_frame.margin_top = 0
                        text_frame.margin_bottom = 0

                        # Clear the text frame by removing all paragraphs.
                        # (Sometimes clear() doesn't remove default paragraphs.)
                        for para in list(text_frame.paragraphs):
                            text_frame._element.remove(para._element)

                        p = text_frame.add_paragraph()
                        p.text = sheet_name
                        p.alignment = PP_ALIGN.LEFT  # Force left-justification
                        p.space_before = Pt(0)
                        p.space_after = Pt(0)

                        run = p.runs[0] if p.runs else p.add_run()
                        run.font.name = "Montserrat"
                        run.font.size = Pt(32)
                        run.font.bold = True
                    else:
                        print(f"Warning: No title placeholder found for sheet '{sheet_name}'.")

                    is_plotting = sheet_name in plot_sheet_names

                    # **Handle Non-Plotting Sheets**
                    if not is_plotting:
                        if processed_data.empty:
                            print(f"Skipping non-plotting sheet '{sheet_name}' due to empty data.")
                            continue

                        
                        # Add the table directly
                        table_width = Inches(13.03)
                        self.add_table_to_slide(slide, full_sample_data, table_width, is_plotting)
                        
                    else:
                        # **Handle Plotting Sheets**
                        table_width = Inches(8.07)
                        self.add_table_to_slide(slide, processed_data, table_width, is_plotting)

                        # Add plots to slide if valid options exist
                        valid_plot_options = processing.get_valid_plot_options(self.plot_options, full_sample_data)
                        if valid_plot_options:
                            self.add_plots_to_slide(slide, sheet_name, full_sample_data, valid_plot_options, images_to_delete)
                        else:
                            print(f"No valid plot options for sheet '{sheet_name}'. Skipping plots.")

                    if slide.shapes.title:
                        title_shape = slide.shapes.title
                        spTree = slide.shapes._spTree
                        spTree.remove(title_shape._element)
                        spTree.append(title_shape._element)

                except Exception as sheet_error:
                    print(f"Error processing sheet '{sheet_name}': {sheet_error}")

            # Save the PowerPoint file
            processing.clean_presentation_tables(prs)
            prs.save(ppt_save_path)
            print(f"PowerPoint report saved successfully at {ppt_save_path}.")

        except Exception as e:
            print(f"Error writing PowerPoint report: {e}")

        finally:
            # Cleanup temporary images
            for image_path in set(images_to_delete):  # Use set to avoid duplicates
                if os.path.exists(image_path):  # Check if the file exists before deleting
                    try:
                        os.remove(image_path)
                    except OSError as cleanup_error:
                        print(f"Error deleting image {image_path}: {cleanup_error}")



    def cleanup_images(self, images_to_delete: list) -> None:
        """Delete temporary images used in the report."""
        for image_path in images_to_delete:
            try:
                os.remove(image_path)
            except OSError as e:
                print(f"Error deleting file {image_path}: {e}")
    
    def add_table_to_slide(self, slide, processed_data, table_width, is_plotting):
        """
        Add the processed table to the slide. If the number of rows exceeds max_rows_per_table,
        split the table into multiple smaller tables and dynamically adjust their width to fit within the slide.
  
        Args:
            slide (pptx.slide.Slide): The PowerPoint slide where the table will be added.
            processed_data (pd.DataFrame): The data to display in the table.
            table_width (float): The width of the allocated table space on the slide (in inches).
        """
        processed_data = processed_data.fillna(' ')
        processed_data = processed_data.astype(str)


        rows, cols = processed_data.shape
        
        # Check limitations: skip table creation if conditions are met
        if cols > 20 and rows > 30 and not is_plotting:
            print(f"Skipping table creation for slide. Number of columns ({cols}) exceeds 20, "
                  f"number of rows ({rows}) exceeds 30, and it is not a plotting sheet.")
            return

        # Define constants for table placement
        
        table_left = Inches(0.15)
        table_top = Inches(1.19)  # Starting vertical position
        max_table_height = Inches(5.97)  # Maximum height for each table
  
        
  
        # Add the table to the slide
        table_shape = slide.shapes.add_table(
            rows+1,  # +1 for the header row
            cols,
            table_left,
            table_top,
            table_width,
            Inches(0.4 * (rows + 1))  # Auto-adjust height
        ).table
  
        # Configure font constants
        header_font = {
            'name': 'Arial',
            'size': Pt(10),
            'bold': True
        }
        
        cell_font = {
            'name': 'Arial',
            'size': Pt(8),
            'bold': False
        }
  
        # Add headers with formatting
        for col_idx, col_name in enumerate(processed_data.columns):
            cell = table_shape.cell(0, col_idx)
            cell.text = str(col_name)
            
            # Format header text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = header_font['name']
                paragraph.font.size = header_font['size']
                paragraph.font.bold = header_font['bold']
                paragraph.alignment = PP_ALIGN.CENTER

        # Add data rows with formatting
        for row_idx in range(rows):
            for col_idx in range(cols):
                cell = table_shape.cell(row_idx + 1, col_idx)
                cell.text = processed_data.iat[row_idx, col_idx]
                
                # Format data cell text
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.name = cell_font['name']
                    paragraph.font.size = cell_font['size']
                    paragraph.font.bold = cell_font['bold']
                    paragraph.alignment = PP_ALIGN.CENTER

        
        return True


    def plot_all_samples(self, frame: ttk.Frame, full_sample_data: pd.DataFrame, num_columns_per_sample: int) -> None:
        """
        Embed the plot in the provided frame using Tkinter's FigureCanvasTkAgg, and add checkboxes.
        Args:
        frame (ttk.Frame): The Tkinter frame to embed the plot into.
        full_sample_data (pd.DataFrame): The sample data to be plotted.
        num_columns_per_sample (int): The number of columns per sample.
        """
        # Clear any existing widgets in the frame
        for widget in frame.winfo_children():
            widget.pack_forget()

        # Check if the data is empty (no valid samples)
        if full_sample_data.empty or full_sample_data.isna().all().all():
            messagebox.showwarning("Warning", "The data is empty or invalid for plotting.")
            return

        # Get the figure and sample names from the processing function
        result = processing.plot_all_samples(full_sample_data, num_columns_per_sample, self.selected_plot_type.get())
        if isinstance(result, tuple):
            fig, sample_names = result  # Extract figure and sample names
        else:
            fig, sample_names = result, None  # Fallback if no sample names are provided

        # Embed the figure into the Tkinter frame using the helper method
        self.canvas = self.embed_plot_in_frame(fig, frame)

        # Pass sample names to add_checkboxes
        self.add_checkboxes(sample_names=sample_names)
        
    def update_plot(self, full_sample_data, num_columns_per_sample, frame=None):
        """Update the plot dynamically based on the selected plot type."""
        if frame is None:
            raise ValueError("Plot frame must be provided.")

        # Clear only the plot area
        for widget in frame.winfo_children():
            widget.pack_forget()

        # Embed the updated plot
        self.plot_all_samples(frame, full_sample_data, num_columns_per_sample)

    def adjust_window_size(self, fixed_width=1500):
        """
        Adjust the window height dynamically to fit the content while keeping the width constant.
        Args:
            fixed_width (int): The fixed width of the window.
        """

        if not isinstance(self.root, (tk.Tk, tk.Toplevel)):
            raise ValueError("Expected 'self.root' to be a tk.Tk or tk.Toplevel instance")

        # Update the geometry manager and calculate the required size
        self.root.update_idletasks()

        # Get the required dimensions for the window height to fit its content
        required_height = self.root.winfo_reqheight()

        # Get the current screen dimensions
        screen_height = self.root.winfo_screenheight()
        screen_width = self.root.winfo_screenwidth()
        # Constrain the height to the screen dimensions if necessary
        final_height = min(required_height, screen_height)

        current_x = self.root.winfo_x()
        current_y = self.root.winfo_y()

        # Set the geometry dynamically with the fixed width
        self.root.geometry(f"{fixed_width}x{final_height}+{current_x}+{current_y}")

    def update_displayed_sheet(self, sheet_name: str) -> None:
        """
        Update the displayed sheet and dynamically manage the plot options and plot type dropdown.
        Clears the plot area if the sheet is empty.
        """
        if not sheet_name or sheet_name not in self.filtered_sheets:
            return

        # Retrieve the sheet data
        sheet_info = self.filtered_sheets.get(sheet_name)
        if not sheet_info:
            messagebox.showerror("Error", f"Sheet '{sheet_name}' not found.")
            return

        data = sheet_info["data"]
        is_empty = sheet_info["is_empty"]
        is_plotting_sheet = plotting_sheet_test(sheet_name, data)

        self.setup_frames(is_plotting_sheet)

        if not hasattr(self, 'image_frame') or self.image_frame is None:
            print("Error: Image frame is missing. Recreating...")
            self.image_frame = ttk.Frame(self.display_frame)
            self.image_frame.pack(fill="both",expand=True)

        self.image_loader = ImageLoader(self.image_frame, is_plotting_sheet)

        # Clear display selectively based on the sheet type
        self.clear_display_frame(is_plotting_sheet=is_plotting_sheet, is_empty_sheet=is_empty)

        if is_empty:
            # Display a message for empty sheets
            empty_label = Label(
                self.display_frame,
                text="This sheet is empty.",
                font=("Arial", 14),
                fg="red"
            )
            empty_label.pack(anchor="center", pady=20)
            return

        # Process data
        
        try:
            process_function = processing.get_processing_function(sheet_name)
            processed_data, _, full_sample_data = process_function(data)
        except Exception as e:
            messagebox.showerror("Processing Error", f"Error processing sheet '{sheet_name}': {e}")
            return

        # Display the table
        self.display_table(self.table_frame, processed_data, sheet_name, is_plotting_sheet)

        if is_plotting_sheet:
            self.display_plot(full_sample_data)

        self.image_loader.add_images()

    def clear_plot_area(self):
        """
        Clear the plot content while preserving the plot frame structure.
        This function ensures proper cleanup of Matplotlib figures and Tkinter widgets
        to prevent memory leaks and improve performance.
        """
        # Clear the plot frame if it exists
        if hasattr(self, 'plot_frame') and self.plot_frame.winfo_exists():
            for widget in self.plot_frame.winfo_children():
                widget.destroy()  # Remove all child widgets from the plot frame

        # Close and release Matplotlib figure resources
        if self.figure:
            plt.close(self.figure)  # Close the figure to free up memory
            self.figure = None  # Clear the reference

        # Destroy the Tkinter canvas widget if it exists
        if self.canvas:
            self.canvas.get_tk_widget().destroy()  # Remove the canvas from the UI
            self.canvas = None  # Clear the reference

        # Reset axes and other plot-related attributes
        self.axes = None
        self.lines = None
        self.line_labels = []
        self.original_lines_data = []
        self.check_buttons = None
        self.checkbox_cid = None

    def display_plot(self, full_sample_data):
        """
        Display the plot in the plot frame based on the current data.
        """
        if not hasattr(self, 'plot_frame') or self.plot_frame is None:
            self.plot_frame = ttk.Frame(self.display_frame)
            self.plot_frame.pack(side="right", fill="both", expand=True, padx=5, pady=5)

        # Generate and embed the plot
        self.plot_all_samples(self.plot_frame, full_sample_data, num_columns_per_sample=12)

        self.add_plot_dropdown(self.plot_frame)

    def add_plot_dropdown(self, frame: ttk.Frame) -> None:
        """Add or update the plot type dropdown dynamically."""
        if self.plot_dropdown and self.plot_dropdown.winfo_exists():
            # Update existing dropdown values dynamically
            self.plot_dropdown["values"] = self.plot_options
        else:
            # Create the dropdown if it doesn't exist or was destroyed
            self.plot_dropdown = ttk.Combobox(
                frame,
                textvariable=self.selected_plot_type,
                values=self.plot_options,
                state="readonly",
                font=('Arial', 12)
            )
            self.plot_dropdown.pack(pady=10)  # Place the dropdown below the plot
            self.plot_dropdown.bind("<<ComboboxSelected>>", self.update_plot_from_dropdown)

        # Set the first plot option as the default if not already set
        if self.selected_plot_type.get() not in self.plot_options:
            self.selected_plot_type.set(self.plot_options[0])

    def update_plot_from_dropdown(self, event) -> None:
        """Update the plot when the dropdown selection changes."""
        selected_plot = self.selected_plot_type.get()
        

        if not selected_plot:
            messagebox.showerror("Error", "No plot type selected.")
            return

        try:
            # Retrieve the current sheet name and data
            current_sheet_name = self.selected_sheet.get()
            if current_sheet_name not in self.filtered_sheets:
                messagebox.showerror("Error", f"Sheet '{current_sheet_name}' not found.")
                return

            sheet_info = self.filtered_sheets[current_sheet_name]
            sheet_data = sheet_info["data"]

            if sheet_data.empty:
                messagebox.showwarning("Warning", "The selected sheet has no data.")
                return

            # Process data for the selected plot option
            process_function = processing.get_processing_function(current_sheet_name)
            processed_data, _, full_sample_data = process_function(sheet_data)

            # Update the plot based on the selected option
            self.plot_all_samples(self.plot_frame, full_sample_data,12)

            # Recreate dropdown if destroyed or inaccessible
            if not self.plot_dropdown or not self.plot_dropdown.winfo_exists():
                self.add_plot_dropdown(self.plot_frame)
            else:
                # Update the dropdown's values if it already exists
                self.plot_dropdown["values"] = self.plot_options
                self.plot_dropdown.set(selected_plot)

            # Update the displayed sheet (e.g., re-render table or UI updates)
            self.update_displayed_sheet(current_sheet_name)

        except Exception as e:
            messagebox.showerror("Error", f"An error occurred while updating the plot: {e}")

    def clear_display_frame(self, is_plotting_sheet: bool, is_empty_sheet: bool) -> None:
        """
        Clear specific parts of the display frame based on the type of sheet being displayed.

        Args:
            is_plotting_sheet (bool): Whether the sheet is a plotting sheet.
            is_empty_sheet (bool): Whether the sheet is empty.
        """
        # Clear the table frame
        if hasattr(self, 'table_frame') and self.table_frame:
            if self.table_frame.winfo_exists():
                for widget in self.table_frame.winfo_children():
                    widget.destroy()
                self.table_frame.destroy()
            self.table_frame = None  # Reset the reference

        
        if hasattr(self, 'plot_frame') and self.plot_frame:
            if self.plot_frame.winfo_exists():
                for widget in self.plot_frame.winfo_children():
                    widget.destroy()
                self.plot_frame.destroy()
            self.plot_frame = None  # Reset the reference

        # Clear any Matplotlib objects if a plotting sheet is being cleared
        if is_plotting_sheet:
            self.figure = None
            self.canvas = None
            self.axes = None

        self.setup_frames(is_plotting_sheet)



    def display_table(self, frame, data, sheet_name, is_plotting_sheet = False):
        for widget in frame.winfo_children():
            widget.destroy()

        # Deduplicate columns
        data = clean_columns(data)
        data.columns = data.columns.map(str) # ensure column headers are strings
        
        if data.empty:
            messagebox.showwarning("Warning", f"Sheet '{sheet_name}' contains no data to display.")
            return

        data = data.astype(str)
        data = data.replace([np.nan, pd.NA], '', regex=True)

        table_frame = ttk.Frame(frame, padding=(2, 1))
        table_frame.pack(fill='both', expand=True)

        # Force the frame to update its geometry
        table_frame.update_idletasks()
        available_width = table_frame.winfo_width()
        num_columns = len(data.columns)

        # Calculate a new cell width based on available width (if table would be too narrow, use the default)
        calculated_cellwidth = max(120, available_width // num_columns)

        v_scrollbar = ttk.Scrollbar(table_frame, orient='vertical')
        h_scrollbar = ttk.Scrollbar(table_frame, orient='horizontal')

        model = TableModel()
        table_data_dict = data.to_dict(orient='index')
        model.importDict(table_data_dict)

        default_cellwidth = 120
        default_rowheight = 30


        if is_plotting_sheet:
            # Adjust scaling factors for plotting sheets
            
            font_height = 20  # Taller rows to accommodate larger content
            char_per_line = 12  # Fewer characters per line for wider spacing

            # Calculate optimal row height
            row_heights = []
            for _, row in data.iterrows():
                max_lines = 1  # Default minimum height for single-line cells
                for cell in row:
                    if cell:
                        # Calculate lines needed for this cell
                        cell_length = len(cell)
                        lines = (cell_length // char_per_line) + 1
                        max_lines = max(max_lines, lines)
                row_heights.append(max_lines * font_height)

            # Use the maximum row height across all rows
            max_row_height = int(max(row_heights))

            # Override default values
            
            default_rowheight = max_row_height


        table_canvas = TableCanvas(table_frame, model=model, cellwidth=calculated_cellwidth, cellbackgr='#4CC9F0',
                                    thefont=('Arial', 12), rowheight=default_rowheight, rowselectedcolor='#FFFFFF',
                                    editable=False, yscrollcommand=v_scrollbar.set, xscrollcommand=h_scrollbar.set, showGrid=True)


        table_canvas.grid(row=0, column=0, sticky='nsew')
        v_scrollbar.grid(row=0, column=1, sticky='ns')
        h_scrollbar.grid(row=1, column=0, sticky='ew')

        table_canvas.show()

        # Add the "View Raw Data" button below the table
        button_frame = ttk.Frame(frame)
        button_frame.pack(fill='x', pady=10)
        view_raw_button = ttk.Button(button_frame, text="View Raw Data", command=lambda: self.open_raw_data_in_excel(sheet_name))
        view_raw_button.pack(anchor='center')

    def reload_excel_file(self):
        """
        Reloads the Excel file into the program, preserving the state of the UI.
        """
        try:
            # Re-load the Excel file
            self.load_excel_file(self.file_path)

        except Exception as e:
            messagebox.showerror("Error", f"Failed to reload the Excel file: {e}")

    def open_raw_data_in_excel(self, sheet_name=None):
        """
        Opens an Excel file using the default system application (Excel).
        Ensures that closing the popup closes Excel, and vice versa.
        """
        try:
            # Use the file path from the class instance
            file_path = self.file_path

            # Ensure the file path is valid
            if not isinstance(file_path, str):
                raise ValueError(f"Invalid file path. Expected string, got {type(file_path).__name__}")

            # Check if the file exists
            if not os.path.exists(file_path):
                raise FileNotFoundError(f"The file {file_path} does not exist.")

            # Use `os.startfile` to open the file
            os.startfile(file_path)

            # Wait for Excel to fully open
            for _ in range(10):  # Check every 0.5 seconds for up to 5 seconds
                if any(proc.name().lower() == "excel.exe" for proc in psutil.process_iter()):
                    break
                time.sleep(0.5)  # Brief delay before re-checking
            else:
                raise TimeoutError("Excel did not open in the expected time.")

            # Track the Excel processes
            excel_processes = [proc for proc in psutil.process_iter() if proc.name().lower() == "excel.exe"]

            # Create a Toplevel window as a custom messagebox
            dialog = Toplevel(self.root)
            dialog.title("Excel Opened")
            dialog.geometry("300x150")
            dialog.grab_set()  # Make it modal
            dialog.protocol("WM_DELETE_WINDOW", lambda: close_excel_and_popup())  # Ensure Excel closes if dialog is closed

            label = Label(
                dialog,
                text=f"The file has been opened in Excel. Please navigate to the sheet '{sheet_name}' if necessary. "
                        "This dialog will close automatically when Excel is closed.",
                wraplength=280,
                justify="center"
            )
            label.pack(pady=20)

            def close_excel_and_popup():
                """Terminate all Excel processes and close the dialog."""
                for proc in excel_processes:
                    if proc.is_running():
                        proc.terminate()
                        proc.wait()  # Ensure the process is fully terminated
                dialog.destroy()

            # Monitor Excel and close the dialog automatically if Excel is closed
            def monitor_excel():
                nonlocal excel_processes

                excel_processes = [proc for proc in excel_processes if proc.is_running()]

                if not any(proc.is_running() for proc in excel_processes):
                    dialog.destroy()  # Automatically close the dialog if Excel is closed
                else:
                    dialog.after(100, monitor_excel)  # Recheck every 0.5 seconds for responsiveness

            # Start monitoring Excel
            dialog.after(100, monitor_excel)

            # Add a Close button
            close_button = Button(dialog, text="Close", command=close_excel_and_popup)
            close_button.pack(pady=10)

        except Exception as e:
            messagebox.showerror("Error", f"Failed to open Excel: {e}")

    def update_plot_dropdown(self):
        """Update the plot type dropdown with only the valid options and manage visibility."""
        if valid_plot_options:
            # Update the plot type dropdown with available plot options
            self.plot_dropdown['values'] = valid_plot_options

            # Check if the selected plot type is valid, set to the first valid option if not
            if self.selected_plot_type.get() not in valid_plot_options:
                self.selected_plot_type.set(valid_plot_options[0])

            # Ensure the plot dropdown is visible and bind its event
            self.plot_dropdown.pack(fill="x", pady=(5, 5))
            self.plot_dropdown.bind(
                "<<ComboboxSelected>>",
                lambda event: self.plot_all_samples(self.display_frame, self.get_full_sample_data(), 12)
            )
        else:
            # Hide the plot dropdown if no valid options are available
            self.plot_dropdown.pack_forget()

    def on_bar_checkbox_click(self, wrapped_label):
        """
        Toggles the visibility of bars when a checkbox is clicked.
        """
        original_label = self.label_mapping.get(wrapped_label)
        if original_label is None:
            return  # Safeguard against unexpected issues

        # Find the bar (patch) corresponding to the original label
        index = self.line_labels.index(original_label)
        bar = self.axes.patches[index]
        bar.set_visible(not bar.get_visible())

        # Redraw the canvas
        if self.canvas:
            self.canvas.draw_idle()

    def add_checkboxes(self, sample_names=None):
        """
        Adds checkboxes to toggle the visibility of lines or bars in the plot.
        Includes text wrapping for labels, a title, and a border above the checkboxes.
        """
        self.line_labels = []
        self.original_lines_data = []
        self.label_mapping = {}
        wrapped_labels = []

        if self.check_buttons:
            self.check_buttons.ax.clear()  # Clear existing checkboxes
            self.check_buttons = None

        # Check if the plot is a bar chart
        is_bar_chart = self.selected_plot_type.get() == "TPM (Bar)"

        if sample_names is None:
            sample_names = self.line_labels

        # Create labels based on the plot type
        if is_bar_chart and sample_names:
            self.line_labels = sample_names
            self.original_lines_data = [(patch.get_x(), patch.get_height()) for patch in self.axes.patches]
        else:
            # For line plots
            if not self.line_labels:
                self.line_labels = [line.get_label() for line in self.lines]
                self.original_lines_data = [(line.get_xdata(), line.get_ydata()) for line in self.lines]


        # Wrap labels for better readability
        
        for label in self.line_labels:
            wrapped_label = wrap_text(text=label, max_width=10)  # Adjust `max_width` as needed
            self.label_mapping[wrapped_label] = label
            wrapped_labels.append(wrapped_label)

        # Create axes for checkboxes
        checkbox_ax = self.figure.add_axes([0.835, 0.38, 0.15, 0.5])
        self.check_buttons = CheckButtons(checkbox_ax, wrapped_labels, [True] * len(self.line_labels))
        
        # Clean up checkbox axes
        checkbox_ax.tick_params(left=False, bottom=False, labelleft=False, labelbottom=False)
        checkbox_ax.grid(False)
        for spine in checkbox_ax.spines.values():
            spine.set_visible(False)

        
        # Adjust font size of checkboxes
        for label in self.check_buttons.labels:
            label.set_fontsize(8)  # Set smaller font size

        # Add a clean border rectangle around the checkbox frame
        rect = plt.Rectangle(
            (0.05, 0.05), 0.9, 0.9,  # Position and size of the rectangle
            transform=checkbox_ax.transAxes,  # Use axes-relative coordinates
            facecolor='none',  # Transparent background
            edgecolor='black',  # Black border color
            linewidth=1.5,  # Border thickness
            zorder=10  # Make sure it appears above other elements
        )
        checkbox_ax.add_patch(rect)
        # Add a title above the checkboxes
        title_text = PLOT_CHECKBOX_TITLE
        title_x = 0.835 + 0.15 / 2  # Center the title horizontally above the checkbox area
        title_y = 0.58 + 0.3 + 0.035  # Slightly above the checkbox area
        self.figure.text(title_x, title_y, title_text,
                            fontsize=8, ha='center', va='center', wrap=True)

        # Add a border rectangle around the title
        title_border_x = 0.835  # Same as checkbox x-axis start
        title_border_y = title_y - 0.03  # Below the title text
        border_width = 0.15
        border_height = 0.065
        self.figure.add_artist(
            plt.Rectangle((title_border_x, title_border_y), border_width, border_height,
                            edgecolor='black', facecolor='white', lw=1, zorder=2)
        )

            # Add functionality to toggle visibility on checkbox click
        if is_bar_chart:
            self.checkbox_cid = self.check_buttons.on_clicked(self.on_bar_checkbox_click)
        else:
            self.checkbox_cid = self.check_buttons.on_clicked(self.on_checkbox_click)

        # Redraw the canvas after adding checkboxes
        if self.canvas:
            self.canvas.draw_idle()

    def zoom(self, event):
        """
        Handles zooming in and out on the plot with the cursor as the zoom origin.
        """
        if event.inaxes is not self.axes:  # Ensure the event is on the correct axes
            return

        # Get the current axis limits
        x_min, x_max = self.axes.get_xlim()
        y_min, y_max = self.axes.get_ylim()

        # Get the cursor position
        cursor_x, cursor_y = event.xdata, event.ydata

        # Define the zoom scale
        zoom_scale = 0.9 if event.button == 'up' else 1.1  # Zoom in for scroll up, zoom out for scroll down

        if cursor_x is None or cursor_y is None:  # Ensure the cursor is within the data range
            return

        # Calculate new axis limits centered on the cursor
        new_x_min = cursor_x - (cursor_x - x_min) * zoom_scale
        new_x_max = cursor_x + (x_max - cursor_x) * zoom_scale
        new_y_min = cursor_y - (cursor_y - y_min) * zoom_scale
        new_y_max = cursor_y + (y_max - cursor_y) * zoom_scale

        # Apply new limits to the axes
        self.axes.set_xlim(new_x_min, new_x_max)
        self.axes.set_ylim(new_y_min, new_y_max)

        # Redraw the canvas to reflect changes
        if self.canvas:
            self.canvas.draw_idle()

    def on_checkbox_click(self, wrapped_label):
        """Toggle the visibility of lines or bars based on checkbox click."""
        original_label = self.label_mapping.get(wrapped_label)
        if original_label is None:
            return

        index = list(self.label_mapping.values()).index(original_label)
        if self.selected_plot_type.get() == "TPM (Bar)":
            bar = self.lines[index]  # Use existing bar references
            bar.set_visible(not bar.get_visible())
        else:
            line = self.lines[index]  # Use existing line references
            line.set_visible(not line.get_visible())

        if self.canvas:
            self.canvas.draw_idle()

    def save_to_new_file(self, data, original_file_path):
        new_file_path = get_save_path
        if not new_file_path:
            return

        with pd.ExcelWriter(new_file_path, engine='xlsxwriter') as writer:
            data.to_excel(writer, index=False, sheet_name='Sheet1')

        messagebox.showinfo("Save Successful", f"Data saved to {new_file_path}")

    def restore_all(self):
        """Restore all lines or bars to their original state."""
        if self.check_buttons:
            self.check_buttons.disconnect(self.checkbox_cid)

        if self.selected_plot_type.get() == "TPM (Bar)":
            # Restore bars
            for patch, original_data in zip(self.axes.patches, self.original_lines_data):
                patch.set_x(original_data[0])
                patch.set_height(original_data[1])
                patch.set_visible(True)
        else:
            # Restore lines
            for line, original_data in zip(self.lines, self.original_lines_data):
                line.set_xdata(original_data[0])
                line.set_ydata(original_data[1])
                line.set_visible(True)

        # Reset all checkboxes to the default "checked" state (forcefully activating each one)
        if self.check_buttons:
            for i in range(len(self.line_labels)):
                if not self.check_buttons.get_status()[i]:  # Only activate if not already active
                    self.check_buttons.set_active(i)  # Reset each checkbox to "checked"

            # Reconnect the callback for CheckButtons after restoring
            self.checkbox_cid = self.check_buttons.on_clicked(
                self.on_bar_checkbox_click if self.selected_plot_type.get() == "TPM (Bar)" else self.on_checkbox_click
            )

        # Redraw the canvas to reflect the restored state
        if self.canvas:
            self.canvas.draw_idle()


    