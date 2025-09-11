# services/report_service.py
"""
services/report_service.py
Consolidated report generation and calculation service.
This consolidates report generation logic from report_generator.py and calculation
services from calculation_service.py into a unified service.
"""

import os
import re
import json
import math
import shutil
import statistics
import tempfile
import traceback
from typing import Optional, Dict, Any, List, Tuple, Union, Callable
from pathlib import Path
from datetime import datetime

# Third party imports
import pandas as pd
import numpy as np
import openpyxl
from openpyxl import Workbook, load_workbook

# PowerPoint imports
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    POWERPOINT_AVAILABLE = True
except ImportError:
    print("WARNING: PowerPoint functionality disabled - python-pptx not installed")
    POWERPOINT_AVAILABLE = False

# Plotting imports
try:
    import matplotlib.pyplot as plt
    import matplotlib.dates as mdates
    from matplotlib.backends.backend_agg import FigureCanvasAgg
    PLOTTING_AVAILABLE = True
except ImportError:
    print("WARNING: Plotting functionality disabled - matplotlib not installed")
    PLOTTING_AVAILABLE = False


def debug_print(message: str):
    """Debug print function for report operations."""
    print(f"DEBUG: ReportService - {message}")


def round_values(value: float, decimal_places: int = 3) -> float:
    """Round values for display consistency."""
    try:
        return round(float(value), decimal_places)
    except (ValueError, TypeError):
        return 0.0


def get_resource_path(relative_path: str) -> str:
    """Get the absolute path to a resource file."""
    try:
        # Try to find the resource file relative to the current directory
        if os.path.exists(relative_path):
            return os.path.abspath(relative_path)
        
        # Try common resource directories
        base_dirs = ['.', 'resources', '../resources', 'assets', '../assets']
        filename = os.path.basename(relative_path)
        
        for base_dir in base_dirs:
            full_path = os.path.join(base_dir, filename)
            if os.path.exists(full_path):
                return os.path.abspath(full_path)
        
        debug_print(f"WARNING: Resource file not found: {relative_path}")
        return relative_path
    except Exception as e:
        debug_print(f"ERROR: Failed to find resource path {relative_path}: {e}")
        return relative_path


def get_save_path(extension: str = ".xlsx") -> Optional[str]:
    """Get save path for reports."""
    try:
        from tkinter import filedialog
        filetypes = [("Excel files", "*.xlsx"), ("All files", "*.*")]
        if extension == ".pptx":
            filetypes = [("PowerPoint files", "*.pptx"), ("All files", "*.*")]
        
        save_path = filedialog.asksaveasfilename(
            defaultextension=extension,
            filetypes=filetypes,
            title=f"Save Report As"
        )
        return save_path if save_path else None
    except ImportError:
        debug_print("WARNING: GUI file dialog not available")
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        return f"report_{timestamp}{extension}"


class HeaderSelectorDialog:
    """Dialog for selecting and ordering report headers."""
    
    def __init__(self, parent=None):
        """Initialize header selector dialog."""
        self.parent = parent
        self.selected_headers = None
        
        # Default headers based on common processing outputs
        self.default_headers = [
            'Sample ID', 'Sample Name', 'Batch ID', 'Test Date',
            'TPM Average (mg)', 'TPM Std Dev (mg)', 'Usage Efficiency (%)',
            'Puff Count', 'Media', 'Viscosity', 'Voltage', 'Resistance', 'Power',
            'Initial Oil Mass (g)', 'Final Oil Mass (g)', 'Oil Mass Consumed (g)',
            'Temperature (°C)', 'Humidity (%)', 'Notes'
        ]
        
        debug_print("HeaderSelectorDialog initialized")
    
    def show(self) -> Optional[List[str]]:
        """Show header selection dialog and return selected headers."""
        try:
            if self.parent:
                # Try to show GUI dialog if tkinter is available
                return self._show_gui_dialog()
            else:
                # Return default headers if no GUI
                debug_print("Using default headers (no GUI available)")
                return self.default_headers[:10]  # Return first 10 default headers
        except Exception as e:
            debug_print(f"ERROR: Header selection failed: {e}")
            return self.default_headers[:10]
    
    def _show_gui_dialog(self) -> Optional[List[str]]:
        """Show GUI dialog for header selection."""
        try:
            import tkinter as tk
            from tkinter import ttk, messagebox
            
            dialog = tk.Toplevel(self.parent)
            dialog.title("Select Report Headers")
            dialog.geometry("600x500")
            dialog.resizable(True, True)
            dialog.transient(self.parent)
            dialog.grab_set()
            
            # Center dialog
            dialog.update_idletasks()
            x = (dialog.winfo_screenwidth() // 2) - (600 // 2)
            y = (dialog.winfo_screenheight() // 2) - (500 // 2)
            dialog.geometry(f"600x500+{x}+{y}")
            
            # Create main frame
            main_frame = ttk.Frame(dialog, padding="10")
            main_frame.grid(row=0, column=0, sticky="nsew")
            
            # Configure grid weights
            dialog.grid_rowconfigure(0, weight=1)
            dialog.grid_columnconfigure(0, weight=1)
            main_frame.grid_rowconfigure(1, weight=1)
            main_frame.grid_columnconfigure(0, weight=1)
            
            # Instructions
            instructions = ttk.Label(main_frame, 
                                   text="Select and reorder headers for your report:",
                                   font=('Arial', 12, 'bold'))
            instructions.grid(row=0, column=0, pady=(0, 10), sticky="w")
            
            # Create listbox with scrollbar
            list_frame = ttk.Frame(main_frame)
            list_frame.grid(row=1, column=0, sticky="nsew", pady=(0, 10))
            list_frame.grid_rowconfigure(0, weight=1)
            list_frame.grid_columnconfigure(0, weight=1)
            
            listbox = tk.Listbox(list_frame, selectmode='extended', font=('Arial', 10))
            scrollbar = ttk.Scrollbar(list_frame, orient="vertical", command=listbox.yview)
            listbox.configure(yscrollcommand=scrollbar.set)
            
            listbox.grid(row=0, column=0, sticky="nsew")
            scrollbar.grid(row=0, column=1, sticky="ns")
            
            # Populate listbox with default headers
            for header in self.default_headers:
                listbox.insert(tk.END, header)
            
            # Select first 8 headers by default
            for i in range(min(8, len(self.default_headers))):
                listbox.selection_set(i)
            
            # Button frame
            button_frame = ttk.Frame(main_frame)
            button_frame.grid(row=2, column=0, pady=(0, 10))
            
            # Buttons for reordering
            ttk.Button(button_frame, text="Move Up", 
                      command=lambda: self._move_selection(listbox, -1)).pack(side='left', padx=5)
            ttk.Button(button_frame, text="Move Down", 
                      command=lambda: self._move_selection(listbox, 1)).pack(side='left', padx=5)
            ttk.Button(button_frame, text="Select All", 
                      command=lambda: listbox.selection_set(0, tk.END)).pack(side='left', padx=5)
            ttk.Button(button_frame, text="Clear All", 
                      command=lambda: listbox.selection_clear(0, tk.END)).pack(side='left', padx=5)
            
            # OK/Cancel buttons
            ok_cancel_frame = ttk.Frame(main_frame)
            ok_cancel_frame.grid(row=3, column=0)
            
            def on_ok():
                selected_indices = listbox.curselection()
                if not selected_indices:
                    messagebox.showwarning("No Selection", "Please select at least one header.")
                    return
                
                self.selected_headers = [listbox.get(i) for i in selected_indices]
                debug_print(f"Selected headers: {self.selected_headers}")
                dialog.destroy()
            
            def on_cancel():
                self.selected_headers = None
                dialog.destroy()
            
            ttk.Button(ok_cancel_frame, text="OK", command=on_ok).pack(side='left', padx=5)
            ttk.Button(ok_cancel_frame, text="Cancel", command=on_cancel).pack(side='left', padx=5)
            
            # Wait for dialog to close
            dialog.wait_window()
            return self.selected_headers
            
        except ImportError:
            debug_print("WARNING: tkinter not available, using default headers")
            return self.default_headers[:10]
        except Exception as e:
            debug_print(f"ERROR: GUI dialog failed: {e}")
            return self.default_headers[:10]
    
    def _move_selection(self, listbox, direction):
        """Move selected items up or down in the listbox."""
        try:
            selection = listbox.curselection()
            if not selection:
                return
            
            items = [listbox.get(i) for i in selection]
            
            # Determine new positions
            if direction == -1:  # Move up
                if selection[0] == 0:
                    return  # Can't move up further
                new_start = selection[0] - 1
            else:  # Move down
                if selection[-1] == listbox.size() - 1:
                    return  # Can't move down further
                new_start = selection[0] + 1
            
            # Remove items
            for i in reversed(selection):
                listbox.delete(i)
            
            # Insert items at new position
            for i, item in enumerate(items):
                listbox.insert(new_start + i, item)
                listbox.selection_set(new_start + i)
                
        except Exception as e:
            debug_print(f"ERROR: Failed to move selection: {e}")


class ReportService:
    """
    Consolidated service for report generation and calculations.
    Handles Excel/PowerPoint report generation, TPM calculations, viscosity analysis,
    and statistical computations.
    """
    
    def __init__(self, file_service=None, database_service=None, calculation_service=None):
        """Initialize the report service."""
        debug_print("Initializing ReportService")
        
        # Service dependencies
        self.file_service = file_service
        self.database_service = database_service
        self.calculation_service = calculation_service
        
        # Report format support
        self.supported_formats = ['excel', 'powerpoint', 'pdf']
        self.template_directory = "templates"
        
        # TPM calculation parameters
        self.default_puff_time = 3.0
        self.default_puffs = 10
        
        # Statistical calculation settings
        self.precision_digits = 3
        self.statistical_threshold = 2
        
        # Formulation database for viscosity calculations
        self.formulation_database = {}
        self.viscosity_models = {}
        self.consolidated_models = {}
        self.base_models = {}
        
        # Report generation caches
        self.header_cache = {}
        self.plot_cache = {}
        self.calculation_cache = {}
        
        # Load existing databases if available
        self._load_formulation_database()
        self._load_viscosity_models()
        
        debug_print("ReportService initialized successfully")
        debug_print(f"Supported formats: {', '.join(self.supported_formats)}")
        debug_print(f"Loaded {len(self.formulation_database)} formulation entries")
        debug_print(f"Loaded {len(self.viscosity_models)} viscosity models")
    
    # ===================== CORE REPORT GENERATION =====================
    
    def generate_test_report(self, selected_sheet: str, sheets: Dict[str, Any], 
                           plot_options: List[str], config: Optional[Dict[str, Any]] = None) -> Tuple[bool, str]:
        """
        Generate an Excel and PowerPoint report for a single test sheet.
        
        Args:
            selected_sheet: Name of the sheet to generate report for
            sheets: Dictionary mapping sheet names to sheet data
            plot_options: List of plot options to include
            config: Optional configuration dictionary
            
        Returns:
            Tuple of (success: bool, message: str)
        """
        debug_print(f"Generating test report for sheet: {selected_sheet}")
        
        try:
            # Validate input
            if selected_sheet not in sheets:
                error_msg = f"Sheet '{selected_sheet}' not found in available sheets"
                debug_print(f"ERROR: {error_msg}")
                return False, error_msg
            
            sheet_info = sheets[selected_sheet]
            if not isinstance(sheet_info, dict) or "data" not in sheet_info:
                error_msg = f"Invalid data structure for sheet '{selected_sheet}'"
                debug_print(f"ERROR: {error_msg}")
                return False, error_msg
            
            # Get header selection
            header_dialog = HeaderSelectorDialog()
            selected_headers = header_dialog.show()
            
            if not selected_headers:
                debug_print("Header selection cancelled")
                return False, "Header selection cancelled"
            
            debug_print(f"Selected headers for test report: {selected_headers}")
            
            # Get save path
            save_path = get_save_path(".xlsx")
            if not save_path:
                return False, "Save cancelled"
            
            # Process sheet data
            data = sheet_info["data"]
            debug_print(f"Processing sheet data with shape: {data.shape}")
            
            # Import processing module for data processing
            try:
                import processing
                process_function = processing.get_processing_function(selected_sheet)
                processed_data, _, full_sample_data = process_function(data)
            except ImportError:
                debug_print("WARNING: processing module not available, using basic processing")
                processed_data = data.copy()
                full_sample_data = data.copy()
            except Exception as e:
                debug_print(f"ERROR: Processing function failed: {e}")
                processed_data = data.copy()
                full_sample_data = data.copy()
            
            if processed_data.empty:
                error_msg = f"No processed data available for sheet {selected_sheet}"
                debug_print(f"ERROR: {error_msg}")
                return False, error_msg
            
            # Reorder data based on selected headers
            processed_data = self._reorder_processed_data(processed_data, selected_headers)
            debug_print(f"Reordered processed data shape: {processed_data.shape}")
            
            # Generate Excel report
            images_to_delete = []
            excel_success = self._generate_excel_report(
                save_path, selected_sheet, processed_data, full_sample_data,
                plot_options, images_to_delete
            )
            
            if not excel_success:
                return False, "Failed to generate Excel report"
            
            # Generate PowerPoint report
            ppt_save_path = save_path.replace('.xlsx', '.pptx')
            ppt_success = self._generate_powerpoint_test_report(
                ppt_save_path, selected_sheet, processed_data, full_sample_data,
                plot_options, images_to_delete
            )
            
            if not ppt_success:
                debug_print("WARNING: PowerPoint report generation failed")
            
            # Cleanup temporary files
            self._cleanup_images(images_to_delete)
            
            success_msg = f"Test report saved successfully to:\nExcel: {save_path}"
            if ppt_success:
                success_msg += f"\nPowerPoint: {ppt_save_path}"
            
            debug_print("Test report generation completed successfully")
            return True, success_msg
            
        except Exception as e:
            error_msg = f"Test report generation failed: {e}"
            debug_print(f"ERROR: {error_msg}")
            traceback.print_exc()
            return False, error_msg
    
    def generate_full_report(self, filtered_sheets: Dict[str, Any], 
                           plot_options: List[str], config: Optional[Dict[str, Any]] = None) -> Tuple[bool, str]:
        """
        Generate a full report (Excel and PowerPoint) for all sheets.
        
        Args:
            filtered_sheets: Dictionary mapping sheet names to sheet info
            plot_options: List of plot options to include
            config: Optional configuration dictionary
            
        Returns:
            Tuple of (success: bool, message: str)
        """
        debug_print(f"Generating full report for {len(filtered_sheets)} sheets")
        
        try:
            # Get header selection
            header_dialog = HeaderSelectorDialog()
            selected_headers = header_dialog.show()
            
            if not selected_headers:
                debug_print("Header selection cancelled")
                return False, "Header selection cancelled"
            
            debug_print(f"Selected headers for full report: {selected_headers}")
            
            # Get save path
            save_path = get_save_path(".xlsx")
            if not save_path:
                return False, "Save cancelled"
            
            ppt_save_path = save_path.replace('.xlsx', '.pptx')
            images_to_delete = []
            total_sheets = len(filtered_sheets)
            processed_count = 0
            
            # Generate Excel report with all sheets
            excel_success = self._generate_full_excel_report(
                save_path, filtered_sheets, plot_options, selected_headers, images_to_delete
            )
            
            if not excel_success:
                return False, "Failed to generate Excel report"
            
            # Generate PowerPoint report
            ppt_success = self._generate_full_powerpoint_report(
                ppt_save_path, filtered_sheets, plot_options, selected_headers, images_to_delete
            )
            
            if not ppt_success:
                debug_print("WARNING: PowerPoint report generation failed")
            
            # Cleanup temporary files
            self._cleanup_images(images_to_delete)
            
            success_msg = f"Full report saved successfully to:\nExcel: {save_path}"
            if ppt_success:
                success_msg += f"\nPowerPoint: {ppt_save_path}"
            
            debug_print("Full report generation completed successfully")
            return True, success_msg
            
        except Exception as e:
            error_msg = f"Full report generation failed: {e}"
            debug_print(f"ERROR: {error_msg}")
            traceback.print_exc()
            return False, error_msg
    
    # ===================== TPM CALCULATIONS =====================
    
    def calculate_tpm_from_weights(self, puffs: pd.Series, before_weights: pd.Series, 
                                   after_weights: pd.Series, test_type: str = "standard") -> pd.Series:
        """
        Calculate TPM (Total Particulate Matter) from weight differences and puffing intervals.
        
        Args:
            puffs: Series of puff counts for each measurement
            before_weights: Series of weights before puffing (g)
            after_weights: Series of weights after puffing (g)
            test_type: Type of test ("standard", "user_simulation", etc.)
            
        Returns:
            Series of TPM values in mg
        """
        debug_print(f"Calculating TPM for {len(puffs)} measurements, test_type: {test_type}")
        
        try:
            # Validate input data
            if puffs.empty or before_weights.empty or after_weights.empty:
                debug_print("WARNING: Empty input data for TPM calculation")
                return pd.Series(dtype=float)
            
            # Ensure all series have the same length
            min_length = min(len(puffs), len(before_weights), len(after_weights))
            puffs_clean = puffs.iloc[:min_length]
            before_clean = before_weights.iloc[:min_length]
            after_clean = after_weights.iloc[:min_length]
            
            # Convert to numeric, handling errors
            puffs_numeric = pd.to_numeric(puffs_clean, errors='coerce')
            before_numeric = pd.to_numeric(before_clean, errors='coerce')
            after_numeric = pd.to_numeric(after_clean, errors='coerce')
            
            # Calculate weight differences (g)
            weight_diff = before_numeric - after_numeric
            
            # Handle invalid puff counts
            puffs_valid = puffs_numeric.copy()
            invalid_puffs = (puffs_valid <= 0) | puffs_valid.isna()
            puffs_valid[invalid_puffs] = self.default_puffs
            
            # Calculate TPM in mg (convert g to mg and normalize by puffs)
            tpm_values = (weight_diff * 1000) / puffs_valid
            
            # Filter out invalid results
            tpm_values[weight_diff < 0] = np.nan  # Negative weight loss
            tpm_values[weight_diff > 1] = np.nan  # Unrealistic weight loss (>1g)
            tpm_values[tpm_values < 0] = np.nan   # Negative TPM
            tpm_values[tpm_values > 100] = np.nan # Unrealistic TPM (>100mg)
            
            debug_print(f"Calculated {(~tpm_values.isna()).sum()} valid TPM values out of {len(tpm_values)}")
            debug_print(f"TPM range: {tpm_values.min():.3f} - {tpm_values.max():.3f} mg")
            
            return tpm_values
            
        except Exception as e:
            debug_print(f"ERROR: TPM calculation failed: {e}")
            return pd.Series(dtype=float)
    
    def calculate_tpm_statistics(self, sample_data: pd.DataFrame, test_type: str = "standard") -> Dict[str, str]:
        """
        Calculate TPM statistics from sample data.
        
        Args:
            sample_data: DataFrame containing puff and weight data
            test_type: Type of test for determining data layout
            
        Returns:
            Dictionary with 'average', 'std_dev', and 'usage_efficiency' keys
        """
        debug_print(f"Calculating TPM statistics for test_type: {test_type}")
        
        try:
            if test_type.lower() in ["user test simulation", "user_simulation"]:
                return self._calculate_user_test_simulation_tpm(sample_data)
            else:
                return self._calculate_standard_tpm(sample_data)
                
        except Exception as e:
            debug_print(f"ERROR: TPM statistics calculation failed: {e}")
            return {'average': 'No data', 'std_dev': 'No data', 'usage_efficiency': ''}
    
    def _calculate_standard_tpm(self, sample_data: pd.DataFrame) -> Dict[str, str]:
        """Calculate TPM statistics for standard test format."""
        try:
            # Standard format: puffs in column A, before weights in B, after weights in C
            # Data starts from row 4 (index 3)
            data_start_row = 3
            
            puffs = pd.to_numeric(sample_data.iloc[data_start_row:, 0], errors='coerce').dropna()
            before_weights = pd.to_numeric(sample_data.iloc[data_start_row:, 1], errors='coerce').dropna()
            after_weights = pd.to_numeric(sample_data.iloc[data_start_row:, 2], errors='coerce').dropna()
            
            if puffs.empty or before_weights.empty or after_weights.empty:
                return {'average': 'No data', 'std_dev': 'No data', 'usage_efficiency': ''}
            
            # Calculate TPM
            tpm_values = self.calculate_tpm_from_weights(puffs, before_weights, after_weights)
            tpm_numeric = pd.to_numeric(tpm_values, errors='coerce').dropna()
            
            if tpm_numeric.empty:
                return {'average': 'No data', 'std_dev': 'No data', 'usage_efficiency': ''}
            
            avg_tpm = round_values(tpm_numeric.mean())
            std_tpm = round_values(tpm_numeric.std())
            
            # Calculate usage efficiency if we have oil mass data
            usage_efficiency = self._calculate_usage_efficiency(sample_data, tpm_numeric)
            
            return {
                'average': f"{avg_tpm:.3f}" if avg_tpm > 0 else "No data",
                'std_dev': f"{std_tpm:.3f}" if std_tpm > 0 else "No data",
                'usage_efficiency': usage_efficiency
            }
            
        except Exception as e:
            debug_print(f"ERROR: Standard TPM calculation failed: {e}")
            return {'average': 'No data', 'std_dev': 'No data', 'usage_efficiency': ''}
    
    def _calculate_user_test_simulation_tpm(self, sample_data: pd.DataFrame) -> Dict[str, str]:
        """Calculate TPM statistics for User Test Simulation format."""
        try:
            # User Test Simulation: puffs in column B, weights in columns C-D
            # Data starts from row 5 (index 4)
            data_start_row = 4
            
            puffs = pd.to_numeric(sample_data.iloc[data_start_row:, 1], errors='coerce').dropna()
            before_weights = pd.to_numeric(sample_data.iloc[data_start_row:, 2], errors='coerce').dropna()
            after_weights = pd.to_numeric(sample_data.iloc[data_start_row:, 3], errors='coerce').dropna()
            
            if puffs.empty or before_weights.empty or after_weights.empty:
                return {'average': 'No data', 'std_dev': 'No data', 'usage_efficiency': ''}
            
            # Calculate TPM
            tpm_values = self.calculate_tpm_from_weights(puffs, before_weights, after_weights, "user_simulation")
            tpm_numeric = pd.to_numeric(tpm_values, errors='coerce').dropna()
            
            if tpm_numeric.empty:
                return {'average': 'No data', 'std_dev': 'No data', 'usage_efficiency': ''}
            
            avg_tpm = round_values(tpm_numeric.mean())
            std_tpm = round_values(tpm_numeric.std())
            
            return {
                'average': f"{avg_tpm:.3f}" if avg_tpm > 0 else "No data",
                'std_dev': f"{std_tpm:.3f}" if std_tpm > 0 else "No data",
                'usage_efficiency': ''  # Not calculated for User Test Simulation
            }
            
        except Exception as e:
            debug_print(f"ERROR: User Test Simulation TPM calculation failed: {e}")
            return {'average': 'No data', 'std_dev': 'No data', 'usage_efficiency': ''}
    
    def _calculate_usage_efficiency(self, sample_data: pd.DataFrame, tpm_values: pd.Series) -> str:
        """Calculate usage efficiency from sample data and TPM values."""
        try:
            # Try to get initial oil mass from header area (H3 position)
            if sample_data.shape[0] > 3 and sample_data.shape[1] > 7:
                initial_oil_mass_val = sample_data.iloc[1, 7]  # Excel H3 = row 1, col 7
                if pd.notna(initial_oil_mass_val):
                    initial_oil_mass = float(initial_oil_mass_val)
                    total_mass_vaporized = tpm_values.sum() / 1000  # Convert mg to g
                    efficiency = (total_mass_vaporized / initial_oil_mass) * 100
                    return f"{round_values(efficiency, 1)}%"
            
            return ""
            
        except Exception as e:
            debug_print(f"WARNING: Could not calculate usage efficiency: {e}")
            return ""
    
    # ===================== VISCOSITY CALCULATIONS =====================
    
    def calculate_viscosity_from_formulation(self, formulation_id: str, temperature: float = 20.0) -> Optional[float]:
        """
        Calculate viscosity for a given formulation at specified temperature.
        
        Args:
            formulation_id: ID of the formulation
            temperature: Temperature in Celsius
            
        Returns:
            Calculated viscosity in cP, or None if calculation fails
        """
        debug_print(f"Calculating viscosity for formulation {formulation_id} at {temperature}°C")
        
        try:
            if formulation_id not in self.formulation_database:
                debug_print(f"WARNING: Formulation {formulation_id} not found in database")
                return None
            
            formulation = self.formulation_database[formulation_id]
            
            # Use existing viscosity model if available
            if formulation_id in self.viscosity_models:
                model = self.viscosity_models[formulation_id]
                viscosity = self._apply_viscosity_model(model, temperature)
                debug_print(f"Calculated viscosity: {viscosity:.2f} cP")
                return viscosity
            
            # Fallback to simple temperature coefficient calculation
            base_viscosity = formulation.get('viscosity', 1.0)
            temp_coefficient = formulation.get('temperature_coefficient', 0.02)
            
            # Apply temperature correction: viscosity decreases with temperature
            viscosity = base_viscosity * (1 - temp_coefficient * (temperature - 20))
            viscosity = max(viscosity, 0.1)  # Minimum viscosity threshold
            
            debug_print(f"Calculated viscosity (fallback): {viscosity:.2f} cP")
            return viscosity
            
        except Exception as e:
            debug_print(f"ERROR: Viscosity calculation failed: {e}")
            return None
    
    def _apply_viscosity_model(self, model: Dict[str, Any], temperature: float) -> float:
        """Apply viscosity model to calculate viscosity at given temperature."""
        try:
            model_type = model.get('type', 'arrhenius')
            
            if model_type == 'arrhenius':
                # Arrhenius model: eta = A * exp(E/(R*T))
                A = model.get('A', 1.0)
                E = model.get('E', 1000.0)  # Activation energy
                R = 8.314  # Gas constant
                T = temperature + 273.15  # Convert to Kelvin
                
                viscosity = A * math.exp(E / (R * T))
                return viscosity
                
            elif model_type == 'polynomial':
                # Polynomial model: eta = a0 + a1*T + a2*T^2 + ...
                coefficients = model.get('coefficients', [1.0])
                viscosity = 0
                for i, coeff in enumerate(coefficients):
                    viscosity += coeff * (temperature ** i)
                return max(viscosity, 0.1)
                
            else:
                debug_print(f"WARNING: Unknown viscosity model type: {model_type}")
                return 1.0
                
        except Exception as e:
            debug_print(f"ERROR: Viscosity model application failed: {e}")
            return 1.0
    
    # ===================== STATISTICAL ANALYSIS =====================
    
    def calculate_statistical_summary(self, data: pd.Series, include_confidence: bool = True) -> Dict[str, float]:
        """
        Calculate comprehensive statistical summary for a data series.
        
        Args:
            data: Pandas series of numeric data
            include_confidence: Whether to include confidence intervals
            
        Returns:
            Dictionary with statistical measures
        """
        debug_print(f"Calculating statistical summary for {len(data)} data points")
        
        try:
            # Convert to numeric and remove NaN values
            numeric_data = pd.to_numeric(data, errors='coerce').dropna()
            
            if len(numeric_data) < self.statistical_threshold:
                debug_print("WARNING: Insufficient data for statistical analysis")
                return {
                    'count': len(numeric_data),
                    'mean': np.nan,
                    'std': np.nan,
                    'min': np.nan,
                    'max': np.nan,
                    'median': np.nan,
                    'cv': np.nan
                }
            
            stats = {
                'count': len(numeric_data),
                'mean': round_values(numeric_data.mean()),
                'std': round_values(numeric_data.std()),
                'min': round_values(numeric_data.min()),
                'max': round_values(numeric_data.max()),
                'median': round_values(numeric_data.median()),
                'cv': round_values((numeric_data.std() / numeric_data.mean()) * 100) if numeric_data.mean() != 0 else np.nan
            }
            
            # Add confidence intervals if requested
            if include_confidence and len(numeric_data) >= 3:
                try:
                    from scipy import stats as scipy_stats
                    confidence_level = 0.95
                    degrees_freedom = len(numeric_data) - 1
                    t_value = scipy_stats.t.ppf((1 + confidence_level) / 2, degrees_freedom)
                    margin_error = t_value * (numeric_data.std() / math.sqrt(len(numeric_data)))
                    
                    stats['ci_lower'] = round_values(numeric_data.mean() - margin_error)
                    stats['ci_upper'] = round_values(numeric_data.mean() + margin_error)
                    stats['margin_error'] = round_values(margin_error)
                    
                except ImportError:
                    debug_print("WARNING: scipy not available for confidence intervals")
                except Exception as e:
                    debug_print(f"WARNING: Could not calculate confidence intervals: {e}")
            
            debug_print(f"Statistical summary: mean={stats['mean']}, std={stats['std']}, cv={stats['cv']}%")
            return stats
            
        except Exception as e:
            debug_print(f"ERROR: Statistical summary calculation failed: {e}")
            return {
                'count': 0,
                'mean': np.nan,
                'std': np.nan,
                'min': np.nan,
                'max': np.nan,
                'median': np.nan,
                'cv': np.nan
            }
    
    def perform_outlier_detection(self, data: pd.Series, method: str = 'iqr', threshold: float = 1.5) -> Tuple[pd.Series, List[int]]:
        """
        Detect outliers in data using specified method.
        
        Args:
            data: Pandas series of numeric data
            method: Method to use ('iqr', 'zscore', 'modified_zscore')
            threshold: Threshold value for outlier detection
            
        Returns:
            Tuple of (cleaned_data, outlier_indices)
        """
        debug_print(f"Performing outlier detection using {method} method")
        
        try:
            numeric_data = pd.to_numeric(data, errors='coerce').dropna()
            
            if len(numeric_data) < 4:
                debug_print("WARNING: Insufficient data for outlier detection")
                return data, []
            
            outlier_indices = []
            
            if method == 'iqr':
                Q1 = numeric_data.quantile(0.25)
                Q3 = numeric_data.quantile(0.75)
                IQR = Q3 - Q1
                lower_bound = Q1 - threshold * IQR
                upper_bound = Q3 + threshold * IQR
                
                outliers_mask = (numeric_data < lower_bound) | (numeric_data > upper_bound)
                outlier_indices = numeric_data[outliers_mask].index.tolist()
                
            elif method == 'zscore':
                z_scores = np.abs((numeric_data - numeric_data.mean()) / numeric_data.std())
                outliers_mask = z_scores > threshold
                outlier_indices = numeric_data[outliers_mask].index.tolist()
                
            elif method == 'modified_zscore':
                median = numeric_data.median()
                mad = np.median(np.abs(numeric_data - median))
                modified_z_scores = 0.6745 * (numeric_data - median) / mad
                outliers_mask = np.abs(modified_z_scores) > threshold
                outlier_indices = numeric_data[outliers_mask].index.tolist()
            
            # Create cleaned data
            cleaned_data = data.copy()
            cleaned_data.iloc[outlier_indices] = np.nan
            
            debug_print(f"Detected {len(outlier_indices)} outliers out of {len(numeric_data)} data points")
            return cleaned_data, outlier_indices
            
        except Exception as e:
            debug_print(f"ERROR: Outlier detection failed: {e}")
            return data, []
    
    # ===================== PRIVATE HELPER METHODS =====================
    
    def _reorder_processed_data(self, processed_data: pd.DataFrame, selected_headers: List[str]) -> pd.DataFrame:
        """Reorder processed data columns based on selected headers."""
        debug_print(f"Reordering columns. Original: {processed_data.columns.tolist()}")
        debug_print(f"Selected headers: {selected_headers}")
        
        try:
            # Create new DataFrame with selected columns in order
            reordered_data = pd.DataFrame()
            
            for header in selected_headers:
                if header in processed_data.columns:
                    reordered_data[header] = processed_data[header]
                    debug_print(f"Added column: {header}")
                else:
                    debug_print(f"Missing column: {header}, adding empty")
                    reordered_data[header] = [''] * len(processed_data)
            
            debug_print(f"Final reordered columns: {reordered_data.columns.tolist()}")
            return reordered_data
            
        except Exception as e:
            debug_print(f"ERROR: Column reordering failed: {e}")
            return processed_data
    
    def _generate_excel_report(self, save_path: str, sheet_name: str, processed_data: pd.DataFrame,
                             full_sample_data: pd.DataFrame, plot_options: List[str], 
                             images_to_delete: List[str]) -> bool:
        """Generate Excel report for a single sheet."""
        try:
            debug_print(f"Generating Excel report: {save_path}")
            
            with pd.ExcelWriter(save_path, engine='xlsxwriter') as writer:
                # Write processed data
                processed_data.astype(str).replace([pd.NA], '')
                processed_data.to_excel(writer, sheet_name=sheet_name, index=False)
                
                # Add plots if this is a plotting sheet
                if self._is_plotting_sheet(sheet_name):
                    self._add_plots_to_excel(writer, sheet_name, full_sample_data, 
                                           images_to_delete, plot_options)
            
            debug_print(f"Excel report saved: {save_path}")
            return True
            
        except Exception as e:
            debug_print(f"ERROR: Excel report generation failed: {e}")
            traceback.print_exc()
            return False
    
    def _generate_full_excel_report(self, save_path: str, filtered_sheets: Dict[str, Any],
                                  plot_options: List[str], selected_headers: List[str],
                                  images_to_delete: List[str]) -> bool:
        """Generate Excel report for all sheets."""
        try:
            debug_print(f"Generating full Excel report: {save_path}")
            
            with pd.ExcelWriter(save_path, engine='xlsxwriter') as writer:
                for sheet_name, sheet_info in filtered_sheets.items():
                    try:
                        if not isinstance(sheet_info, dict) or "data" not in sheet_info:
                            debug_print(f"Skipping sheet '{sheet_name}': Invalid data structure")
                            continue
                        
                        data = sheet_info["data"]
                        
                        # Process data
                        try:
                            import processing
                            process_function = processing.get_processing_function(sheet_name)
                            processed_data, _, full_sample_data = process_function(data)
                        except:
                            processed_data = data.copy()
                            full_sample_data = data.copy()
                        
                        if processed_data.empty:
                            debug_print(f"Skipping sheet '{sheet_name}': Empty processed data")
                            continue
                        
                        # Apply header reordering
                        processed_data = self._reorder_processed_data(processed_data, selected_headers)
                        
                        # Write to Excel
                        processed_data.astype(str).replace([pd.NA], '')
                        processed_data.to_excel(writer, sheet_name=sheet_name, index=False)
                        
                        # Add plots if applicable
                        if self._is_plotting_sheet(sheet_name):
                            self._add_plots_to_excel(writer, sheet_name, full_sample_data,
                                                   images_to_delete, plot_options)
                        
                    except Exception as e:
                        debug_print(f"ERROR: Processing sheet '{sheet_name}': {e}")
                        continue
            
            debug_print(f"Full Excel report saved: {save_path}")
            return True
            
        except Exception as e:
            debug_print(f"ERROR: Full Excel report generation failed: {e}")
            traceback.print_exc()
            return False
    
    def _generate_powerpoint_test_report(self, ppt_save_path: str, sheet_name: str,
                                       processed_data: pd.DataFrame, full_sample_data: pd.DataFrame,
                                       plot_options: List[str], images_to_delete: List[str]) -> bool:
        """Generate PowerPoint report for a single test sheet."""
        if not POWERPOINT_AVAILABLE:
            debug_print("WARNING: PowerPoint functionality not available")
            return False
        
        try:
            debug_print(f"Generating PowerPoint test report: {ppt_save_path}")
            
            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            
            # Create main slide
            main_slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # Add background and logo
            self._setup_slide_background(main_slide, prs)
            
            # Add title
            self._add_slide_title(main_slide, sheet_name)
            
            # Add table and plots
            is_plotting = self._is_plotting_sheet(sheet_name)
            table_width = Inches(8.07) if is_plotting else Inches(13.03)
            
            self._add_table_to_slide(main_slide, processed_data, table_width, is_plotting)
            
            if is_plotting:
                try:
                    import processing
                    valid_plot_options = processing.get_valid_plot_options(plot_options, full_sample_data)
                    if valid_plot_options:
                        self._add_plots_to_slide(main_slide, sheet_name, full_sample_data,
                                               valid_plot_options, images_to_delete)
                except:
                    debug_print("WARNING: Could not add plots to PowerPoint slide")
            
            prs.save(ppt_save_path)
            debug_print(f"PowerPoint test report saved: {ppt_save_path}")
            return True
            
        except Exception as e:
            debug_print(f"ERROR: PowerPoint test report generation failed: {e}")
            traceback.print_exc()
            return False
    
    def _generate_full_powerpoint_report(self, ppt_save_path: str, filtered_sheets: Dict[str, Any],
                                       plot_options: List[str], selected_headers: List[str],
                                       images_to_delete: List[str]) -> bool:
        """Generate PowerPoint report for all sheets."""
        if not POWERPOINT_AVAILABLE:
            debug_print("WARNING: PowerPoint functionality not available")
            return False
        
        try:
            debug_print(f"Generating full PowerPoint report: {ppt_save_path}")
            
            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            
            for sheet_name, sheet_info in filtered_sheets.items():
                try:
                    if not isinstance(sheet_info, dict) or "data" not in sheet_info:
                        debug_print(f"Skipping sheet '{sheet_name}': Invalid data structure")
                        continue
                    
                    data = sheet_info["data"]
                    
                    # Process data
                    try:
                        import processing
                        process_function = processing.get_processing_function(sheet_name)
                        processed_data, _, full_sample_data = process_function(data)
                    except:
                        processed_data = data.copy()
                        full_sample_data = data.copy()
                    
                    if processed_data.empty:
                        debug_print(f"Skipping sheet '{sheet_name}': Empty processed data")
                        continue
                    
                    # Apply header reordering
                    processed_data = self._reorder_processed_data(processed_data, selected_headers)
                    
                    # Create slide for this sheet
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    self._setup_slide_background(slide, prs)
                    self._add_slide_title(slide, sheet_name)
                    
                    # Add content based on sheet type
                    is_plotting = self._is_plotting_sheet(sheet_name)
                    table_width = Inches(8.07) if is_plotting else Inches(13.03)
                    
                    self._add_table_to_slide(slide, processed_data, table_width, is_plotting)
                    
                    if is_plotting:
                        try:
                            import processing
                            valid_plot_options = processing.get_valid_plot_options(plot_options, full_sample_data)
                            if valid_plot_options:
                                self._add_plots_to_slide(slide, sheet_name, full_sample_data,
                                                       valid_plot_options, images_to_delete)
                        except:
                            debug_print(f"WARNING: Could not add plots for sheet '{sheet_name}'")
                    
                except Exception as e:
                    debug_print(f"ERROR: Processing PowerPoint slide for '{sheet_name}': {e}")
                    continue
            
            prs.save(ppt_save_path)
            debug_print(f"Full PowerPoint report saved: {ppt_save_path}")
            return True
            
        except Exception as e:
            debug_print(f"ERROR: Full PowerPoint report generation failed: {e}")
            traceback.print_exc()
            return False
    
    def _setup_slide_background(self, slide, prs):
        """Setup slide background and logo."""
        try:
            background_path = get_resource_path("resources/ccell_background.png")
            if os.path.exists(background_path):
                slide.shapes.add_picture(background_path, Inches(0), Inches(0),
                                       width=prs.slide_width, height=prs.slide_height)
            
            logo_path = get_resource_path("resources/ccell_logo_full.png")
            if os.path.exists(logo_path):
                slide.shapes.add_picture(logo_path, Inches(11.21), Inches(0.43),
                                       width=Inches(1.57), height=Inches(0.53))
        except Exception as e:
            debug_print(f"WARNING: Could not setup slide background: {e}")
    
    def _add_slide_title(self, slide, title_text: str):
        """Add title to slide."""
        try:
            title_shape = slide.shapes.add_textbox(Inches(0.45), Inches(-0.04),
                                                 Inches(10.72), Inches(0.64))
            text_frame = title_shape.text_frame
            text_frame.clear()
            p = text_frame.add_paragraph()
            p.text = title_text
            p.alignment = PP_ALIGN.LEFT
            run = p.runs[0]
            run.font.name = "Montserrat"
            run.font.size = Pt(32)
            run.font.bold = True
        except Exception as e:
            debug_print(f"WARNING: Could not add slide title: {e}")
    
    def _add_table_to_slide(self, slide, data: pd.DataFrame, table_width, is_plotting: bool):
        """Add data table to slide."""
        try:
            if data.empty:
                debug_print("WARNING: No data to add to table")
                return
            
            # Table positioning
            table_left = Inches(0.58)
            table_top = Inches(0.93)
            table_height = Inches(4.18)
            
            # Limit columns to fit on slide
            max_columns = 8 if is_plotting else 12
            display_data = data.iloc[:, :max_columns] if len(data.columns) > max_columns else data
            
            # Create table
            rows = min(len(display_data) + 1, 20)  # +1 for header, max 20 rows
            cols = len(display_data.columns)
            
            table = slide.shapes.add_table(rows, cols, table_left, table_top,
                                         table_width, table_height).table
            
            # Add headers
            for col_idx, column_name in enumerate(display_data.columns):
                cell = table.cell(0, col_idx)
                cell.text = str(column_name)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(79, 129, 189)  # Blue header
                
                # Header text formatting
                if cell.text_frame and cell.text_frame.paragraphs:
                    paragraph = cell.text_frame.paragraphs[0]
                    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                    run.font.color.rgb = RGBColor(255, 255, 255)  # White text
                    run.font.bold = True
                    run.font.size = Pt(10)
            
            # Add data rows
            for row_idx in range(1, min(rows, len(display_data) + 1)):
                data_row_idx = row_idx - 1
                if data_row_idx < len(display_data):
                    for col_idx, column_name in enumerate(display_data.columns):
                        cell = table.cell(row_idx, col_idx)
                        cell_value = display_data.iloc[data_row_idx, col_idx]
                        cell.text = str(cell_value) if pd.notna(cell_value) else ""
                        
                        # Data text formatting
                        if cell.text_frame and cell.text_frame.paragraphs:
                            paragraph = cell.text_frame.paragraphs[0]
                            run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                            run.font.size = Pt(9)
                            run.font.color.rgb = RGBColor(0, 0, 0)  # Black text
            
            debug_print(f"Added table with {rows} rows and {cols} columns")
            
        except Exception as e:
            debug_print(f"ERROR: Could not add table to slide: {e}")
    
    def _add_plots_to_slide(self, slide, sheet_name: str, full_sample_data: pd.DataFrame,
                          valid_plot_options: List[str], images_to_delete: List[str]):
        """Add plots to slide."""
        if not PLOTTING_AVAILABLE:
            debug_print("WARNING: Plotting functionality not available")
            return
        
        try:
            # Plot positioning - cascading layout
            plot_start_left = Inches(0.02)
            plot_top = Inches(5.26)
            plot_height = Inches(2.0)
            current_left = plot_start_left
            
            debug_print(f"Adding {len(valid_plot_options)} plots to slide")
            
            # Convert data to numeric for plotting
            numeric_data = full_sample_data.apply(pd.to_numeric, errors='coerce')
            if numeric_data.isna().all(axis=0).all():
                debug_print(f"No numeric data available for plotting in sheet '{sheet_name}'")
                return
            
            # Determine sample configuration
            is_user_test_simulation = sheet_name in ["User Test Simulation", "User Simulation Test"]
            num_columns_per_sample = 8 if is_user_test_simulation else 12
            
            for i, plot_option in enumerate(valid_plot_options):
                plot_image_path = f"{sheet_name}_{plot_option}_plot.png"
                
                try:
                    # Create plot using processing module
                    import processing
                    fig, sample_names_returned = processing.plot_all_samples(
                        numeric_data, plot_option, num_columns_per_sample, None
                    )
                    
                    # Calculate plot width based on aspect ratio
                    if is_user_test_simulation and hasattr(fig, 'is_split_plot') and fig.is_split_plot:
                        # Split plots are wider
                        plot_width = Inches(3.5)
                    else:
                        # Standard plots
                        plot_width = Inches(2.29)
                    
                    # Save plot
                    plt.savefig(plot_image_path, dpi=150, bbox_inches='tight', 
                              facecolor='white', edgecolor='none')
                    
                    # Add to slide
                    slide.shapes.add_picture(plot_image_path, current_left, plot_top,
                                           plot_width, plot_height)
                    
                    # Update position for next plot
                    current_left += plot_width + Inches(0.1)  # Small gap between plots
                    
                    # Add to cleanup list
                    if plot_image_path not in images_to_delete:
                        images_to_delete.append(plot_image_path)
                    
                    plt.close(fig)
                    debug_print(f"Added plot: {plot_option}")
                    
                except Exception as e:
                    debug_print(f"ERROR: Could not create plot '{plot_option}': {e}")
                    continue
            
        except Exception as e:
            debug_print(f"ERROR: Could not add plots to slide: {e}")
    
    def _add_plots_to_excel(self, writer, sheet_name: str, full_sample_data: pd.DataFrame,
                          images_to_delete: List[str], plot_options: List[str]):
        """Add plots to Excel worksheet."""
        if not PLOTTING_AVAILABLE:
            debug_print("WARNING: Plotting functionality not available")
            return
        
        try:
            # Get worksheet
            workbook = writer.book
            worksheet = writer.sheets[sheet_name]
            
            # Convert data to numeric for plotting
            numeric_data = full_sample_data.apply(pd.to_numeric, errors='coerce')
            if numeric_data.isna().all(axis=0).all():
                debug_print(f"No numeric data available for plotting in sheet '{sheet_name}'")
                return
            
            # Get valid plot options
            try:
                import processing
                valid_plot_options = processing.get_valid_plot_options(plot_options, full_sample_data)
            except:
                valid_plot_options = plot_options
            
            # Determine sample configuration
            is_user_test_simulation = sheet_name in ["User Test Simulation", "User Simulation Test"]
            num_columns_per_sample = 8 if is_user_test_simulation else 12
            
            # Add plots below the data table
            row_offset = len(full_sample_data) + 5  # Start below data with some spacing
            
            for i, plot_option in enumerate(valid_plot_options):
                plot_image_path = f"{sheet_name}_{plot_option}_excel_plot.png"
                
                try:
                    # Create plot
                    import processing
                    fig, sample_names_returned = processing.plot_all_samples(
                        numeric_data, plot_option, num_columns_per_sample, None
                    )
                    
                    # Save plot with appropriate size for Excel
                    plt.savefig(plot_image_path, dpi=100, bbox_inches='tight',
                              facecolor='white', edgecolor='none')
                    
                    # Insert plot into Excel
                    worksheet.insert_image(f'A{row_offset + i * 25}', plot_image_path)
                    
                    # Add to cleanup list
                    if plot_image_path not in images_to_delete:
                        images_to_delete.append(plot_image_path)
                    
                    plt.close(fig)
                    debug_print(f"Added Excel plot: {plot_option}")
                    
                except Exception as e:
                    debug_print(f"ERROR: Could not create Excel plot '{plot_option}': {e}")
                    continue
            
        except Exception as e:
            debug_print(f"ERROR: Could not add plots to Excel: {e}")
    
    def _is_plotting_sheet(self, sheet_name: str) -> bool:
        """Determine if a sheet should include plots."""
        plotting_sheet_names = [
            "Test Plan", "Lifetime Test", "Device Life Test", "Quick Screening",
            "Aerosol Temperature", "User Test", "Horizontal Puffing", "Extended Test",
            "User Test Simulation", "User Simulation Test"
        ]
        return sheet_name in plotting_sheet_names
    
    def _cleanup_images(self, images_to_delete: List[str]):
        """Clean up temporary image files."""
        debug_print(f"Cleaning up {len(images_to_delete)} temporary image files")
        
        for image_path in images_to_delete:
            try:
                if os.path.exists(image_path):
                    os.remove(image_path)
                    debug_print(f"Removed temporary file: {image_path}")
            except Exception as e:
                debug_print(f"WARNING: Could not remove temporary file {image_path}: {e}")
    
    # ===================== DATABASE OPERATIONS =====================
    
    def _load_formulation_database(self):
        """Load formulation database from file."""
        try:
            db_path = "data/formulation_database.json"
            if os.path.exists(db_path):
                with open(db_path, 'r') as f:
                    self.formulation_database = json.load(f)
                debug_print(f"Loaded formulation database with {len(self.formulation_database)} entries")
            else:
                debug_print("No formulation database found, starting with empty database")
                self.formulation_database = {}
        except Exception as e:
            debug_print(f"ERROR: Could not load formulation database: {e}")
            self.formulation_database = {}
    
    def _load_viscosity_models(self):
        """Load viscosity models from file."""
        try:
            models_path = "data/viscosity_models.json"
            if os.path.exists(models_path):
                with open(models_path, 'r') as f:
                    self.viscosity_models = json.load(f)
                debug_print(f"Loaded viscosity models with {len(self.viscosity_models)} entries")
            else:
                debug_print("No viscosity models found, starting with empty models")
                self.viscosity_models = {}
        except Exception as e:
            debug_print(f"ERROR: Could not load viscosity models: {e}")
            self.viscosity_models = {}
    
    def save_formulation_database(self):
        """Save formulation database to file."""
        try:
            os.makedirs("data", exist_ok=True)
            db_path = "data/formulation_database.json"
            with open(db_path, 'w') as f:
                json.dump(self.formulation_database, f, indent=2)
            debug_print(f"Saved formulation database to {db_path}")
        except Exception as e:
            debug_print(f"ERROR: Could not save formulation database: {e}")
    
    def save_viscosity_models(self):
        """Save viscosity models to file."""
        try:
            os.makedirs("data", exist_ok=True)
            models_path = "data/viscosity_models.json"
            with open(models_path, 'w') as f:
                json.dump(self.viscosity_models, f, indent=2)
            debug_print(f"Saved viscosity models to {models_path}")
        except Exception as e:
            debug_print(f"ERROR: Could not save viscosity models: {e}")
    
    # ===================== PUBLIC API METHODS =====================
    
    def get_supported_calculations(self) -> List[str]:
        """Get list of supported calculation types."""
        return [
            'tpm_standard', 'tpm_user_simulation', 'viscosity_arrhenius',
            'viscosity_polynomial', 'statistical_summary', 'outlier_detection'
        ]
    
    def get_report_formats(self) -> List[str]:
        """Get list of supported report formats."""
        formats = ['excel']
        if POWERPOINT_AVAILABLE:
            formats.append('powerpoint')
        return formats
    
    def validate_data_for_calculations(self, data: pd.DataFrame, calculation_type: str) -> Tuple[bool, str]:
        """Validate data for specific calculation types."""
        try:
            if data.empty:
                return False, "Data is empty"
            
            if calculation_type in ['tpm_standard', 'tpm_user_simulation']:
                # Check for required columns
                required_cols = 3  # At least puffs, before_weight, after_weight
                if data.shape[1] < required_cols:
                    return False, f"Insufficient columns for TPM calculation (need {required_cols}, got {data.shape[1]})"
                
                # Check for numeric data in key columns
                numeric_cols = pd.to_numeric(data.iloc[:, :required_cols], errors='coerce')
                if numeric_cols.isna().all().any():
                    return False, "Required columns contain no numeric data"
            
            elif calculation_type == 'statistical_summary':
                # Check for at least some numeric data
                numeric_data = data.apply(pd.to_numeric, errors='coerce')
                if numeric_data.isna().all().all():
                    return False, "No numeric data found for statistical analysis"
            
            return True, "Data validation passed"
            
        except Exception as e:
            return False, f"Data validation failed: {e}"
    
    def get_calculation_parameters(self) -> Dict[str, Any]:
        """Get current calculation parameters."""
        return {
            'default_puff_time': self.default_puff_time,
            'default_puffs': self.default_puffs,
            'precision_digits': self.precision_digits,
            'statistical_threshold': self.statistical_threshold
        }
    
    def update_calculation_parameters(self, parameters: Dict[str, Any]):
        """Update calculation parameters."""
        for key, value in parameters.items():
            if hasattr(self, key):
                setattr(self, key, value)
                debug_print(f"Updated parameter {key} = {value}")
            else:
                debug_print(f"WARNING: Unknown parameter {key}")