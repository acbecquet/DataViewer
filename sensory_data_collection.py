"""
sensory_data_collection.py
Sensory Data Collection Window for DataViewer Application
Developed by Charlie Becquet
"""

import tkinter as tk
from tkinter import ttk, messagebox, filedialog
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from matplotlib.backends.backend_tkagg import FigureCanvasTkAgg
from matplotlib.patches import Circle
import math
from datetime import datetime
import json
import os
from utils import APP_BACKGROUND_COLOR, BUTTON_COLOR, FONT, debug_print

from sensory_ml_training import SensoryMLTrainer, SensoryAIProcessor


def _lazy_import_cv2():
    """Lazy import opencv."""
    try:
        import cv2
        debug_print("TIMING: Lazy loaded cv2 for image processing")
        return cv2
    except ImportError as e:
        debug_print(f"Error importing cv2: {e}")
        messagebox.showerror("Missing Dependency", 
                            "OpenCV is required for image processing.\nPlease install: pip install opencv-python")
        return None

def _lazy_import_pil():
    """Lazy import PIL."""
    try:
        from PIL import Image, ImageTk
        debug_print("TIMING: Lazy loaded PIL for image processing")
        return Image, ImageTk
    except ImportError as e:
        debug_print(f"Error importing PIL: {e}")
        messagebox.showerror("Missing Dependency", 
                            "PIL is required for image processing.\nPlease install: pip install Pillow")
        return None, None

def _lazy_import_sklearn():
    """Lazy import scikit-learn."""
    try:
        from sklearn.cluster import DBSCAN
        debug_print("TIMING: Lazy loaded sklearn for image processing")
        return DBSCAN
    except ImportError as e:
        debug_print(f"Error importing sklearn: {e}")
        messagebox.showerror("Missing Dependency", 
                            "Scikit-learn is required for advanced image processing.\nPlease install: pip install scikit-learn")
        return None

def _lazy_import_pytesseract():
    """Lazy import pytesseract."""
    try:
        import pytesseract
        debug_print("TIMING: Lazy loaded pytesseract for OCR")
        return pytesseract
    except ImportError as e:
        debug_print(f"Error importing pytesseract: {e}")
        messagebox.showerror("Missing Dependency", 
                            "Pytesseract is required for text recognition.\nPlease install: pip install pytesseract")
        return None


class SensoryDataCollectionWindow:
    """Main window for sensory data collection and visualization."""
    
    def __init__(self, parent, close_callback=None):
        self.parent = parent
        self.close_callback = close_callback
        self.window = None
        self.data = {}
        self.sessions = {}  # {'session_id': {'header': {}, 'samples': {}, 'timestamp': '', 'source_image': ''}}
        self.current_session_id = None
        self.session_counter = 1
        self.samples = {}
        self.current_sample = None
        debug_print("DEBUG: Initialized session-based data structure")
        debug_print(f"DEBUG: self.sessions = {self.sessions}")
        debug_print(f"DEBUG: self.current_session_id = {self.current_session_id}")
        debug_print(f"DEBUG: self.session_counter = {self.session_counter}")
        # Sensory metrics (5 attributes)
        self.metrics = [
            "Burnt Taste",
            "Vapor Volume", 
            "Overall Flavor",
            "Smoothness",
            "Overall Liking"
        ]

        self.current_mode = "collection"
        self.all_sessions_data = {}
        self.average_samples = {}
        debug_print("Initialized dual-mode functionality")

        self.ml_trainer = SensoryMLTrainer(self)
        self.ai_processor = SensoryAIProcessor(self)
        debug_print("DEBUG: Initialized ML trainer and AI processor")
        
        # Header data fields
        self.header_fields = [
            "Assessor Name", 
            "Media",
            "Puff Length",
            "Date"
        ]
        debug_print("Updated header fields to 4 essential fields only")
        debug_print(f"New header fields: {self.header_fields}")
        # SOP text
        self.sop_text = """
SENSORY EVALUATION STANDARD OPERATING PROCEDURE

1. PREPARATION:
   - Ensure all cartridges are at room temperature
   - Use clean, odor-free environment, ideally in a fume hood

2. EVALUATION:
   - Take 2-3 moderate puffs per sample
   - Be sure to compare flavor with original oil if available
   - Rate each attribute on 1-9 scale (1=poor, 9=excellent)

3. SCALE INTERPRETATION:
   - Burnt Taste: 1=Very Burnt, 9=No Burnt Taste
   - Vapor Volume: 1=Very Low, 9=Very High
   - Overall Flavor: 1=Poor, 9=Excellent
   - Smoothness: 1=Very Harsh, 9=Very Smooth
   - Overall Liking: 1=Dislike Extremely, 9=Like Extremely

4. NOTES:
   - Record any unusual observations
   - Note any technical issues with samples
        """
    # Lazy import functions for image processing
    

    def show(self):
        """Create and display the sensory data collection window with optimized sizing."""
        self.window = tk.Toplevel(self.parent)
        self.window.title("Sensory Data Collection")
        self.window.resizable(True, True)
    
        # Start with a reasonable initial size (will be optimized after content loads)
        initial_width = 1000
        initial_height = 600
        self.window.geometry(f"{initial_width}x{initial_height}")
        self.window.configure(bg=APP_BACKGROUND_COLOR)
        #self.window.transient(self.parent)
    
        debug_print(f"DEBUG: Initial window size set to {initial_width}x{initial_height}")

        self.window.protocol("WM_DELETE_WINDOW", self.on_window_close)
    
        # Create main layout (this will trigger size optimization)
        self.setup_layout()
        self.setup_menu()       
        self.center_window()
        
    def on_window_close(self):
        """Handle window close event and call the callback if provided."""
        debug_print("DEBUG: Sensory window close event triggered")
    
        # Check for unsaved changes and handle auto-save if needed
        try:
            # Your existing save logic here if any
            debug_print("DEBUG: Performing cleanup before window close")
        
            # Close the window
            self.window.destroy()
            debug_print("DEBUG: Sensory window destroyed")
        
            # Call the callback to restore main window
            if self.close_callback:
                debug_print("DEBUG: Calling close callback to restore main window")
                self.close_callback()
            else:
                debug_print("DEBUG: No close callback provided")
            
        except Exception as e:
            debug_print(f"DEBUG: Error during window close: {e}")
            # Still try to close and call callback even if there's an error
            try:
                self.window.destroy()
                if self.close_callback:
                    self.close_callback()
            except:
                pass

    def center_window(self):
        """Center the window on screen."""
        self.window.update_idletasks()
        width = self.window.winfo_width()
        height = self.window.winfo_height()
        x = (self.window.winfo_screenwidth() // 2) - (width // 2)
        y = (self.window.winfo_screenheight() // 2) - (height // 2)
        self.window.geometry(f"{width}x{height}+{x}+{y}")
        
    def setup_menu(self):
        """Create the menu bar."""
        menubar = tk.Menu(self.window)
        self.window.config(menu=menubar)
        
        # File menu
        file_menu = tk.Menu(menubar, tearoff=0)
        file_menu.add_command(label="New Session", command=self.add_new_session)
        file_menu.add_command(label="Load Session", command=self.load_session)
        file_menu.add_command(label="Save Session", command=self.save_session)
        file_menu.add_separator()
        file_menu.add_command(label="Merge Sessions from Files", command=self.merge_sessions_from_files)
        file_menu.add_separator()
        file_menu.add_command(label="Load from Image (ML)", command=self.ai_processor.load_from_image_enhanced) #add this back once new function added
        file_menu.add_command(label="Load with AI (Claude)", command=self.ai_processor.load_from_image_with_ai)
        file_menu.add_command(label="Batch Process with AI", command=self.ai_processor.batch_process_with_ai)
        file_menu.add_separator()
        file_menu.add_command(label="Export to Excel", command=self.export_to_excel)
        file_menu.add_separator()
        file_menu.add_command(label="Close", command=self.window.destroy)
        menubar.add_cascade(label="File", menu=file_menu)
        
        # Sample menu
        sample_menu = tk.Menu(menubar, tearoff=0)
        sample_menu.add_command(label="Add Sample", command=self.add_sample)
        sample_menu.add_command(label="Remove Sample", command=self.remove_sample)
        sample_menu.add_command(label="Clear Sample Data", command=self.clear_current_sample)
        sample_menu.add_separator()  # NEW
        sample_menu.add_command(label="Rename Current Sample", command=self.rename_current_sample)  # NEW
        sample_menu.add_command(label="Batch Rename Samples", command=self.batch_rename_samples)  # NEW
        menubar.add_cascade(label="Sample", menu=sample_menu)

        # Export menu
        export_menu = tk.Menu(menubar, tearoff=0)
        export_menu.add_command(label="Save Plot as Image", command=self.save_plot_as_image)
        export_menu.add_command(label="Save Table as Image", command=self.save_table_as_image)
        export_menu.add_separator()
        export_menu.add_command(label="Generate PowerPoint Report", command=self.generate_powerpoint_report)
        menubar.add_cascade(label="Export", menu=export_menu)
        
        # View menu
        view_menu = tk.Menu(menubar, tearoff=0)
        view_menu.add_command(label="Show SOP", command=self.show_sop)
        view_menu.add_command(label="Refresh Plot", command=self.update_plot)
        menubar.add_cascade(label="View", menu=view_menu)
        
        ml_menu = tk.Menu(menubar, tearoff=0)
    
       # Core ML workflow - no training structure creation (separate scripts handle this)
        ml_menu.add_command(label="Check Enhanced Data Balance", command=self.ml_trainer.check_enhanced_data_balance)
        ml_menu.add_command(label="Train Enhanced Model", command=self.ml_trainer.train_enhanced_model)
        ml_menu.add_separator()

        # Testing and validation
        ml_menu.add_command(label="Test Enhanced Model", command=self.ml_trainer.test_enhanced_model)
        ml_menu.add_command(label="Validate Model Performance", command=self.ml_trainer.validate_enhanced_performance)
        ml_menu.add_separator()

        # Configuration management
        ml_menu.add_command(label="Update Processor Configuration", command=self.ml_trainer.update_processor_config)

        menubar.add_cascade(label="Enhanced ML", menu=ml_menu)

    def on_window_resize(self, event):
        """Handle general window resize events with dynamic sash positioning."""
        # Only process if this is the main window resize, not child widgets
        if event.widget != self.window:
            return
    
        debug_print("DEBUG: === GENERAL WINDOW RESIZE HANDLER ===")
    
        # Get current window dimensions
        window_width = self.window.winfo_width()
        window_height = self.window.winfo_height()
        debug_print(f"DEBUG: Window resized to: {window_width}x{window_height}px")
    
        # Update paned window sash position proportionally
        if hasattr(self, 'main_paned'):
            # Force geometry update first
            self.main_paned.update_idletasks()
        
            left_panel_proportion = 0.3  # 30% for left panel, 70% for right panel, this is updated later
            new_sash_position = int(window_width * left_panel_proportion)
        
            # Apply minimum and maximum constraints to keep usable
            min_left_width = 350  # Minimum space for left panel functionality
            max_left_width = window_width - 400  # Leave at least 400px for right panel
        
            new_sash_position = max(min_left_width, min(new_sash_position, max_left_width))
        
            debug_print(f"DEBUG: Setting proportional sash position to: {new_sash_position}px ({left_panel_proportion*100:.0f}% of window width)")
        
            # Update the sash position
            try:
                self.main_paned.sash_place(0, new_sash_position, 0)
                debug_print("DEBUG: Sash position updated successfully")
            except Exception as e:
                debug_print(f"DEBUG: Error updating sash position: {e}")
    
        # Force all frames to update their geometry
        debug_print("DEBUG: Forcing frame geometry updates...")
        self.window.update_idletasks()
    
        # Equalize canvas heights
        self.equalize_canvas_heights()
    
        # Force right canvas to update its size
        if hasattr(self, 'right_canvas'):
            self.right_canvas.update_idletasks()
            # Trigger the canvas configure event to update interior frame
            self.right_canvas.event_generate('<Configure>')
    
        # Trigger plot-specific resize with a slight delay to allow frame updates
        if hasattr(self, 'on_window_resize_plot'):
            # Add a small delay to let the sash repositioning complete
            self.window.after(50, lambda: self.on_window_resize_plot(event))
    
        debug_print("DEBUG: === END GENERAL WINDOW RESIZE ===")

    def merge_sessions_from_files(self):
        """Merge multiple session JSON files into a new session."""
    
        debug_print("DEBUG: Starting merge sessions from files")
    
        # Select multiple JSON files
        filenames = filedialog.askopenfilenames(
            filetypes=[("JSON files", "*.json"), ("All files", "*.*")],
            title="Select Session Files to Merge (Hold Ctrl to select multiple)"
        )
    
        if not filenames or len(filenames) < 2:
            messagebox.showinfo("Insufficient Files", 
                              "Please select at least 2 session files to merge.")
            return
    
        debug_print(f"DEBUG: Selected {len(filenames)} files for merging")
    
        # Load and validate all session files
        loaded_sessions = {}
        failed_files = []
    
        for filename in filenames:
            try:
                with open(filename, 'r') as f:
                    session_data = json.load(f)
            
                # Validate session data structure
                if self.validate_session_data(session_data):
                    session_name = os.path.splitext(os.path.basename(filename))[0]
                    loaded_sessions[session_name] = {
                        'file_path': filename,
                        'data': session_data
                    }
                    debug_print(f"DEBUG: Successfully loaded session from {filename}")
                else:
                    failed_files.append(filename)
                    debug_print(f"DEBUG: Invalid session format in {filename}")
                
            except Exception as e:
                failed_files.append(filename)
                debug_print(f"DEBUG: Failed to load {filename}: {e}")
    
        if not loaded_sessions:
            messagebox.showerror("Load Error", 
                               "No valid session files could be loaded.\n"
                               "Ensure files are in the correct JSON format.")
            return
    
        if failed_files:
            failed_list = '\n'.join([os.path.basename(f) for f in failed_files])
            messagebox.showwarning("Some Files Failed", 
                                 f"Failed to load {len(failed_files)} files:\n{failed_list}\n\n"
                                 f"Continuing with {len(loaded_sessions)} valid files.")
    
        # Show merge configuration dialog
        self.show_merge_sessions_dialog(loaded_sessions)

    def validate_session_data(self, session_data):
        """Validate that the JSON file has the correct session format."""
    
        debug_print("DEBUG: Validating session data structure")
    
        if not isinstance(session_data, dict):
            debug_print("DEBUG: Session data is not a dictionary")
            return False
    
        # Check for required top-level keys
        required_keys = ['samples']
        if not all(key in session_data for key in required_keys):
            debug_print(f"DEBUG: Missing required keys. Found: {list(session_data.keys())}")
            return False
    
        # Check samples structure
        samples = session_data.get('samples', {})
        if not isinstance(samples, dict):
            debug_print("DEBUG: Samples is not a dictionary")
            return False
    
        # Validate sample data structure
        for sample_name, sample_data in samples.items():
            if not isinstance(sample_data, dict):
                debug_print(f"DEBUG: Sample {sample_name} data is not a dictionary")
                return False
        
            # Check for expected metrics (at least some should be present)
            metrics_found = sum(1 for metric in self.metrics if metric in sample_data)
            if metrics_found == 0:
                debug_print(f"DEBUG: No valid metrics found in sample {sample_name}")
                return False
    
        debug_print("DEBUG: Session data validation passed")
        return True

    def show_merge_sessions_dialog(self, loaded_sessions):
        """Show dialog to configure session merging."""
    
        debug_print(f"DEBUG: Showing merge dialog for {len(loaded_sessions)} sessions")
    
        # Create dialog window
        merge_window = tk.Toplevel(self.window)
        merge_window.title("Merge Session Files")
        merge_window.geometry("600x500")
        merge_window.transient(self.window)
        merge_window.grab_set()
    
        # Main frame
        main_frame = ttk.Frame(merge_window)
        main_frame.pack(fill='both', expand=True, padx=10, pady=10)
    
        # Title
        title_label = ttk.Label(main_frame, 
                               text="Configure Session Merge", 
                               font=('Arial', 14, 'bold'))
        title_label.pack(pady=(0, 10))
    
        # Instructions
        instructions = ("Select which sessions to include in the merge.\n"
                       "Conflicting sample names will be automatically renamed.\n"
                       "Header information will be merged where possible.")
    
        ttk.Label(main_frame, text=instructions, 
                 font=('Arial', 9), justify='left').pack(pady=(0, 15))
    
        # Session selection frame with scrollbar
        selection_frame = ttk.LabelFrame(main_frame, text="Sessions to Merge", padding=10)
        selection_frame.pack(fill='both', expand=True, pady=(0, 10))
    
        # Create scrollable frame
        canvas = tk.Canvas(selection_frame, height=200)
        scrollbar = ttk.Scrollbar(selection_frame, orient="vertical", command=canvas.yview)
        scrollable_frame = ttk.Frame(canvas)
    
        scrollable_frame.bind(
            "<Configure>",
            lambda e: canvas.configure(scrollregion=canvas.bbox("all"))
        )
    
        canvas.create_window((0, 0), window=scrollable_frame, anchor="nw")
        canvas.configure(yscrollcommand=scrollbar.set)
    
        canvas.pack(side="left", fill="both", expand=True)
        scrollbar.pack(side="right", fill="y")
    
        # Session checkboxes with details
        session_vars = {}
    
        for session_name, session_info in loaded_sessions.items():
            session_data = session_info['data']
            sample_count = len(session_data.get('samples', {}))
        
            # Extract some header info for display
            header = session_data.get('header', {})
            assessor = header.get('Assessor Name', 'Unknown')
            date = header.get('Date', 'Unknown')
        
            var = tk.BooleanVar(value=True)  # Default to checked
            session_vars[session_name] = {
                'var': var,
                'data': session_data,
                'file_path': session_info['file_path']
            }
        
            # Create session info frame
            session_frame = ttk.Frame(scrollable_frame)
            session_frame.pack(fill='x', pady=2)
        
            # Checkbox and main info
            info_text = f"{session_name} ({sample_count} samples)"
            ttk.Checkbutton(session_frame, text=info_text, 
                       variable=var).pack(anchor='w')
        
            # Additional details
            details = f"   Assessor: {assessor} | Date: {date} | File: {os.path.basename(session_info['file_path'])}"
            ttk.Label(session_frame, text=details, 
                     font=('Arial', 8), foreground='gray').pack(anchor='w', padx=(20, 0))
    
        # Merge options frame
        options_frame = ttk.LabelFrame(main_frame, text="Merge Options", padding=10)
        options_frame.pack(fill='x', pady=(0, 10))
    
        # New session name
        name_frame = ttk.Frame(options_frame)
        name_frame.pack(fill='x', pady=(0, 5))
    
        ttk.Label(name_frame, text="Merged session name:").pack(side='left')
        default_name = f"Merged_Session_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
        name_var = tk.StringVar(value=default_name)
        name_entry = ttk.Entry(name_frame, textvariable=name_var, width=30)
        name_entry.pack(side='right')
    
        # Header merge strategy
        header_frame = ttk.Frame(options_frame)
        header_frame.pack(fill='x', pady=5)
    
        ttk.Label(header_frame, text="Header merge strategy:").pack(side='left')
        header_strategy = tk.StringVar(value="first")
        header_combo = ttk.Combobox(header_frame, textvariable=header_strategy, 
                                   values=["first", "most_recent", "manual"], 
                                   state='readonly', width=15)
        header_combo.pack(side='right')
    
        # Sample naming strategy
        naming_frame = ttk.Frame(options_frame)
        naming_frame.pack(fill='x', pady=5)
    
        ttk.Label(naming_frame, text="Duplicate sample naming:").pack(side='left')
        naming_strategy = tk.StringVar(value="auto_rename")
        naming_combo = ttk.Combobox(naming_frame, textvariable=naming_strategy,
                                   values=["auto_rename", "prefix_session", "skip_duplicates"],
                                   state='readonly', width=15)
        naming_combo.pack(side='right')
    
        # Buttons frame
        button_frame = ttk.Frame(main_frame)
        button_frame.pack(fill='x', pady=(10, 0))
    
        def perform_merge():
            debug_print("DEBUG: Starting merge process")
        
            # Get selected sessions
            selected_sessions = {}
            for session_name, session_info in session_vars.items():
                if session_info['var'].get():
                    selected_sessions[session_name] = session_info
        
            if len(selected_sessions) < 2:
                messagebox.showwarning("Insufficient Selection", 
                                     "Please select at least 2 sessions to merge.")
                return
        
            new_session_name = name_var.get().strip()
            if not new_session_name:
                messagebox.showwarning("Invalid Name", "Please enter a name for the merged session.")
                return
        
            debug_print(f"DEBUG: Merging {len(selected_sessions)} sessions into '{new_session_name}'")
        
            # Perform the merge
            success = self.execute_session_merge(
                selected_sessions, 
                new_session_name, 
                header_strategy.get(), 
                naming_strategy.get()
            )
        
            if success:
                merge_window.destroy()
                messagebox.showinfo("Merge Complete", 
                                  f"Successfully merged {len(selected_sessions)} sessions!\n"
                                  f"New session: {new_session_name}")
        
        def select_all():
            for session_info in session_vars.values():
                session_info['var'].set(True)
    
        def select_none():
            for session_info in session_vars.values():
                session_info['var'].set(False)
    
        # Selection buttons
        select_frame = ttk.Frame(button_frame)
        select_frame.pack(side='left')
    
        ttk.Button(select_frame, text="Select All", command=select_all).pack(side='left', padx=2)
        ttk.Button(select_frame, text="Select None", command=select_none).pack(side='left', padx=2)
    
        # Action buttons  
        action_frame = ttk.Frame(button_frame)
        action_frame.pack(side='right')
    
        ttk.Button(action_frame, text="Cancel", 
                   command=merge_window.destroy).pack(side='left', padx=5)
        ttk.Button(action_frame, text="Merge Sessions", 
                   command=perform_merge).pack(side='left', padx=5)
    
        debug_print("DEBUG: Merge dialog created successfully")

    def execute_session_merge(self, selected_sessions, new_session_name, header_strategy, naming_strategy):
        """Execute the actual merging of sessions."""
    
        debug_print(f"DEBUG: Executing merge with strategy - header: {header_strategy}, naming: {naming_strategy}")
    
        try:
            # Initialize merged session structure
            merged_session = {
                'header': {},
                'samples': {},
                'timestamp': datetime.now().isoformat(),
                'merge_info': {
                    'source_sessions': list(selected_sessions.keys()),
                    'merge_date': datetime.now().isoformat(),
                    'header_strategy': header_strategy,
                    'naming_strategy': naming_strategy
                }
            }
        
            # Merge headers based on strategy
            merged_session['header'] = self.merge_headers(selected_sessions, header_strategy)
        
            # Merge samples with conflict resolution
            merged_samples, conflicts_resolved = self.merge_samples(selected_sessions, naming_strategy)
            merged_session['samples'] = merged_samples
        
            # Create the new session in memory
            if hasattr(self, 'sessions'):
                # Using session-based structure
                self.sessions[new_session_name] = merged_session
                self.switch_to_session(new_session_name)
                if hasattr(self, 'session_var'):
                    self.session_var.set(new_session_name)
                self.update_session_combo()
            else:
                # Fallback to old structure
                self.samples = merged_samples
            
                # Update header fields
                for field, value in merged_session['header'].items():
                    if field in self.header_vars:
                        self.header_vars[field].set(value)
        
            # Update UI
            self.update_sample_combo()
            self.update_sample_checkboxes()
        
            # Select first sample if available
            if merged_samples:
                first_sample = list(merged_samples.keys())[0]
                self.sample_var.set(first_sample)
                self.load_sample_data(first_sample)
        
            self.update_plot()
        
            # Log merge details
            total_samples = len(merged_samples)
            total_sessions = len(selected_sessions)
        
            debug_print(f"DEBUG: Merge completed successfully")
            debug_print(f"DEBUG: Total samples: {total_samples}")
            debug_print(f"DEBUG: Conflicts resolved: {conflicts_resolved}")
            debug_print(f"DEBUG: Sessions merged: {total_sessions}")
        
            return True
        
        except Exception as e:
            debug_print(f"DEBUG: Merge execution failed: {e}")
            messagebox.showerror("Merge Error", f"Failed to merge sessions: {e}")
            return False

    def merge_headers(self, selected_sessions, strategy):
        """Merge header information based on the selected strategy."""
    
        debug_print(f"DEBUG: Merging headers with strategy: {strategy}")
    
        all_headers = []
        for session_name, session_info in selected_sessions.items():
            header = session_info['data'].get('header', {})
            if header:
                all_headers.append((session_name, header))
    
        if not all_headers:
            return {}
    
        if strategy == "first":
            return all_headers[0][1].copy()
    
        elif strategy == "most_recent":
            # Find header with most recent timestamp
            most_recent = None
            most_recent_time = None
        
            for session_name, header in all_headers:
                session_data = selected_sessions[session_name]['data']
                timestamp_str = session_data.get('timestamp', '')
            
                try:
                    if timestamp_str:
                        timestamp = datetime.fromisoformat(timestamp_str.replace('Z', '+00:00'))
                        if most_recent_time is None or timestamp > most_recent_time:
                            most_recent_time = timestamp
                            most_recent = header
                except:
                    pass
        
            return most_recent.copy() if most_recent else all_headers[0][1].copy()
    
        elif strategy == "manual":
            # For now, return first header - could be extended to show manual selection dialog
            return all_headers[0][1].copy()
    
        return {}

    def merge_samples(self, selected_sessions, naming_strategy):
        """Merge samples with conflict resolution."""
    
        debug_print(f"DEBUG: Merging samples with naming strategy: {naming_strategy}")
    
        merged_samples = {}
        conflicts_resolved = 0
    
        for session_name, session_info in selected_sessions.items():
            session_samples = session_info['data'].get('samples', {})
        
            for original_sample_name, sample_data in session_samples.items():
                final_sample_name = original_sample_name
            
                # Handle naming conflicts
                if final_sample_name in merged_samples:
                    conflicts_resolved += 1
                
                    if naming_strategy == "auto_rename":
                        counter = 1
                        while f"{original_sample_name}_{counter}" in merged_samples:
                            counter += 1
                        final_sample_name = f"{original_sample_name}_{counter}"
                
                    elif naming_strategy == "prefix_session":
                        final_sample_name = f"{session_name}_{original_sample_name}"
                        counter = 1
                        while final_sample_name in merged_samples:
                            final_sample_name = f"{session_name}_{original_sample_name}_{counter}"
                            counter += 1
                
                    elif naming_strategy == "skip_duplicates":
                        debug_print(f"DEBUG: Skipping duplicate sample: {original_sample_name}")
                        continue
            
                # Copy sample data
                merged_samples[final_sample_name] = sample_data.copy()
            
                debug_print(f"DEBUG: Added sample {original_sample_name} as {final_sample_name}")
    
        debug_print(f"DEBUG: Sample merge complete - {len(merged_samples)} total samples, {conflicts_resolved} conflicts resolved")
    
        return merged_samples, conflicts_resolved

    def create_new_session(self, session_name=None, source_image=None):
        """Create a new session for data collection."""
        if session_name is None:
            session_name = f"Session_{self.session_counter}"
            self.session_counter += 1
    
        debug_print(f"DEBUG: Creating new session: {session_name}")
        debug_print(f"DEBUG: Source image: {source_image}")
    
        # Create new session structure
        self.sessions[session_name] = {
            'header': {field: var.get() for field, var in self.header_vars.items()},
            'samples': {},
            'timestamp': datetime.now().isoformat(),
            'source_image': source_image or ''
        }
    
        # Switch to new session
        self.current_session_id = session_name
        self.samples = self.sessions[session_name]['samples']
    
        debug_print(f"DEBUG: Session created successfully")
        debug_print(f"DEBUG: Current session ID: {self.current_session_id}")
        debug_print(f"DEBUG: Session structure: {self.sessions[session_name]}")
    
        self.update_session_combo()
        self.update_sample_combo()
        self.update_sample_checkboxes()
    
        return session_name

    def switch_to_session(self, session_id):
        """Switch to a specific session."""
        if session_id not in self.sessions:
            debug_print(f"DEBUG: Session {session_id} not found")
            return False

        debug_print(f"DEBUG: Switching from session {self.current_session_id} to {session_id}")

        # Save current session data before switching
        if self.current_session_id and self.current_session_id in self.sessions:
            self.sessions[self.current_session_id]['samples'] = self.samples
            self.sessions[self.current_session_id]['header'] = {field: var.get() for field, var in self.header_vars.items()}
            debug_print(f"DEBUG: Saved {len(self.samples)} samples to previous session")

        # Switch to new session
        self.current_session_id = session_id
        self.samples = self.sessions[session_id]['samples']

        # Update header fields with session data
        session_header = self.sessions[session_id]['header']
        for field, var in self.header_vars.items():
            if field in session_header:
                var.set(session_header[field])
            else:
                if field == "Date":
                    var.set(datetime.now().strftime("%Y-%m-%d"))
                else:
                    var.set('')

        debug_print(f"DEBUG: Switched to session {session_id} with {len(self.samples)} samples")

        # Update UI components
        self.update_sample_combo()
        self.update_sample_checkboxes()

        # Select first sample if available
        if self.samples:
            first_sample = list(self.samples.keys())[0]
            self.sample_var.set(first_sample)
            self.load_sample_data(first_sample)
            # REFRESH DISPLAYS AFTER LOADING DATA
            self.refresh_value_displays()
        else:
            self.sample_var.set('')
            self.clear_form()

        self.update_plot()
        debug_print("DEBUG: Session switch completed with display refresh")
        return True

    def update_session_combo(self):
        """Update the session selection combo box."""
        session_names = list(self.sessions.keys())
        if hasattr(self, 'session_combo'):
            self.session_combo['values'] = session_names
            debug_print(f"DEBUG: Updated session combo with {len(session_names)} sessions")

 

    def setup_layout(self):
        """Create the main layout with proper canvas sizing."""
        debug_print("DEBUG: Setting up layout with enhanced canvas sizing")

        # Create main paned window
        main_paned = tk.PanedWindow(self.window, orient='horizontal', sashrelief='raised', sashwidth=4)
        main_paned.pack(fill='both', expand=True, padx=5, pady=5)
    
        # Store reference to main_paned for resize handling
        self.main_paned = main_paned

        # === LEFT PANEL SETUP ===
        left_canvas = tk.Canvas(main_paned, bg=APP_BACKGROUND_COLOR, highlightthickness=0)
        self.left_canvas = left_canvas
        left_scrollbar = ttk.Scrollbar(main_paned, orient="vertical", command=left_canvas.yview)
        self.left_frame = ttk.Frame(left_canvas)

        # Better scroll configuration
        def configure_left_scroll(event=None):
            """Configure scroll region and handle resizing."""
            self.left_frame.update_idletasks()
        
            # Update scroll region
            left_canvas.configure(scrollregion=left_canvas.bbox("all"))
        
            # Get the current canvas size
            canvas_width = left_canvas.winfo_width()
            if canvas_width > 50:  # Valid width
                # Configure the interior frame to fill the canvas width
                left_canvas.itemconfig(left_canvas.find_all()[0], width=canvas_width-4)

        self.left_frame.bind("<Configure>", configure_left_scroll)

        # Create the canvas window
        left_canvas.create_window((0, 0), window=self.left_frame, anchor="nw")
        left_canvas.configure(yscrollcommand=left_scrollbar.set)

        # Better canvas resize handling
        def on_left_canvas_configure(event):
            """Handle left canvas resize events."""
            if event.widget == left_canvas:
                # Update scroll region
                left_canvas.configure(scrollregion=left_canvas.bbox("all"))
            
                # Ensure interior frame fills width
                canvas_width = event.width
                if canvas_width > 50 and left_canvas.find_all():
                    left_canvas.itemconfig(left_canvas.find_all()[0], width=canvas_width-4)

        left_canvas.bind('<Configure>', on_left_canvas_configure)

        # Mouse wheel scrolling
        def _on_mousewheel_left(event):
            left_canvas.yview_scroll(int(-1*(event.delta/120)), "units")
        left_canvas.bind("<MouseWheel>", _on_mousewheel_left)

        # Add to paned window
        main_paned.add(left_canvas, stretch="always")

        # === RIGHT PANEL SETUP ===
        
        right_canvas = tk.Canvas(main_paned, bg=APP_BACKGROUND_COLOR, highlightthickness=0)
        right_scrollbar = ttk.Scrollbar(main_paned, orient="vertical", command=right_canvas.yview)
        self.right_frame = ttk.Frame(right_canvas)
    
        
        self.right_canvas = right_canvas

        # Right panel configuration
        def configure_right_scroll(event=None):
            """Configure right scroll region."""
            self.right_frame.update_idletasks()
            bbox = right_canvas.bbox("all")
            if bbox:
                right_canvas.configure(scrollregion=bbox)
            
                
                # Get the height of the paned window
                if hasattr(self, 'main_paned'):
                    paned_height = self.main_paned.winfo_height()
                    if paned_height > 100:  # Valid height
                        # Set canvas to use full paned window height
                        right_canvas.configure(height=paned_height)
                        debug_print(f"DEBUG: Set right canvas height to match paned window: {paned_height}px")

        self.right_frame.bind("<Configure>", configure_right_scroll)

        # Create window with proper anchor
        self.right_canvas_window = right_canvas.create_window((0, 0), window=self.right_frame, anchor="nw")
        right_canvas.configure(yscrollcommand=right_scrollbar.set)
    
        # Configure right canvas to expand interior frame and match height
        def on_right_canvas_configure(event):
            """Handle right canvas resize events and update interior frame."""
            if event.widget == right_canvas:
                canvas_width = event.width
                canvas_height = event.height
                debug_print(f"DEBUG: Right canvas resized to: {canvas_width}x{canvas_height}")
        
                if canvas_width > 50 and canvas_height > 50:
                    # CRITICAL: Set both width AND height for the interior frame
                    right_canvas.itemconfig(self.right_canvas_window, width=canvas_width-4, height=canvas_height-4)
            
                    # Force the right frame to update its size
                    self.right_frame.configure(width=canvas_width-4, height=canvas_height-4)
                    debug_print(f"DEBUG: Updated right_frame to {canvas_width-4}x{canvas_height-4}")
            
                    # Force update of all children
                    self.right_frame.update_idletasks()
            
                    # Trigger plot resize after canvas updates (only if we have valid dimensions)
                    if hasattr(self, 'update_plot_size_for_resize') and canvas_width > 200 and canvas_height > 200:
                        self.window.after(1000, self.update_plot_size_for_resize)
    
        right_canvas.bind('<Configure>', on_right_canvas_configure)

        # Mouse wheel scrolling
        def _on_mousewheel_right(event):
            right_canvas.yview_scroll(int(-1*(event.delta/120)), "units")
        right_canvas.bind("<MouseWheel>", _on_mousewheel_right)

        # Add right canvas with stretch="always" like the left canvas
        main_paned.add(right_canvas, stretch="always")

        # Store the sash position function
        def set_initial_sash_position():
            try:
                window_width = self.window.winfo_width()
                if window_width > 100:
                    sash_position = int(window_width * 0.40)
                    main_paned.sash_place(0, sash_position, 0)
                    debug_print(f"DEBUG: Set sash position to {sash_position}")
                
                    # Force canvas height update after sash positioning
                    self.window.after(50, self.equalize_canvas_heights)
            except Exception as e:
                debug_print(f"DEBUG: Sash positioning failed: {e}")

        self.set_initial_sash_position = set_initial_sash_position

        debug_print("DEBUG: Enhanced layout setup complete")

        # Add session management and panels
        self.setup_session_selector(self.left_frame)
        self.setup_data_entry_panel()
        self.setup_plot_panel()

        # Apply sizing optimization after content is added
        self.window.after(100, self.optimize_window_size)

        # Initialize default session
        if not self.sessions:
            self.create_new_session("Default_Session")

    def equalize_canvas_heights(self):
        """Ensure both canvases have the same height."""
        if hasattr(self, 'left_canvas') and hasattr(self, 'right_canvas'):
            # Get the paned window height
            if hasattr(self, 'main_paned'):
                paned_height = self.main_paned.winfo_height()
            
                # Set both canvases to the same height
                if paned_height > 100:
                    self.left_canvas.configure(height=paned_height - 10) 
                    self.right_canvas.configure(height=paned_height - 10)
                    debug_print(f"DEBUG: Set both canvases to height: {paned_height - 10}px")
                
                    # Force update
                    self.left_canvas.update_idletasks()
                    self.right_canvas.update_idletasks()

    def optimize_window_size(self):
        """Calculate window size based on actual frame dimensions after layout."""
        debug_print("DEBUG: Starting precise window size optimization")
    
        # Force complete layout update
        self.window.update_idletasks()
        self.window.update()
    
        # Measure what the left_frame actually uses after layout
        self.left_frame.update_idletasks()
        actual_left_frame_height = self.left_frame.winfo_reqheight() 
    
        # Also measure right frame actual height for comparison
        self.right_frame.update_idletasks()
        actual_right_frame_height = self.right_frame.winfo_reqheight() 
    
        debug_print(f"DEBUG: Actual frame heights - Left: {actual_left_frame_height}px, Right: {actual_right_frame_height}px")
    
        # Width calculations (existing logic)
        left_frame_width = self.left_frame.winfo_reqwidth()
        right_frame_width = self.right_frame.winfo_reqwidth()
    
        min_plot_width = 500
        optimal_left_width = max(left_frame_width + 40, 450)
        optimal_right_width = max(min_plot_width, right_frame_width + 20)
        total_optimal_width = optimal_left_width + optimal_right_width + 50
    
        # Reduce window chrome overhead and use required height instead of actual height
        governing_content_height = max(actual_left_frame_height, actual_right_frame_height)
        window_chrome = 10  # REDUCED from 120 to 30 - just enough for title bar and borders
        total_optimal_height = governing_content_height + window_chrome
    
        debug_print(f"DEBUG: Window sized for actual frame height: {governing_content_height}px")
    
        # Screen constraints
        screen_width = self.window.winfo_screenwidth()
        screen_height = self.window.winfo_screenheight()
    
        max_usable_width = int(screen_width * 0.9)
        max_usable_height = int(screen_height * 0.85)  
    
        final_width = min(total_optimal_width, max_usable_width)
        final_height = min(total_optimal_height, max_usable_height)
    
        final_width = max(final_width, 800)
        final_height = max(final_height, 500)
    
        
        if final_height > screen_height*0.91:  
            final_height = screen_height*0.91
    
        debug_print(f"DEBUG: Final window size matching actual content: {final_width}x{final_height}")
    
        # Apply the sizing
        self.window.geometry(f"{final_width}x{final_height}")
    
        # Pass the actual required height, not the full window height
        available_height = governing_content_height  # Don't add window_chrome here
        self.window.after(50, lambda: self.configure_canvas_sizing(available_height))
    
        self.center_window()
    
        if hasattr(self, 'set_initial_sash_position'):
            self.window.after(200, self.set_initial_sash_position)

    def configure_canvas_sizing(self, available_content_height):
        """Configure canvas sizing using the frame's actual rendered size."""
        debug_print(f"DEBUG: Configuring canvas sizing")
    
        self.window.update_idletasks()
    
        # Get the required height of the left frame content
        self.left_frame.update_idletasks()
        required_frame_height = self.left_frame.winfo_reqheight()
    
        debug_print(f"DEBUG: left_frame required height: {required_frame_height}px")
    
        # Set canvas to exactly match what the frame requires
        if hasattr(self, 'left_canvas'):
            # Add small padding but not excessive
            canvas_height = required_frame_height  # Just 10px padding instead of full height
            self.left_canvas.configure(height=canvas_height)
            debug_print(f"DEBUG: Canvas set to frame's required height + padding: {canvas_height}px")
        
            # Update scroll region
            self.left_canvas.configure(scrollregion=self.left_canvas.bbox("all"))
    
        debug_print("DEBUG: Canvas sized to exact frame height - minimal gray space")

    def coordinate_panel_heights(self, final_window_height, window_chrome_height):
        """Coordinate both panels to work optimally at the chosen window height."""
        debug_print("DEBUG: Coordinating panel heights for optimal layout")
    
        # Calculate available height for panel content
        available_panel_height = final_window_height - window_chrome_height
    
        # Left Panel Strategy: Optimize scrolling behavior
        # If content fits, disable scrolling; if not, optimize scroll region
        left_content_height = self.left_frame.winfo_reqheight()
    
        if left_content_height <= available_panel_height:
            # Content fits! Set canvas to exact content height to eliminate gray space
            optimal_left_height = left_content_height + 5  # Small buffer
            debug_print(f"DEBUG: Left panel content fits - setting canvas to {optimal_left_height}px")
        else:
            # Content is taller - use available height and enable smooth scrolling
            optimal_left_height = available_panel_height - 10  # Account for scrollbar
            debug_print(f"DEBUG: Left panel content scrollable - setting canvas to {optimal_left_height}px")
    
        # Apply the left panel height optimization
        if hasattr(self, 'left_canvas'):
            self.left_canvas.configure(height=optimal_left_height)
    
        # Right Panel Strategy: Ensure plot area uses available space efficiently
        right_content_height = self.right_frame.winfo_reqheight()
    
        if right_content_height < available_panel_height:
            # Right panel has extra space - we could expand plot or center it
            extra_space = available_panel_height - right_content_height
            debug_print(f"DEBUG: Right panel has {extra_space}px extra space - content will be naturally centered")
        else:
            debug_print(f"DEBUG: Right panel content fits exactly in available space")
    
        debug_print(f"DEBUG: Panel coordination complete - both panels optimized for {final_window_height}px window")
        
    def setup_data_entry_panel(self):
        """Setup the left panel for data entry."""
        # Header section
        header_frame = ttk.LabelFrame(self.left_frame, text="Session Information", padding=10)
        header_frame.pack(fill='x', padx=5, pady=5)
        
        # Configure header_frame for 2x2 + button layout
        header_frame.grid_columnconfigure(0, weight=1)
        header_frame.grid_columnconfigure(1, weight=1) 
        header_frame.grid_columnconfigure(2, weight=1)
        header_frame.grid_columnconfigure(3, weight=1)
        debug_print("DEBUG: Configured header_frame for optimized 2x2 layout")
        
        self.header_vars = {}
        
        # Row 0: Assessor Name and Media
        ttk.Label(header_frame, text="Assessor Name:", font=FONT).grid(
            row=0, column=0, sticky='e', padx=5, pady=2)
        assessor_var = tk.StringVar()
        ttk.Entry(header_frame, textvariable=assessor_var, font=FONT, width=15).grid(
            row=0, column=1, sticky='w', padx=5, pady=2)
        self.header_vars["Assessor Name"] = assessor_var
        debug_print("DEBUG: Added Assessor Name to row 0, column 0-1")

        ttk.Label(header_frame, text="Media:", font=FONT).grid(
            row=0, column=2, sticky='e', padx=5, pady=2)
        media_var = tk.StringVar()
        ttk.Entry(header_frame, textvariable=media_var, font=FONT, width=15).grid(
            row=0, column=3, sticky='w', padx=5, pady=2)
        self.header_vars["Media"] = media_var
        debug_print("DEBUG: Added Media to row 0, column 2-3")

        # Row 1: Puff Length and Date
        ttk.Label(header_frame, text="Puff Length:", font=FONT).grid(
            row=1, column=0, sticky='e', padx=5, pady=2)
        puff_var = tk.StringVar()
        ttk.Entry(header_frame, textvariable=puff_var, font=FONT, width=15).grid(
            row=1, column=1, sticky='w', padx=5, pady=2)
        self.header_vars["Puff Length"] = puff_var
        debug_print("DEBUG: Added Puff Length to row 1, column 0-1")

        ttk.Label(header_frame, text="Date:", font=FONT).grid(
            row=1, column=2, sticky='e', padx=5, pady=2)
        date_var = tk.StringVar(value=datetime.now().strftime("%Y-%m-%d"))
        ttk.Entry(header_frame, textvariable=date_var, font=FONT, width=15).grid(
            row=1, column=3, sticky='w', padx=5, pady=2)
        self.header_vars["Date"] = date_var
        debug_print("DEBUG: Added Date to row 1, column 2-3")
            
        # Row 2: Mode switch button (centered across all columns)
        mode_button_frame = ttk.Frame(header_frame)
        mode_button_frame.grid(row=2, column=0, columnspan=4, pady=10)
        
        self.mode_button = ttk.Button(mode_button_frame, text="Switch to Comparison Mode", 
                                     command=self.toggle_mode, width=25)
        self.mode_button.pack()
        debug_print("Added mode switch button to header section")

        # Sample management section - SIMPLE NATURAL CENTERING
        sample_frame = ttk.LabelFrame(self.left_frame, text="Sample Management", padding=10)
        sample_frame.pack(fill='x', padx=5, pady=5)
        
        debug_print("Setting up sample management with simple centering")
        
        # ROW 1: Sample selection - simple center using expand
        sample_select_outer = ttk.Frame(sample_frame)
        sample_select_outer.pack(fill='x', pady=5)
        
        sample_select_frame = ttk.Frame(sample_select_outer)
        sample_select_frame.pack(expand=True)  # This centers it naturally
        
        ttk.Label(sample_select_frame, text="Current Sample:", font=FONT).pack(side='left')
        self.sample_var = tk.StringVar()
        self.sample_combo = ttk.Combobox(sample_select_frame, textvariable=self.sample_var,
                                        state="readonly", width=15)
        self.sample_combo.pack(side='left', padx=5)
        self.sample_combo.bind('<<ComboboxSelected>>', self.on_sample_changed)
        
        debug_print("Sample selection centered with expand=True")
        
        # ROW 2: Buttons - simple center using expand
        button_outer = ttk.Frame(sample_frame)
        button_outer.pack(fill='x', pady=5)
        
        button_frame = ttk.Frame(button_outer)
        button_frame.pack(expand=True)  # This centers it naturally
        
        ttk.Button(button_frame, text="Add Sample", 
                  command=self.add_sample).pack(side='left', padx=2)
        ttk.Button(button_frame, text="Remove Sample", 
                  command=self.remove_sample).pack(side='left', padx=2)
        ttk.Button(button_frame, text="Clear Data", 
                  command=self.clear_current_sample).pack(side='left', padx=2)
        
        debug_print("Buttons centered with expand=True")
                  
        # Sensory evaluation section
        eval_frame = ttk.LabelFrame(self.left_frame, text="Sensory Evaluation", padding=10)
        eval_frame.pack(fill='both', expand=True, padx=5, pady=5)
        
        # Create rating scales for each metric
        self.rating_vars = {}
        self.value_labels = {}

        for i, metric in enumerate(self.metrics):
            # Create centered container for each metric row
            metric_container = ttk.Frame(eval_frame)
            metric_container.pack(fill='x', pady=4)
    
            # Center the metric frame within the container
            metric_frame = ttk.Frame(metric_container)
            metric_frame.pack(anchor='center')
            
            # Metric label
            ttk.Label(metric_frame, text=f"{metric}:", font=FONT, width=12).pack(side='left')
            
            # Container for scale and value with fixed width (50% reduction)
            scale_container = ttk.Frame(metric_frame)
            scale_container.pack(side='left', padx=5)
            
            # Rating scale (1-9) 
            self.rating_vars[metric] = tk.IntVar(value=5)
            scale = tk.Scale(scale_container, from_=1, to=9, orient='horizontal',
                           variable=self.rating_vars[metric], font=FONT, 
                           length=300, showvalue=0, tickinterval=1,
                           sliderlength=20, sliderrelief='raised', width=15)  
            scale.pack(side='left')
            debug_print(f"DEBUG: Created centered scale for {metric} with length=200, smaller pointer (sliderlength=15), and tickmarks every 1 point")
            
            # Current value display
            value_label = ttk.Label(metric_frame, text="5", width=2)
            value_label.pack(side='left', padx=(10, 0))
            
            # STORE REFERENCE TO LABEL
            self.value_labels[metric] = value_label
            debug_print(f"DEBUG: Stored reference to value label for {metric}")

            # Update value display AND plot when scale changes (LIVE UPDATES)
            def update_live(val, label=value_label, var=self.rating_vars[metric], metric_name=metric):
                label.config(text=str(var.get()))
                self.auto_save_and_update()  # Add this method call
            scale.config(command=update_live)
            debug_print(f"DEBUG: Centered scale for {metric} configured with smaller pointer and tickmarks from 1-9")
            
        # Comments section
        comments_frame = ttk.Frame(eval_frame)
        comments_frame.pack(fill='x', pady=10)
        
        ttk.Label(comments_frame, text="Additional Comments:", font=FONT).pack(anchor='w')
        self.comments_text = tk.Text(comments_frame, height=4, font=FONT)
        self.comments_text.pack(fill='x', pady=2)
        
       # Auto-save comments when user types
        def on_comment_change(event=None):
            """Auto-save comments when user types."""
            current_sample = self.sample_var.get()
            if current_sample and current_sample in self.samples:
                comments = self.comments_text.get('1.0', tk.END).strip()
                self.samples[current_sample]['comments'] = comments
                debug_print(f"DEBUG: Auto-saved comments for {current_sample}: '{comments[:50]}...'")

        # Bind to text change events
        self.comments_text.bind('<KeyRelease>', on_comment_change)
        self.comments_text.bind('<FocusOut>', on_comment_change)
        self.comments_text.bind('<Button-1>', lambda e: self.window.after(100, on_comment_change))
         
    def on_window_resize_plot(self, event):
        """Handle window resize events to update plot size dynamically - ENHANCED DEBUG VERSION."""
        debug_print(f"DEBUG: RESIZE EVENT DETECTED - Widget: {event.widget}, Window: {self.window}")
        debug_print(f"DEBUG: Event widget type: {type(event.widget)}")
        debug_print(f"DEBUG: Window type: {type(self.window)}")
        debug_print(f"DEBUG: Event widget == window? {event.widget == self.window}")
        debug_print(f"DEBUG: Event widget is window? {event.widget is self.window}")
    
        # Only handle main window resize events, not child widgets
        if event.widget != self.window:
            debug_print(f"DEBUG: Ignoring resize event from child widget: {event.widget}")
            return
        
        debug_print("DEBUG: MAIN WINDOW RESIZE CONFIRMED - Processing...")
    
        # Get current window dimensions for verification
        current_width = self.window.winfo_width()
        current_height = self.window.winfo_height()
        debug_print(f"DEBUG: Current window dimensions: {current_width}x{current_height}")
    
        # Debounce rapid resize events
        if hasattr(self, '_resize_timer'):
            self.window.after_cancel(self._resize_timer)
            debug_print("DEBUG: Cancelled previous resize timer")
    
        # Schedule plot size update with a small delay to avoid excessive updates
        self._resize_timer = self.window.after(1000, self.update_plot_size_for_resize)
        debug_print("DEBUG: Scheduled plot resize update in 150ms")

    def update_plot_size_for_resize(self):
        """Update plot size with artifact prevention and frame validation."""
        try:
            debug_print("DEBUG: === PLOT SIZE UPDATE WITH VALIDATION ===")
        
            # Check if we have the necessary components
            if not hasattr(self, 'canvas_frame') or not self.canvas_frame.winfo_exists():
                debug_print("DEBUG: Canvas frame not available, skipping resize")
                return
        
            if not hasattr(self, 'fig') or not self.fig:
                debug_print("DEBUG: Figure not available, skipping resize")
                return
        
            # Wait for frame geometry to stabilize
            self.window.update_idletasks()
        
            # Use the actual plot container for sizing
            if hasattr(self, 'plot_container') and self.plot_container.winfo_exists():
                parent_for_sizing = self.plot_container
            else:
                parent_for_sizing = self.canvas_frame.master
        
            parent_for_sizing.update_idletasks()
        
            # Validate that frames have reasonable dimensions before proceeding
            parent_width = parent_for_sizing.winfo_width()
            parent_height = parent_for_sizing.winfo_height()
        
            debug_print(f"DEBUG: Parent frame size: {parent_width}x{parent_height}px")
        
            # Don't defer if size is small - just skip this update
            if parent_width < 200 or parent_height < 200:
                debug_print("DEBUG: Parent frame size too small, skipping this resize update")
                return  # Just return, don't schedule another call
        
            # Calculate new size based on validated frame dimensions
            new_width, new_height = self.calculate_dynamic_plot_size(parent_for_sizing)
            debug_print(f"DEBUG: Calculated plot size: {new_width:.2f}x{new_height:.2f} inches")
        
            # Get current figure size for comparison
            current_width, current_height = self.fig.get_size_inches()
            debug_print(f"DEBUG: Current plot size: {current_width:.2f}x{current_height:.2f} inches")
        
            # ARTIFACT PREVENTION: Only update if change is significant
            width_diff = abs(new_width - current_width)
            height_diff = abs(new_height - current_height)
            threshold = 1  # Threshold to reduce excessive updates
        
            if width_diff > threshold or height_diff > threshold:
                debug_print(f"DEBUG: Significant size change detected - updating plot")
            
                # Apply the new size
                self.fig.set_size_inches(new_width, new_height)
                            
                self.canvas.draw_idle()
                debug_print("DEBUG: Plot resize scheduled with draw_idle() to prevent artifacts")
            else:
                debug_print("DEBUG: Size change below threshold, skipping update to prevent artifacts")
        
            debug_print("DEBUG: === END PLOT SIZE UPDATE ===")
        
        except Exception as e:
            debug_print(f"DEBUG: Error during plot resize: {str(e)}")
            import traceback
            debug_print(f"DEBUG: Full traceback: {traceback.format_exc()}")

    def setup_plot_panel(self):
        """Setup the right panel for spider plot visualization with proper resizing."""
        debug_print("DEBUG: Setting up plot panel with enhanced resizing support")
   
        # Create the main plot frame with proper expansion settings
        plot_frame = ttk.LabelFrame(self.right_frame, text="Sensory Profile Comparison", padding=10)
        plot_frame.pack(fill='both', expand=True, padx=5, pady=5)
    
        # Sample selection for plotting (fixed height)
        control_frame = ttk.Frame(plot_frame)
        control_frame.pack(side='top', fill='x', pady=(0, 5))
    
        ttk.Label(control_frame, text="Select Samples to Display:", font=FONT).pack(anchor='w')
    
        # Checkboxes frame (will be populated when samples are added)
        self.checkbox_frame = ttk.Frame(control_frame)
        self.checkbox_frame.pack(fill='x', pady=5)
    
        self.sample_checkboxes = {}
    
        # Canvas container frame (expandable)
        canvas_container = ttk.Frame(plot_frame)
        canvas_container.pack(side='top', fill='both', expand=True)
    
        debug_print("DEBUG: Plot panel frame hierarchy configured for proper expansion")
    
        # Store reference to the container for proper sizing
        self.plot_container = canvas_container
    
        # Plot canvas (pass the expandable container)
        self.setup_plot_canvas(canvas_container)
        
    def setup_plot_canvas(self, parent):
        """Create the matplotlib canvas for the spider plot with dynamic responsive sizing."""
        debug_print("DEBUG: Setting up enhanced plot canvas with dynamic sizing")
    
        # Calculate dynamic plot size based on available space
        dynamic_width, dynamic_height = self.calculate_dynamic_plot_size(parent)

        # Create figure with calculated responsive sizing
        self.fig, self.ax = plt.subplots(figsize=(dynamic_width, dynamic_height), subplot_kw=dict(projection='polar'))
        self.fig.patch.set_facecolor('white')
        debug_print(f"DEBUG: Created spider plot with dynamic size: {dynamic_width:.2f}x{dynamic_height:.2f} inches")

        # Adjust subplot to use more of the available figure space
        
        self.fig.subplots_adjust(left=0.1, right=0.9, top=0.85, bottom=0.1)
        debug_print("DEBUG: Applied subplot adjustments for dynamic plot")
    
        # Create canvas with proper expansion configuration
        canvas_frame = ttk.Frame(parent)
        canvas_frame.pack(fill='both', expand=True)  # CRITICAL: fill='both', expand=True
    
        # Store reference to canvas_frame for resize handling
        self.canvas_frame = canvas_frame

        # Create canvas
        self.canvas = FigureCanvasTkAgg(self.fig, canvas_frame)
        canvas_widget = self.canvas.get_tk_widget()
        canvas_widget.pack(fill='both', expand=True) 

        # Add toolbar for additional functionality
        self.setup_plot_context_menu(canvas_widget)
        debug_print("DEBUG: Plot canvas setup complete with dynamic sizing and right-click context menu")

        # Initialize empty plot
        self.update_plot()

        # Bind window resize events to update plot size
        self.window.bind('<Configure>', self.on_window_resize_plot, add=True)
        debug_print("DEBUG: Window resize binding added for dynamic plot updates")

    def ensure_canvas_expansion(self):
        """Ensure canvas frame expands to use full available height."""
        if hasattr(self, 'canvas_frame') and self.canvas_frame.winfo_exists():
            # Force the canvas frame to expand
            self.canvas_frame.update_idletasks()
        
            # Get parent dimensions
            parent = self.canvas_frame.master
            parent_height = parent.winfo_height()
        
            # Check if canvas is using full height
            canvas_height = self.canvas_frame.winfo_height()
            debug_print(f"DEBUG: Canvas frame height: {canvas_height}px, Parent height: {parent_height}px")
        
            if canvas_height < parent_height - 100:  # If significantly smaller
                debug_print("DEBUG: Canvas not using full height - forcing expansion")
                # Force redraw
                if hasattr(self, 'update_plot_size_for_resize'):
                    self.update_plot_size_for_resize()

    def setup_plot_context_menu(self, canvas_widget):
        """Set up right-click context menu for plot with essential functionality."""
        from matplotlib.backends.backend_tkagg import NavigationToolbar2Tk
    
        # Create a hidden toolbar to access the navigation functions
        self.hidden_toolbar = NavigationToolbar2Tk(self.canvas, canvas_widget.master)
        self.hidden_toolbar.pack_forget()  # Hide the toolbar but keep functionality
    
        def show_context_menu(event):
            """Show simplified right-click context menu with essential plot options."""
            context_menu = tk.Menu(self.window, tearoff=0)
        
            # Essential options only
            context_menu.add_command(label="⚙️ Configure Plot", 
                                   command=self.hidden_toolbar.configure_subplots)
            context_menu.add_separator()
            context_menu.add_command(label="💾 Save Plot...", 
                                   command=self.hidden_toolbar.save_figure)
        
            try:
                context_menu.tk_popup(event.x_root, event.y_root)
            finally:
                context_menu.grab_release()
    
        # Bind right-click to show context menu
        canvas_widget.bind("<Button-3>", show_context_menu)  
        canvas_widget.bind("<Button-2>", show_context_menu)  
    
        debug_print("Simplified right-click context menu set up for plot canvas")
        
    def toggle_mode(self):
        """Toggle between collection mode and comparison mode."""
        if self.current_mode == "collection":
            # Switch to comparison mode
            self.switch_to_comparison_mode()
        else:
            # Switch to collection mode
            self.switch_to_collection_mode()
            
    def switch_to_comparison_mode(self):
        """Switch to comparison mode - show averages across users."""
        debug_print("DEBUG: Switching to comparison mode")
    
        self.current_mode = "comparison"
        self.mode_button.config(text="Switch to Collection Mode")
    
        # Change to bright white background for comparison mode
        self.update_window_background('#FFFFFF')
    
        # Add comparison title
        self.setup_comparison_title()
    
        # Gray out sensory evaluation panel
        self.disable_sensory_evaluation()
    
        # Load multiple sessions if needed (use your existing logic)
        if not self.all_sessions_data:
            self.load_multiple_sessions()
    
        # Calculate averages (use your existing method)
        self.calculate_sample_averages()
    
        # Update plot with averages
        self.update_comparison_plot()
    
        # Bring to front after mode switch
        self.bring_to_front()
    
        debug_print("Switched to comparison mode - showing averaged data across users")
        messagebox.showinfo("Comparison Mode", "Now showing averaged data across multiple users.\nSensory evaluation is disabled in this mode.")

    def switch_to_collection_mode(self):
        """Switch to collection mode - normal single user operation."""
        debug_print("DEBUG: Switching to collection mode")
    
        self.current_mode = "collection" 
        self.mode_button.config(text="Switch to Comparison Mode")
    
        # Change back to light gray background for collection mode
        self.update_window_background(APP_BACKGROUND_COLOR)
    
        # Remove comparison title if it exists
        if hasattr(self, 'comparison_title_frame'):
            self.comparison_title_frame.destroy()
    
        # Re-enable sensory evaluation panel
        self.enable_sensory_evaluation()
    
        # Update plot with current user's data
        self.update_plot()
    
        # Bring to front after mode switch
        self.bring_to_front()
    
        debug_print("Switched to collection mode - showing single user data")
        messagebox.showinfo("Collection Mode", "Now showing single user data collection mode.\nSensory evaluation is enabled.")

    def update_widget_backgrounds(self, parent, color):
        """Recursively update background colors for tkinter widgets."""
        for child in parent.winfo_children():
            try:
                # Only update standard tkinter widgets (not ttk)
                widget_class = child.winfo_class()
                if widget_class in ['Frame', 'Label', 'Button', 'Entry', 'Text', 'Listbox']:
                    child.configure(bg=color)
            except:
                pass  # Skip widgets that don't support bg parameter
        
            # Recursively update children
            self.update_widget_backgrounds(child, color)

    def disable_sensory_evaluation(self):
        """Gray out and disable all sensory evaluation controls."""
        # Find the sensory evaluation frame and disable all children
        for widget in self.left_frame.winfo_children():
            if isinstance(widget, ttk.LabelFrame) and widget.cget('text') == 'Sensory Evaluation':
                self.set_widget_state(widget, 'disabled')
                widget.configure(style='Disabled.TLabelframe')
        debug_print("Disabled sensory evaluation panel for comparison mode")
        
    def enable_sensory_evaluation(self):
        """Re-enable all sensory evaluation controls."""
        # Find the sensory evaluation frame and enable all children
        for widget in self.left_frame.winfo_children():
            if isinstance(widget, ttk.LabelFrame) and widget.cget('text') == 'Sensory Evaluation':
                self.set_widget_state(widget, 'normal')
                widget.configure(style='TLabelframe')
        debug_print("Enabled sensory evaluation panel for collection mode")

    def bring_to_front(self):
        """Bring the sensory window to front after user actions."""
        if self.window and self.window.winfo_exists():
            self.window.lift()
            self.window.focus_set()
            debug_print("DEBUG: Brought sensory window to front")

    def update_window_background(self, color):
        """Update the window background color and all child widgets."""
        debug_print(f"DEBUG: Updating window background to {color}")
    
        # Update the main window background
        self.window.configure(bg=color)
    
        # Update ttk style for the new background
        style = ttk.Style()
        style.configure('TLabel', background=color)
        style.configure('TLabelFrame', background=color)
        style.configure('TLabelFrame.Label', background=color)
        style.configure('TFrame', background=color)
    
        # Update any direct tkinter widgets recursively
        self.update_widget_backgrounds(self.window, color)
    
        # Force a redraw
        self.window.update_idletasks()
    
    
        debug_print(f"DEBUG: Window background updated to {color} and brought to front")

    def set_widget_state(self, parent, state):
        """Recursively set state of all child widgets."""
        try:
            parent.configure(state=state)
        except:
            pass  # Some widgets don't support state
            
        for child in parent.winfo_children():
            self.set_widget_state(child, state)

    def load_multiple_sessions(self):
        """Enhanced method to load multiple session files for comparison."""
        debug_print("DEBUG: Loading multiple sessions for comparison mode")
    
        filenames = filedialog.askopenfilenames(
            title="Select Session Files for Comparison",
            filetypes=[("JSON files", "*.json"), ("All files", "*.*")]
        )
    
        if not filenames:
            debug_print("DEBUG: No files selected for comparison")
            return False
    
        if len(filenames) < 2:
            messagebox.showwarning("Warning", "Please select at least 2 session files for comparison.")
            return False
    
        successful_loads = 0
    
        for filename in filenames:
            try:
                debug_print(f"DEBUG: Loading session file: {filename}")
            
                with open(filename, 'r') as f:
                    session_data = json.load(f)
            
                # Create session name from filename
                base_filename = os.path.splitext(os.path.basename(filename))[0]
                session_name = base_filename
            
                # Ensure unique session name
                counter = 1
                original_name = session_name
                while session_name in self.sessions:
                    session_name = f"{original_name}_{counter}"
                    counter += 1
            
                # Create new session with loaded data
                self.sessions[session_name] = {
                    'header': session_data.get('header', {}),
                    'samples': session_data.get('samples', {}),
                    'timestamp': session_data.get('timestamp', datetime.now().isoformat()),
                    'source_file': filename
                }
            
                successful_loads += 1
                debug_print(f"DEBUG: Successfully loaded session {session_name} with {len(self.sessions[session_name]['samples'])} samples")
            
            except Exception as e:
                debug_print(f"DEBUG: Error loading session from {filename}: {e}")
                messagebox.showerror("Error", f"Failed to load session from {os.path.basename(filename)}: {e}")
    
        if successful_loads >= 2:
            debug_print(f"DEBUG: Successfully loaded {successful_loads} sessions for comparison")
            messagebox.showinfo("Success", f"Loaded {successful_loads} sessions for comparison.")
            return True
        else:
            messagebox.showerror("Error", "Failed to load enough sessions for comparison (minimum 2 required).")
            return False

        self.bring_to_front()
            
    def calculate_sample_averages(self):
        """Calculate averages for each sample across all loaded sessions."""
        debug_print("DEBUG: Calculating sample averages across all sessions")
    
        if len(self.sessions) < 2:
            debug_print("DEBUG: Not enough sessions for comparison")
            return
    
        sample_data = {}  # {sample_name: {metric: [values], 'comments': [comments]}}
    
        # Collect all values for each sample/metric combination
        for session_name, session_info in self.sessions.items():
            samples = session_info.get('samples', {})
            header = session_info.get('header', {})
            assessor_name = header.get('Assessor Name', session_name)
        
            debug_print(f"DEBUG: Processing session {session_name} with assessor {assessor_name}")
        
            for sample_name, sample_values in samples.items():
                if sample_name not in sample_data:
                    sample_data[sample_name] = {metric: [] for metric in self.metrics}
                    sample_data[sample_name]['comments'] = []
                    sample_data[sample_name]['assessors'] = []
            
                # Collect metric values
                for metric in self.metrics:
                    if metric in sample_values and sample_values[metric] is not None:
                        try:
                            value = float(sample_values[metric])
                            sample_data[sample_name][metric].append(value)
                        except (ValueError, TypeError):
                            debug_print(f"DEBUG: Invalid value for {metric} in {sample_name}: {sample_values[metric]}")
            
                # Collect comments
                if 'comments' in sample_values and sample_values['comments'].strip():
                    sample_data[sample_name]['comments'].append(f"{assessor_name}: {sample_values['comments']}")
            
                sample_data[sample_name]['assessors'].append(assessor_name)
    
        # Calculate averages
        self.average_samples = {}
        for sample_name, data in sample_data.items():
            self.average_samples[sample_name] = {}
        
            for metric in self.metrics:
                if data[metric]:  # If we have values
                    avg_value = sum(data[metric]) / len(data[metric])
                    self.average_samples[sample_name][metric] = round(avg_value, 1)
                    debug_print(f"DEBUG: {sample_name} {metric} average: {avg_value:.1f} (from {len(data[metric])} values)")
                else:
                    self.average_samples[sample_name][metric] = 5  # Default middle value
        
            # Combine comments
            self.average_samples[sample_name]['comments'] = '\n'.join(data['comments'])
            self.average_samples[sample_name]['assessor_count'] = len(set(data['assessors']))
    
        debug_print(f"DEBUG: Calculated averages for {len(self.average_samples)} samples across {len(self.sessions)} sessions")
        
    def update_comparison_plot(self):
        """Update plot to show averaged data across users."""
        if not self.average_samples:
            return
            
        # Temporarily replace samples with averages for plotting
        original_samples = self.samples.copy()
        original_checkboxes = self.sample_checkboxes.copy()
        
        # Set up average samples for plotting
        self.samples = self.average_samples.copy()
        
        # Update sample checkboxes to show average samples
        self.sample_checkboxes = {}
        for sample_name in self.average_samples.keys():
            var = tk.BooleanVar(value=True)  # Show all by default
            self.sample_checkboxes[sample_name] = var
            
        # Update the checkbox display
        self.update_sample_checkboxes()
        
        # Update the plot
        self.create_spider_plot()
        
        debug_print("Updated plot with averaged comparison data")

        self.bring_to_front()

    def rename_current_sample(self):
        """Rename the currently selected sample."""
        current_sample = self.sample_var.get()
        if not current_sample:
            messagebox.showwarning("No Sample Selected", "Please select a sample to rename.")
            return
    
        if current_sample not in self.samples:
            messagebox.showwarning("Sample Not Found", f"Sample '{current_sample}' not found.")
            return
    
        # Get new name from user
        new_name = tk.simpledialog.askstring(
            "Rename Sample", 
            f"Current name: {current_sample}\n\nEnter new name:",
            initialvalue=current_sample
        )
    
        if not new_name or new_name.strip() == "":
            return  # User cancelled or entered empty name
    
        new_name = new_name.strip()
    
        # Check if new name already exists
        if new_name in self.samples and new_name != current_sample:
            messagebox.showerror("Name Conflict", 
                               f"A sample named '{new_name}' already exists.\n"
                               f"Please choose a different name.")
            return
    
        # Perform the rename
        if new_name != current_sample:
            # Copy data to new key
            self.samples[new_name] = self.samples[current_sample]
        
            # Remove old key
            del self.samples[current_sample]
        
            # Update UI components
            self.update_sample_combo()
            self.update_sample_checkboxes()
        
            # Select the renamed sample
            self.sample_var.set(new_name)
            self.load_sample_data(new_name)
        
            # Update plot
            self.update_plot()
        
            debug_print(f"Renamed sample '{current_sample}' to '{new_name}'")
            messagebox.showinfo("Success", f"Sample renamed to '{new_name}'")

    def batch_rename_samples(self):
        """Rename multiple samples at once."""
        if not self.samples:
            messagebox.showwarning("No Samples", "No samples available to rename.")
            return
    
        # Create batch rename dialog
        rename_window = tk.Toplevel(self.window)
        rename_window.title("Batch Rename Samples")
        rename_window.geometry("600x400")
        rename_window.transient(self.window)
        rename_window.grab_set()
    
        # Center the window
        rename_window.update_idletasks()
        x = (rename_window.winfo_screenwidth() // 2) - (300)
        y = (rename_window.winfo_screenheight() // 2) - (200)
        rename_window.geometry(f"600x400+{x}+{y}")
    
        main_frame = ttk.Frame(rename_window, padding=10)
        main_frame.pack(fill='both', expand=True)
    
        ttk.Label(main_frame, text="Batch Rename Samples", font=('Arial', 14, 'bold')).pack(pady=5)
        ttk.Label(main_frame, text="Edit the names below, then click Apply Changes", font=('Arial', 10)).pack(pady=2)
    
        # Create scrollable frame for sample entries
        canvas = tk.Canvas(main_frame)
        scrollbar = ttk.Scrollbar(main_frame, orient="vertical", command=canvas.yview)
        scrollable_frame = ttk.Frame(canvas)
    
        canvas.configure(yscrollcommand=scrollbar.set)
        canvas.bind('<Configure>', lambda e: canvas.configure(scrollregion=canvas.bbox("all")))
    
        canvas.create_window((0, 0), window=scrollable_frame, anchor="nw")
    
        # Store the name variables
        name_vars = {}
        original_names = list(self.samples.keys())
    
        # Create entry fields for each sample
        for i, sample_name in enumerate(original_names):
            row_frame = ttk.Frame(scrollable_frame)
            row_frame.pack(fill='x', pady=2, padx=5)
        
            ttk.Label(row_frame, text=f"Sample {i+1}:", width=10).pack(side='left')
        
            name_var = tk.StringVar(value=sample_name)
            name_vars[sample_name] = name_var
        
            entry = ttk.Entry(row_frame, textvariable=name_var, width=40)
            entry.pack(side='left', padx=5, fill='x', expand=True)
    
        canvas.pack(side="left", fill="both", expand=True)
        scrollbar.pack(side="right", fill="y")
    
        # Quick action buttons
        quick_frame = ttk.Frame(main_frame)
        quick_frame.pack(fill='x', pady=10)
    
        def apply_prefix():
            prefix = tk.simpledialog.askstring("Add Prefix", "Enter prefix to add:")
            if prefix:
                for sample_name in original_names:
                    current = name_vars[sample_name].get()
                    name_vars[sample_name].set(f"{prefix}{current}")
    
        def apply_suffix():
            suffix = tk.simpledialog.askstring("Add Suffix", "Enter suffix to add:")
            if suffix:
                for sample_name in original_names:
                    current = name_vars[sample_name].get()
                    name_vars[sample_name].set(f"{current}{suffix}")
    
        def number_samples():
            base = tk.simpledialog.askstring("Number Samples", "Enter base name (e.g., 'Test'):")
            if base:
                for i, sample_name in enumerate(original_names):
                    name_vars[sample_name].set(f"{base} {i+1}")
    
        ttk.Button(quick_frame, text="Add Prefix", command=apply_prefix).pack(side='left', padx=5)
        ttk.Button(quick_frame, text="Add Suffix", command=apply_suffix).pack(side='left', padx=5)
        ttk.Button(quick_frame, text="Number Samples", command=number_samples).pack(side='left', padx=5)
    
        # Apply/Cancel buttons
        button_frame = ttk.Frame(main_frame)
        button_frame.pack(fill='x', pady=10)
    
        def apply_changes():
            # Collect all new names
            new_names = {}
            for original_name in original_names:
                new_name = name_vars[original_name].get().strip()
                if not new_name:
                    messagebox.showerror("Empty Name", "All samples must have names.")
                    return
                new_names[original_name] = new_name
        
            # Check for duplicates
            name_counts = {}
            for new_name in new_names.values():
                name_counts[new_name] = name_counts.get(new_name, 0) + 1
        
            duplicates = [name for name, count in name_counts.items() if count > 1]
            if duplicates:
                messagebox.showerror("Duplicate Names", 
                                   f"The following names are used more than once:\n{', '.join(duplicates)}\n\n"
                                   f"Please ensure all names are unique.")
                return
        
            # Apply the changes
            new_samples = {}
            for original_name in original_names:
                new_name = new_names[original_name]
                new_samples[new_name] = self.samples[original_name]
        
            self.samples = new_samples
        
            # Update UI
            current_selection = self.sample_var.get()
            if current_selection in new_names:
                new_selection = new_names[current_selection]
            else:
                new_selection = list(self.samples.keys())[0] if self.samples else ""
        
            self.update_sample_combo()
            self.update_sample_checkboxes()
        
            if new_selection:
                self.sample_var.set(new_selection)
                self.load_sample_data(new_selection)
        
            self.update_plot()
        
            rename_window.destroy()
            messagebox.showinfo("Success", f"Successfully renamed {len(original_names)} samples.")
    
        ttk.Button(button_frame, text="Apply Changes", command=apply_changes).pack(side='left', padx=5)
        ttk.Button(button_frame, text="Cancel", command=rename_window.destroy).pack(side='right', padx=5)

    def calculate_dynamic_plot_size(self, parent_frame):
        """Calculate plot size that directly scales with window size - with enhanced debugging."""
        debug_print("DEBUG: Starting SIMPLE dynamic plot size calculation")

        # Force geometry update to get current dimensions
        self.window.update_idletasks()
    
        
        if hasattr(self, 'plot_container') and self.plot_container.winfo_exists():
            parent_frame = self.plot_container
            parent_frame.update_idletasks()
    
        # Get the actual dimensions
        available_width = parent_frame.winfo_width()
        available_height = parent_frame.winfo_height()

        aspect_ratio = available_height/available_width
    
        debug_print(f"DEBUG: Parent frame: {parent_frame.__class__.__name__}")
        debug_print(f"DEBUG: Available dimensions - Width: {available_width}px, Height: {available_height}px")
    
        # Simple fallback for initial sizing or very small windows
        if available_width < 200 or available_height < 200:
            debug_print("DEBUG: Using fallback size for small window")
            return (6, 4.8)
    
        # Reserve space for controls (checkboxes, labels, etc.)
        control_space = 100  # Space for labelframe title and checkbox row
        scrollbar_space = 20  # Account for potential scrollbar
    
        # Calculate available space more accurately
        plot_height_available = available_height - control_space - scrollbar_space
        plot_width_available = available_width - scrollbar_space
    
        # Use most of the available space for the plot
        plot_width_pixels = plot_width_available - 20  # 100px margin on each side for the legend
        plot_height_pixels = plot_height_available - 10  # 5px margin top/bottom
    
        debug_print(f"DEBUG: Plot space in pixels - Width: {plot_width_pixels}px, Height: {plot_height_pixels}px")
    
        # Convert to inches for matplotlib (using standard 100 DPI)
        plot_width_inches = plot_width_pixels / 100.0
        plot_height_inches = plot_width_pixels*aspect_ratio/100
    
        # Apply minimum sizes to prevent too small plots
        plot_width_inches = max(plot_width_inches, 4.0)
        plot_height_inches = max(plot_height_inches, 3.0)
    
        debug_print(f"DEBUG: FINAL plot size - Width: {plot_width_inches:.2f} inches, Height: {plot_height_inches:.2f} inches")
    
        return (plot_width_inches, plot_height_inches)

    def initialize_plot_size(self):
        """Initialize plot size after window is fully rendered to prevent startup artifacts."""
        debug_print("DEBUG: Initializing plot size after window render")
    
        # Give window time to fully render
        self.window.update_idletasks()
    
        # Force canvas updates
        if hasattr(self, 'right_canvas'):
            self.right_canvas.update_idletasks()
    
        # Update plot size
        if hasattr(self, 'update_plot_size_for_resize'):
            self.update_plot_size_for_resize()
    
        debug_print("DEBUG: Initial plot size set")

    def create_spider_plot(self):
        """Create a spider/radar plot of sensory data."""
        self.ax.clear()
        
        # Check if we have any data to plot
        if not self.samples:
            self.ax.text(0.5, 0.5, 'No samples to display\nAdd samples to begin evaluation', 
                        transform=self.ax.transAxes, ha='center', va='center', 
                        fontsize=12, color='gray')
            self.canvas.draw()
            return
            
        # Get selected samples for plotting
        selected_samples = []
        for sample_name, checkbox_var in self.sample_checkboxes.items():
            if checkbox_var.get() and sample_name in self.samples:
                selected_samples.append(sample_name)
                    
        if not selected_samples:
            self.ax.text(0.5, 0.5, 'No samples selected\nUse checkboxes to select samples', 
                        transform=self.ax.transAxes, ha='center', va='center', 
                        fontsize=12, color='gray')
            self.canvas.draw()
            return
        
        # Setup the spider plot
        num_metrics = len(self.metrics)
        angles = np.linspace(0, 2 * np.pi, num_metrics, endpoint=False).tolist()
        angles += angles[:1]  # Complete the circle
        
        # Set up the plot
        self.ax.set_theta_offset(np.pi / 2)
        self.ax.set_theta_direction(-1)
        self.ax.set_thetagrids(np.degrees(angles[:-1]), self.metrics, fontsize=10)
        self.ax.set_ylim(0, 9)  
        self.ax.set_yticks(range(1, 10))  
        self.ax.set_yticklabels(range(1, 10))
        self.ax.grid(True, alpha=0.3)
        
        # Colors for different samples
        colors = ['red', 'green', 'blue', 'orange', 'purple', 'brown', 'pink', 'cyan', 'magenta', 'lime']
        line_styles = ['-', '--', '-.', ':']
        
        # Plot each selected sample
        for i, sample_name in enumerate(selected_samples):
            sample_data = self.samples[sample_name]
            values = [sample_data.get(metric, 5) for metric in self.metrics]  # Default to 5 if no data
            values += values[:1]  # Complete the circle
            
            color = colors[i % len(colors)]
            line_style = line_styles[i % len(line_styles)]
            
            # Plot the line and markers
            self.ax.plot(angles, values, 'o', linewidth=2.5, label=sample_name, 
                        color=color, linestyle=line_style, markersize=8, alpha=0.8)
            # Fill the area
            self.ax.fill(angles, values, alpha=0.1, color=color)
            
        # Add legend
        if selected_samples:
            self.ax.legend(loc='upper right', bbox_to_anchor=(1.1, 1.2), fontsize=8)
            
        # Set title
        self.ax.set_title('Sensory Profile Comparison', fontsize=12, fontweight='bold', pad=15, ha = 'center', y = 1.08)
        
        # Force canvas update
        self.canvas.draw_idle()
        
    def add_sample(self):
        """Add a new sample for evaluation."""
        import tkinter.simpledialog
        sample_name = tk.simpledialog.askstring("Add Sample", "Enter sample name:")
        if sample_name and sample_name not in self.samples:
            # Initialize sample data
            self.samples[sample_name] = {metric: 0 for metric in self.metrics}
            self.samples[sample_name]['comments'] = ''
            
            # Update UI
            self.update_sample_combo()
            self.update_sample_checkboxes()
            self.sample_var.set(sample_name)
            self.load_sample_data(sample_name)
            
            # force plot update
            self.update_plot()

            debug_print(f"Added sample: {sample_name}")
        elif sample_name in self.samples:
            messagebox.showwarning("Warning", "Sample name already exists!")
            
    def remove_sample(self):
        """Remove the currently selected sample."""
        current_sample = self.sample_var.get()
        if current_sample and current_sample in self.samples:
            if messagebox.askyesno("Confirm", f"Remove sample '{current_sample}'?"):
                del self.samples[current_sample]
                self.update_sample_combo()
                self.update_sample_checkboxes()
                
                # Select first available sample or clear
                if self.samples:
                    first_sample = list(self.samples.keys())[0]
                    self.sample_var.set(first_sample)
                    self.load_sample_data(first_sample)
                else:
                    self.sample_var.set('')
                    self.clear_form()
                    
                self.update_plot()
                debug_print(f"Removed sample: {current_sample}")
                
    def clear_current_sample(self):
        """Clear data for the current sample."""
        current_sample = self.sample_var.get()
        if current_sample and current_sample in self.samples:
            if messagebox.askyesno("Confirm", f"Clear data for '{current_sample}'?"):
                # Reset all ratings to 5 (neutral)
                for metric in self.metrics:
                    self.samples[current_sample][metric] = 5
                    self.rating_vars[metric].set(5)
                
                self.samples[current_sample]['comments'] = ''
                self.comments_text.delete('1.0', tk.END)
                
                self.update_plot()
                debug_print(f"Cleared data for sample: {current_sample}")
                
    def update_sample_combo(self):
        """Update the sample selection combo box."""
        sample_names = list(self.samples.keys())
        self.sample_combo['values'] = sample_names
        
    def update_sample_checkboxes(self):
        """Update the checkboxes for sample selection in plotting."""
        # Clear existing checkboxes
        for widget in self.checkbox_frame.winfo_children():
            widget.destroy()

        # Show different header based on mode
        if self.current_mode == "comparison":
            header_text = "Select Averaged Samples to Display:"
        else:
            header_text = "Select Samples to Display:"
            
        # Update the header label
        for widget in self.checkbox_frame.master.winfo_children():
            if isinstance(widget, ttk.Label):
                widget.config(text=header_text)
                break
            
        self.sample_checkboxes = {}
        
        # Create checkboxes for each sample
        for i, sample_name in enumerate(self.samples.keys()):
            var = tk.BooleanVar(value=True)  # Default to checked
            checkbox = ttk.Checkbutton(self.checkbox_frame, text=sample_name, 
                                     variable=var, command=self.update_plot)
            checkbox.grid(row=i//3, column=i%3, sticky='w', padx=5, pady=2)
            self.sample_checkboxes[sample_name] = var
            
    def on_sample_changed(self, event=None):
        """Handle sample selection change."""
        selected_sample = self.sample_var.get()
        if selected_sample in self.samples:
            self.load_sample_data(selected_sample)
            
    def load_sample_data(self, sample_name):
        """Load data for the specified sample into the form."""
        debug_print(f"DEBUG: Loading sample data for: {sample_name}")
    
        if sample_name in self.samples:
            sample_data = self.samples[sample_name]
            debug_print(f"DEBUG: Found sample data: {sample_data}")
        
            # Load ratings and update both sliders AND display labels
            for metric in self.metrics:
                value = sample_data.get(metric, 5)
                debug_print(f"DEBUG: Setting {metric} to {value}")
            
                # Update the slider value
                self.rating_vars[metric].set(value)
            
                # MANUALLY UPDATE THE DISPLAY LABEL
                if hasattr(self, 'value_labels') and metric in self.value_labels:
                    self.value_labels[metric].config(text=str(value))
                    debug_print(f"DEBUG: Updated display label for {metric} to {value}")
                else:
                    debug_print(f"DEBUG: No value label found for {metric}")
            
            # Load comments
            comments = sample_data.get('comments', '')
            self.comments_text.delete('1.0', tk.END)
            self.comments_text.insert('1.0', comments)
            debug_print(f"DEBUG: Loaded comments: '{comments[:50]}...'")
        
            debug_print(f"DEBUG: Successfully loaded all data for {sample_name}")
        else:
            debug_print(f"DEBUG: Sample {sample_name} not found in samples")

    def refresh_value_displays(self):
        """Refresh all value display labels to match current slider values."""
        debug_print("DEBUG: Refreshing all value displays")
    
        if not hasattr(self, 'value_labels'):
            debug_print("DEBUG: No value labels found, skipping refresh")
            return
        
        for metric in self.metrics:
            if metric in self.value_labels and metric in self.rating_vars:
                current_value = self.rating_vars[metric].get()
                self.value_labels[metric].config(text=str(current_value))
                debug_print(f"DEBUG: Refreshed {metric} display to {current_value}")
         
    def save_plot_as_image(self):
        """Save the current spider plot as an image file."""
        debug_print("DEBUG: Starting plot image save")
    
        if not self.samples:
            messagebox.showwarning("Warning", "No samples to save! Please add samples first.")
            return
        
        try:
            filename = filedialog.asksaveasfilename(
                title="Save Plot as Image",
                defaultextension=".png",
                filetypes=[
                    ("PNG files", "*.png"),
                    ("JPG files", "*.jpg"),
                    ("PDF files", "*.pdf"),
                    ("SVG files", "*.svg"),
                    ("All files", "*.*")
                ]
            )
        
            if filename:
                # Ensure we have the latest plot
                self.update_plot()
            
                # Save the figure with high DPI for quality
                self.fig.savefig(filename, dpi=300, bbox_inches='tight', 
                               facecolor='white', edgecolor='none')
            
                debug_print(f"DEBUG: Plot saved successfully to {filename}")
                messagebox.showinfo("Success", f"Plot saved successfully as {os.path.basename(filename)}")
            
        except Exception as e:
            debug_print(f"DEBUG: Error saving plot: {e}")
            messagebox.showerror("Error", f"Failed to save plot: {str(e)}")

    def save_table_as_image(self):
        """Save the sensory data table as an image."""
        debug_print("DEBUG: Starting table image save with comments")
    
        if not self.samples:
            messagebox.showwarning("Warning", "No data to save! Please add samples first.")
            return
        
        try:
            filename = filedialog.asksaveasfilename(
                title="Save Table as Image",
                defaultextension=".png",
                filetypes=[
                    ("PNG files", "*.png"),
                    ("JPG files", "*.jpg"),
                    ("PDF files", "*.pdf"),
                    ("All files", "*.*")
                ]
            )
        
            if filename:
                # Create table data with attributes as headers and samples as rows
                table_data = []
            
                # Header row with attributes + comments
                headers = ["Sample"] + self.metrics + ["Additional Comments"]
                table_data.append(headers)
            
                # Data rows - one per sample
                for sample_name, sample_data in self.samples.items():
                    row = [sample_name]
                    for metric in self.metrics:
                        row.append(str(sample_data.get(metric, "N/A")))
                    # Add comments - get from sample data or leave blank if empty
                    comments = sample_data.get("comments", "").strip()
                    row.append(comments if comments else "")
                    table_data.append(row)
            
                # Create figure for table with wider width to accommodate comments
                fig, ax = plt.subplots(figsize=(16, max(6, len(self.samples) * 0.5)))
                ax.axis('tight')
                ax.axis('off')
            
                # Create table
                table = ax.table(cellText=table_data[1:], colLabels=table_data[0],
                               cellLoc='center', loc='center')
                table.auto_set_font_size(False)
                table.set_fontsize(9)  # Slightly smaller font to fit more content
                table.scale(1.4, 2.2)  # Wider scale to accommodate comments
            
                # Style the table
                for i in range(len(headers)):
                    table[(0, i)].set_facecolor('#4CAF50')
                    table[(0, i)].set_text_props(weight='bold', color='white')
            
                # Make comments column wider and left-aligned
                comments_col_idx = len(headers) - 1
                for row_idx in range(len(table_data)):
                    if row_idx == 0:  # Header
                        continue
                    cell = table[(row_idx, comments_col_idx)]
                    cell.set_width(0.3)  # Make comments column wider
                    cell.set_text_props(ha='left', va='top', wrap=True)  # Left align and wrap text
            
                # Add title
                ax.set_title("Sensory Evaluation Results", fontsize=16, fontweight='bold', pad=20)
            
                # Save with high quality
                fig.savefig(filename, dpi=300, bbox_inches='tight', 
                           facecolor='white', edgecolor='none')
                plt.close(fig)
            
                debug_print(f"DEBUG: Table with comments saved successfully to {filename}")
                messagebox.showinfo("Success", f"Table saved successfully as {os.path.basename(filename)}")
            
        except Exception as e:
            debug_print(f"DEBUG: Error saving table: {e}")
            import traceback
            traceback.print_exc()
            messagebox.showerror("Error", f"Failed to save table: {str(e)}")

    def generate_powerpoint_report(self):
        """Generate a PowerPoint report using the same template as generate_test_report."""
        debug_print("DEBUG: Starting PowerPoint report generation")

        if not self.samples:
            messagebox.showwarning("Warning", "No data to export! Please add samples first.")
            return
    
        try:
            # Get save location
            filename = filedialog.asksaveasfilename(
                title="Save PowerPoint Report",
                defaultextension=".pptx",
                filetypes=[("PowerPoint files", "*.pptx"), ("All files", "*.*")]
            )
    
            if not filename:
                return
        
            debug_print(f"DEBUG: Creating PowerPoint report at {filename}")
    
            # Import required modules for PowerPoint generation
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor
            from datetime import datetime
            import tempfile
    
            # Create presentation with same template structure as existing reports
            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
    
            # Create main content slide
            main_slide = prs.slides.add_slide(prs.slide_layouts[6])
    
            # Add background and logo using existing template structure
            from resource_utils import get_resource_path
    
            background_path = get_resource_path("resources/ccell_background.png")
            if os.path.exists(background_path):
                main_slide.shapes.add_picture(background_path, Inches(0), Inches(0),
                                            width=prs.slide_width, height=prs.slide_height)
                debug_print("DEBUG: Background added successfully")
            else:
                debug_print("DEBUG: Background not found, using plain slide")
        
            logo_path = get_resource_path("resources/ccell_logo_full.png")
            if os.path.exists(logo_path):
                main_slide.shapes.add_picture(logo_path, Inches(11.21), Inches(0.43),
                                            width=Inches(1.57), height=Inches(0.53))
                debug_print("DEBUG: Logo added successfully")
    
            # Add title
            title_shape = main_slide.shapes.add_textbox(Inches(0.45), Inches(-0.04), 
                                                       Inches(10.72), Inches(0.64))
            text_frame = title_shape.text_frame
            text_frame.clear()
    
            p = text_frame.add_paragraph()
            p.text = "Sensory Evaluation Report"
            p.font.name = "Montserrat"
            p.font.size = Pt(32)
            p.font.bold = True
    
            # Create table data with proper structure (attributes as headers + comments)
            table_data = []
            headers = ["Sample"] + self.metrics + ["Additional Comments"]
    
            # Add header data if available
            header_info = []
            for field in self.header_fields:
                if field in self.header_vars and self.header_vars[field].get():
                    header_info.append(f"{field}: {self.header_vars[field].get()}")
    
            # Add data rows with current comments (including any just typed)
            for sample_name, sample_data in self.samples.items():
                row = [sample_name]
                for metric in self.metrics:
                    row.append(str(sample_data.get(metric, "N/A")))
                # Include current comments
                comments = sample_data.get("comments", "").strip()
                row.append(comments if comments else "")
                table_data.append(row)
    
            # Better layout - smaller table, larger plot area for legend
            if table_data:
                table_shape = main_slide.shapes.add_table(
                    len(table_data) + 1, len(headers),  # +1 for header row
                    Inches(0.45), Inches(1.5),
                    Inches(6.5), Inches(4.5)  
                )
                table = table_shape.table
        
                # Set header row
                for col_idx, header in enumerate(headers):
                    cell = table.cell(0, col_idx)
                    cell.text = header
                    cell.text_frame.paragraphs[0].font.bold = True
                    cell.text_frame.paragraphs[0].font.size = Pt(10) 
        
                # Set data rows
                for row_idx, row_data in enumerate(table_data, 1):
                    for col_idx, cell_value in enumerate(row_data):
                        cell = table.cell(row_idx, col_idx)
                        cell.text = str(cell_value)
                    
                        # Special formatting for comments column
                        if col_idx == len(headers) - 1:  
                            cell.text_frame.paragraphs[0].font.size = Pt(8)
                            # Set text alignment for comments
                            for paragraph in cell.text_frame.paragraphs:
                                paragraph.alignment = PP_ALIGN.LEFT
                        else:
                            cell.text_frame.paragraphs[0].font.size = Pt(9)
        
                debug_print(f"DEBUG: Table with comments added - {len(table_data)} rows and {len(headers)} columns")
    
            # Create plot with better legend positioning
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
                plot_image_path = tmp_file.name
        
            try:
                # Create a special figure for PowerPoint with better legend handling
                fig_ppt, ax_ppt = plt.subplots(figsize=(8, 6), subplot_kw=dict(projection='polar'))
                fig_ppt.patch.set_facecolor('white')
            
                # Get selected samples for plotting
                selected_samples = []
                for sample_name, checkbox_var in self.sample_checkboxes.items():
                    if checkbox_var.get() and sample_name in self.samples:
                        selected_samples.append(sample_name)
            
                # If no samples selected, select all
                if not selected_samples:
                    selected_samples = list(self.samples.keys())
            
                if selected_samples:
                    # Setup the spider plot
                    num_metrics = len(self.metrics)
                    angles = np.linspace(0, 2 * np.pi, num_metrics, endpoint=False).tolist()
                    angles += angles[:1]  # Complete the circle
                
                    # Set up the plot
                    ax_ppt.set_theta_offset(np.pi / 2)
                    ax_ppt.set_theta_direction(-1)
                    ax_ppt.set_thetagrids(np.degrees(angles[:-1]), self.metrics, fontsize=10)
                    ax_ppt.set_ylim(0, 9)
                    ax_ppt.set_yticks(range(1, 10))
                    ax_ppt.set_yticklabels(range(1, 10))
                    ax_ppt.grid(True, alpha=0.3)
                
                    # Colors for different samples
                    colors = ['red', 'green', 'blue', 'orange', 'purple', 'brown', 'pink', 'cyan', 'magenta', 'lime']
                    line_styles = ['-', '--', '-.', ':']
                
                    # Plot each selected sample
                    for i, sample_name in enumerate(selected_samples):
                        sample_data = self.samples[sample_name]
                        values = [sample_data.get(metric, 5) for metric in self.metrics]
                        values += values[:1]  # Complete the circle
                    
                        color = colors[i % len(colors)]
                        line_style = line_styles[i % len(line_styles)]
                    
                        # Plot the line and markers
                        ax_ppt.plot(angles, values, 'o', linewidth=2.5, label=sample_name, 
                                   color=color, linestyle=line_style, markersize=8, alpha=0.8)
                        # Fill the area
                        ax_ppt.fill(angles, values, alpha=0.1, color=color)
                
                    # legend positioning - inside the plot area
                    ax_ppt.legend(loc='upper right', bbox_to_anchor=(1.1, 1.1), fontsize=9)
                
                    # Set title
                    ax_ppt.set_title('Sensory Profile Comparison', fontsize=12, fontweight='bold', pad=15)
            
                # Save the PowerPoint-specific plot
                fig_ppt.savefig(plot_image_path, dpi=300, bbox_inches='tight',
                               facecolor='white', edgecolor='none')
                plt.close(fig_ppt)
        
                # Better plot positioning and sizing - more space, legend won't be cut off
                main_slide.shapes.add_picture(plot_image_path, 
                                            Inches(7.2), Inches(1.5),    # MOVED: Further left
                                            Inches(5.8), Inches(4.5))    # INCREASED: Wider plot
                debug_print("DEBUG: Plot with proper legend positioning added to PowerPoint slide")
        
            finally:
                # Clean up temporary file
                if os.path.exists(plot_image_path):
                    os.remove(plot_image_path)
    
            # Add header information as text box if available
            if header_info:
                info_shape = main_slide.shapes.add_textbox(Inches(0.45), Inches(6.2),
                                                         Inches(12.0), Inches(1.0))
                info_frame = info_shape.text_frame
                info_frame.clear()
        
                p = info_frame.add_paragraph()
                p.text = " | ".join(header_info)
                p.font.size = Pt(10)
                p.font.name = "Montserrat"
    
            # Save the presentation
            prs.save(filename)
            debug_print(f"DEBUG: PowerPoint report with auto-saved comments and proper legend saved successfully to {filename}")
            messagebox.showinfo("Success", f"PowerPoint report saved successfully as {os.path.basename(filename)}")
    
        except Exception as e:
            debug_print(f"DEBUG: Error generating PowerPoint report: {e}")
            import traceback
            traceback.print_exc()
            messagebox.showerror("Error", f"Failed to generate PowerPoint report: {str(e)}")

    def save_current_sample(self):
        """Save the current form data to the selected sample."""
        current_sample = self.sample_var.get()
        if not current_sample:
            messagebox.showwarning("Warning", "No sample selected!")
            return
            
        if current_sample not in self.samples:
            # Create new sample if it doesn't exist
            self.samples[current_sample] = {}
            
        # Save ratings
        for metric in self.metrics:
            self.samples[current_sample][metric] = self.rating_vars[metric].get()
            
        # Save comments
        comments = self.comments_text.get('1.0', tk.END).strip()
        self.samples[current_sample]['comments'] = comments
        
        self.update_plot()
        debug_print(f"Saved data for sample: {current_sample}")
        messagebox.showinfo("Success", f"Data saved for {current_sample}")
        
    def clear_form(self):
        """Clear all form fields."""
        debug_print("DEBUG: Clearing form and refreshing displays")
    
        for metric in self.metrics:
            self.rating_vars[metric].set(5)
        
            # Also update the display labels
            if hasattr(self, 'value_labels') and metric in self.value_labels:
                self.value_labels[metric].config(text="5")
                debug_print(f"DEBUG: Reset {metric} display to 5")
            
        self.comments_text.delete('1.0', tk.END)
        debug_print("DEBUG: Form cleared and displays refreshed")
        
    def update_plot(self):
        """Update the spider plot."""
        self.create_spider_plot()
        self.bring_to_front()
        
    def show_sop(self):
        """Display the Standard Operating Procedure."""
        sop_window = tk.Toplevel(self.window)
        sop_window.title("Sensory Evaluation SOP")
        sop_window.geometry("600x500")
        sop_window.configure(bg='white')
        
        text_widget = tk.Text(sop_window, wrap='word', font=('Arial', 11), 
                             bg='white', fg='black', padx=20, pady=20)
        text_widget.pack(fill='both', expand=True)
        text_widget.insert('1.0', self.sop_text)
        text_widget.config(state='disabled')
        
        # Center the SOP window
        sop_window.update_idletasks()
        x = (sop_window.winfo_screenwidth() // 2) - (300)
        y = (sop_window.winfo_screenheight() // 2) - (250)
        sop_window.geometry(f"600x500+{x}+{y}")
      
    def save_session(self):
        """Save the current session to a JSON file."""
        if not self.current_session_id or not self.sessions:
            messagebox.showwarning("Warning", "No session to save!")
            return
    
        # Make sure current samples are saved to the session
        if self.current_session_id in self.sessions:
            self.sessions[self.current_session_id]['samples'] = self.samples
            self.sessions[self.current_session_id]['header'] = {field: var.get() for field, var in self.header_vars.items()}
        
        current_session = self.sessions[self.current_session_id]
    
        if not current_session['samples']:
            messagebox.showwarning("Warning", "No sample data to save!")
            return
        
        # Default filename based on session name and assessor
        assessor_name = current_session['header'].get('Assessor Name', 'Unknown')
        safe_assessor = "".join(c for c in assessor_name if c.isalnum() or c in (' ', '-', '_')).strip()
        safe_session = "".join(c for c in self.current_session_id if c.isalnum() or c in (' ', '-', '_')).strip()
    
        default_filename = f"{safe_assessor}_{safe_session}_sensory_session.json"
    
        filename = filedialog.asksaveasfilename(
            defaultextension=".json",
            filetypes=[("JSON files", "*.json"), ("All files", "*.*")],
            title="Save Sensory Session",
            initialfile=default_filename
        )
    
        if filename:
            try:
                # Create session data with updated timestamp
                session_data = {
                    'header': current_session['header'],
                    'samples': current_session['samples'],
                    'timestamp': datetime.now().isoformat(),
                    'session_name': self.current_session_id,
                    'source_file': current_session.get('source_file', ''),
                    'source_image': current_session.get('source_image', '')
                }
            
                with open(filename, 'w') as f:
                    json.dump(session_data, f, indent=2)
                
                debug_print(f"DEBUG: Saved session {self.current_session_id} to {filename}")
                messagebox.showinfo("Success", 
                                  f"Session '{self.current_session_id}' saved to {os.path.basename(filename)}\n"
                                  f"Saved {len(current_session['samples'])} samples")
                debug_print(f"Saved sensory session to: {filename}")
            
            except Exception as e:
                debug_print(f"DEBUG: Error saving session: {e}")
                messagebox.showerror("Error", f"Failed to save session: {e}")

        self.bring_to_front()
                
    def load_session(self):
        """Load one or more sessions from JSON files as new sessions."""
        debug_print("DEBUG: Starting load session with multiple file selection")
    
        # Use askopenfilenames to allow multiple file selection
        filenames = filedialog.askopenfilenames(
            filetypes=[("JSON files", "*.json"), ("All files", "*.*")],
            title="Load Sensory Sessions (Hold Ctrl to select multiple)"
        )

        if not filenames:
            debug_print("DEBUG: No files selected for loading")
            return

        debug_print(f"DEBUG: Selected {len(filenames)} files for loading")
    
        successful_loads = 0
        failed_loads = []
        loaded_session_names = []

        for filename in filenames:
            try:
                debug_print(f"DEBUG: Processing file: {filename}")
            
                with open(filename, 'r') as f:
                    session_data = json.load(f)

                debug_print(f"DEBUG: Loading session from {filename}")
                debug_print(f"DEBUG: Session data keys: {list(session_data.keys())}")

                # Validate session data
                if not self.validate_session_data(session_data):
                    debug_print(f"DEBUG: Invalid session data in {filename}")
                    failed_loads.append(f"{os.path.basename(filename)} - Invalid format")
                    continue

                # Create session name from filename
                base_filename = os.path.splitext(os.path.basename(filename))[0]
                session_name = base_filename

                # Ensure unique session name
                counter = 1
                original_name = session_name
                while session_name in self.sessions:
                    session_name = f"{original_name}_{counter}"
                    counter += 1

                debug_print(f"DEBUG: Creating new session: {session_name}")

                # Create new session with loaded data
                self.sessions[session_name] = {
                    'header': session_data.get('header', {}),
                    'samples': session_data.get('samples', {}),
                    'timestamp': session_data.get('timestamp', datetime.now().isoformat()),
                    'source_file': filename
                }

                debug_print(f"DEBUG: Session created with {len(self.sessions[session_name]['samples'])} samples")
                successful_loads += 1
                loaded_session_names.append(session_name)

            except Exception as e:
                debug_print(f"DEBUG: Error loading session from {filename}: {e}")
                import traceback
                traceback.print_exc()
                failed_loads.append(f"{os.path.basename(filename)} - {str(e)}")

        # Report results to user
        if successful_loads > 0:
            # Switch to the last loaded session
            last_session = loaded_session_names[-1]
            self.switch_to_session(last_session)

            # Update session selector UI
            self.update_session_combo()
            if hasattr(self, 'session_var'):
                self.session_var.set(last_session)

            # Update other UI components
            self.update_sample_combo()
            self.update_sample_checkboxes()

            # Select first sample if available
            if self.samples:
                first_sample = list(self.samples.keys())[0]
                self.sample_var.set(first_sample)
                self.load_sample_data(first_sample)
            else:
                self.sample_var.set('')
                self.clear_form()
        
            self.update_plot()

            # Create success message
            success_msg = f"Successfully loaded {successful_loads} session(s):\n"
            success_msg += "\n".join([f"• {name}" for name in loaded_session_names])
            success_msg += f"\n\nCurrently viewing: {last_session}"
            success_msg += "\nUse session selector to switch between sessions."

            debug_print(f"DEBUG: Successfully loaded {successful_loads} sessions")
            messagebox.showinfo("Sessions Loaded", success_msg)

        # Report any failures
        if failed_loads:
            failure_msg = f"Failed to load {len(failed_loads)} file(s):\n"
            failure_msg += "\n".join([f"• {fail}" for fail in failed_loads])
            debug_print(f"DEBUG: Failed to load {len(failed_loads)} files")
            messagebox.showwarning("Load Errors", failure_msg)

        # Overall result
        if successful_loads == 0:
            messagebox.showerror("Load Failed", "No valid session files could be loaded.")
        else:
            debug_print(f"DEBUG: Load session completed - {successful_loads} successful, {len(failed_loads)} failed")

        self.bring_to_front()
                
    def export_to_excel(self):
        """Export the sensory data to an Excel file."""
        if not self.samples:
            messagebox.showwarning("Warning", "No data to export!")
            return
            
        filename = filedialog.asksaveasfilename(
            defaultextension=".xlsx",
            filetypes=[("Excel files", "*.xlsx"), ("All files", "*.*")],
            title="Export Sensory Data"
        )
        
        if filename:
            try:
                # Create DataFrame for sensory data
                data_rows = []
                for sample_name, sample_data in self.samples.items():
                    row = {'Sample': sample_name}
                    
                    # Add header information
                    for field, var in self.header_vars.items():
                        row[field] = var.get()
                        
                    # Add sensory ratings
                    for metric in self.metrics:
                        row[metric] = sample_data.get(metric, 0)
                        
                    row['Comments'] = sample_data.get('comments', '')
                    data_rows.append(row)
                    
                df = pd.DataFrame(data_rows)
                
                # Save to Excel
                with pd.ExcelWriter(filename, engine='openpyxl') as writer:
                    df.to_excel(writer, sheet_name='Sensory Data', index=False)
                    
                    # Save spider plot as image
                    plot_filename = filename.replace('.xlsx', '_spider_plot.png')
                    self.fig.savefig(plot_filename, dpi=300, bbox_inches='tight')
                    
                messagebox.showinfo("Success", f"Data exported to {filename}\nSpider plot saved as {plot_filename}")
                debug_print(f"Exported sensory data to: {filename}")
                
            except Exception as e:
                messagebox.showerror("Error", f"Failed to export data: {e}")

        self.bring_to_front()

    def auto_save_and_update(self):
        """Automatically save current sample data and update plot."""
        current_sample = self.sample_var.get()
        if current_sample and current_sample in self.samples:
            # Auto-save ratings
            for metric in self.metrics:
                self.samples[current_sample][metric] = self.rating_vars[metric].get()
        
            # Also auto-save comments
            comments = self.comments_text.get('1.0', tk.END).strip()
            self.samples[current_sample]['comments'] = comments
        
            # Update plot immediately
            self.update_plot()
        
            debug_print(f"DEBUG: Auto-saved all data for {current_sample}")

    def setup_comparison_title(self):
        """Add or update the comparison mode title."""
        # Remove existing title if it exists
        if hasattr(self, 'comparison_title_frame'):
            self.comparison_title_frame.destroy()
    
        if self.current_mode == "comparison":
            # Create title frame at the top of the window
            self.comparison_title_frame = ttk.Frame(self.window)
            self.comparison_title_frame.pack(side='top', fill='x', pady=10)
        
            # Add the title label with white background
            title_label = ttk.Label(
                self.comparison_title_frame, 
                text="Comparing Average Sensory Results",
                font=("Arial", 16, "bold"),
                background='#FFFFFF',  # Now bright white
                foreground="black",
                anchor='center'
            )
            title_label.pack(expand=True)
        
            # Ensure window stays on top after adding title
            self.window.update_idletasks()
            self.bring_to_front()
        
            debug_print("DEBUG: Added comparison mode title with white background and brought to front")

    def setup_session_selector(self, parent_frame):
        """Add session selector to the interface."""
        # Add session selector frame with reduced width
        session_frame = ttk.LabelFrame(parent_frame, text="Session Management", padding=10)
        session_frame.pack(fill='x', padx=5, pady=5)
    
        # Configure session_frame for centered grid layout
        session_frame.grid_columnconfigure(0, weight=1)
        session_frame.grid_columnconfigure(1, weight=1)
        debug_print("DEBUG: Configured session_frame for centered layout")

        # Top row - Session selection centered
        top_frame = ttk.Frame(session_frame)
        top_frame.grid(row=0, column=0, columnspan=2, pady=(0, 5))
    
        session_label = ttk.Label(top_frame, text="Current Session:", font=FONT)
        session_label.pack(side='left', padx=(0, 5))

        self.session_var = tk.StringVar()
        self.session_combo = ttk.Combobox(top_frame, textvariable=self.session_var, 
                                         font=FONT, state='readonly', width=15)
        self.session_combo.pack(side='left', padx=(0, 10))
        self.session_combo.bind('<<ComboboxSelected>>', self.on_session_selected)
        debug_print("DEBUG: Session dropdown centered on top row")

        # Second row - Session management buttons centered
        button_frame = ttk.Frame(session_frame)
        button_frame.grid(row=1, column=0, columnspan=2)
    
        ttk.Button(button_frame, text="New Session", 
                   command=self.add_new_session).pack(side='left', padx=2)
        ttk.Button(button_frame, text="Combine Sessions", 
                   command=self.show_combine_sessions_dialog).pack(side='left', padx=2)
        ttk.Button(button_frame, text="Delete Session", 
                   command=self.delete_current_session).pack(side='left', padx=2)
        debug_print("DEBUG: Session management buttons centered on second row")

        debug_print("DEBUG: Session selector UI setup complete with centered layout")

    def on_session_selected(self, event=None):
        """Handle session selection change."""
        selected_session = self.session_var.get()
        if selected_session and selected_session != self.current_session_id:
            debug_print(f"DEBUG: Session selection changed to: {selected_session}")
            self.switch_to_session(selected_session)

    def add_new_session(self):
        """Add a new empty session."""
        session_name = tk.simpledialog.askstring("New Session", "Enter session name:")
        if session_name and session_name.strip():
            session_name = session_name.strip()
            if session_name in self.sessions:
                messagebox.showerror("Session Exists", f"Session '{session_name}' already exists.")
                return
        
            debug_print(f"DEBUG: Creating new session: {session_name}")
            self.create_new_session(session_name)
            self.session_var.set(session_name)
            messagebox.showinfo("Success", f"Created new session: {session_name}")

    def delete_current_session(self):
        """Delete the current session."""
        if not self.current_session_id:
            messagebox.showwarning("No Session", "No session selected to delete.")
            return
    
        if len(self.sessions) <= 1:
            messagebox.showwarning("Cannot Delete", "Cannot delete the last session.")
            return
    
        if messagebox.askyesno("Confirm Delete", 
                              f"Delete session '{self.current_session_id}'?\n"
                              f"This will permanently remove all data in this session."):
        
            session_to_delete = self.current_session_id
            debug_print(f"DEBUG: Deleting session: {session_to_delete}")
        
            # Switch to another session first
            remaining_sessions = [s for s in self.sessions.keys() if s != session_to_delete]
            if remaining_sessions:
                self.switch_to_session(remaining_sessions[0])
        
            # Delete the session
            del self.sessions[session_to_delete]
            self.update_session_combo()
        
            debug_print(f"DEBUG: Session {session_to_delete} deleted successfully")
            messagebox.showinfo("Success", f"Session '{session_to_delete}' deleted.")

    def show_combine_sessions_dialog(self):
        """Show dialog to select and combine multiple sessions."""
        if len(self.sessions) < 2:
            messagebox.showinfo("Insufficient Sessions", 
                              "Need at least 2 sessions to combine.")
            return
    
        # Create dialog window
        combine_window = tk.Toplevel(self.window)
        combine_window.title("Combine Sessions")
        combine_window.geometry("400x300")
        combine_window.transient(self.window)
        combine_window.grab_set()
    
        # Instructions
        ttk.Label(combine_window, 
                 text="Select sessions to combine into a new session:",
                 font=FONT).pack(pady=10)
    
        # Session selection with checkboxes
        selection_frame = ttk.Frame(combine_window)
        selection_frame.pack(fill='both', expand=True, padx=20, pady=10)
    
        session_vars = {}
        for session_id in self.sessions.keys():
            var = tk.BooleanVar()
            session_vars[session_id] = var
        
            # Create checkbox with session info
            sample_count = len(self.sessions[session_id]['samples'])
            source_image = self.sessions[session_id].get('source_image', 'Manual')
            source_name = os.path.basename(source_image) if source_image else 'Manual'
        
            checkbox_text = f"{session_id} ({sample_count} samples) - {source_name}"
            ttk.Checkbutton(selection_frame, text=checkbox_text, 
                           variable=var).pack(anchor='w', pady=2)
    
        # New session name
        name_frame = ttk.Frame(combine_window)
        name_frame.pack(fill='x', padx=20, pady=10)
    
        ttk.Label(name_frame, text="New session name:").pack(side='left')
        name_var = tk.StringVar(value=f"Combined_Session_{self.session_counter}")
        name_entry = ttk.Entry(name_frame, textvariable=name_var, width=20)
        name_entry.pack(side='right')
    
        # Buttons
        button_frame = ttk.Frame(combine_window)
        button_frame.pack(fill='x', padx=20, pady=20)
    
        def combine_selected_sessions():
            selected_sessions = [sid for sid, var in session_vars.items() if var.get()]
            new_session_name = name_var.get().strip()
        
            debug_print(f"DEBUG: Combining sessions: {selected_sessions}")
            debug_print(f"DEBUG: New session name: {new_session_name}")
        
            if len(selected_sessions) < 2:
                messagebox.showwarning("Insufficient Selection", 
                                     "Select at least 2 sessions to combine.")
                return
        
            if not new_session_name:
                messagebox.showwarning("Invalid Name", "Enter a name for the new session.")
                return
        
            if new_session_name in self.sessions:
                messagebox.showerror("Name Exists", 
                                   f"Session '{new_session_name}' already exists.")
                return
        
            # Combine sessions
            combined_samples = {}
            total_sample_count = 0
        
            for session_id in selected_sessions:
                session_samples = self.sessions[session_id]['samples']
                for sample_name, sample_data in session_samples.items():
                    # Create unique sample name if conflict
                    unique_name = sample_name
                    counter = 1
                    while unique_name in combined_samples:
                        unique_name = f"{sample_name}_{counter}"
                        counter += 1
                
                    combined_samples[unique_name] = sample_data
                    total_sample_count += 1
                    debug_print(f"DEBUG: Added sample {unique_name} from session {session_id}")
        
            # Create new combined session
            self.create_new_session(new_session_name)
            self.sessions[new_session_name]['samples'] = combined_samples
            self.samples = combined_samples
        
            # Update UI
            self.session_var.set(new_session_name)
            self.update_sample_combo()
            self.update_sample_checkboxes()
            self.update_plot()
        
            combine_window.destroy()
        
            debug_print(f"DEBUG: Successfully combined {len(selected_sessions)} sessions")
            debug_print(f"DEBUG: New session has {total_sample_count} samples")
        
            messagebox.showinfo("Success", 
                              f"Combined {len(selected_sessions)} sessions into '{new_session_name}'!\n"
                              f"Total samples: {total_sample_count}")
    
        ttk.Button(button_frame, text="Combine Sessions", 
                   command=combine_selected_sessions).pack(side='left', padx=5)
        ttk.Button(button_frame, text="Cancel", 
                   command=combine_window.destroy).pack(side='right', padx=5)
    
        debug_print("DEBUG: Combine sessions dialog created")

    def setup_interface(self):
        """Set up the main interface with session management."""
        # Create main frames
        main_frame = ttk.Frame(self.window)
        main_frame.pack(fill='both', expand=True, padx=10, pady=10)
    
        # Add session management at the top
        self.setup_session_selector(main_frame)
    
    
        # Initialize with default session
        if not self.sessions:
            self.create_new_session("Default_Session")
